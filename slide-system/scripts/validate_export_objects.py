#!/usr/bin/env python3
"""validate_export_objects.py — the SINGLE QA gate for PPTX export (plan §3.3).

Every quality verdict lives here, in one exit code. Build and compare never
judge; capture/build crash only on operational errors.

Checks:
  1. export-manifest.json against the documented contract
     (_reference/export-manifest.schema.json — implemented in code, no new deps).
  2. PPTX structure vs manifest: picture count per slide (FAIL when a slide
     declares overlays but holds a single picture), overlay shape names,
     overlay bounds within tolerance, text box count, z-order (shape sequence
     follows the manifest's merged z list).
  3. Parity report.json files (from compare_renders.py) vs the thresholds in
     slide-system/registries/export-qa-thresholds.json.
  4. Native text geometry: no text box off-slide, and no non-wrapping box whose
     line is wider than the box (PowerPoint does not clip text to its shape, so
     that overruns the neighbouring shape — invisible to browser/PDF parity).

Exit 0 = PASS, 1 = FAIL. Writes validation-report.json next to --pptx.

Usage:
    validate_export_objects.py --pptx deck.pptx --manifest export-manifest.json \
        [--parity-dir qa/parity] [--thresholds <json>] [--tolerance-in 0.02]
"""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path

from _common import SYSTEM_ROOT, load_json, write_json

EMU_PER_IN = 914400
DEFAULT_THRESHOLDS = SYSTEM_ROOT / "registries" / "export-qa-thresholds.json"


def expected_svg_blips(overlays: list[dict], vector_root: Path) -> int:
    """Count only vector sources the builder could actually embed.

    Missing relative SVGs intentionally fall back to an overlay PNG. That PPTX
    remains layered and visually verifiable, so the validator must not demand
    an svgBlip the builder could never have written.
    """
    return sum(
        1 for overlay in overlays
        if overlay.get("vector_source")
        and not overlay.get("css_effects")
        and (vector_root / overlay["vector_source"]).is_file()
    )


def check_manifest_contract(manifest: dict, failures: list[str]) -> None:
    for field in ("manifest_version", "mode", "slides"):
        if field not in manifest:
            failures.append(f"manifest: missing required field '{field}'")
            return
    if manifest["manifest_version"] != 2:
        failures.append(f"manifest: manifest_version must be 2, got {manifest['manifest_version']!r}")
    if manifest["mode"] not in ("layered", "flat"):
        failures.append(f"manifest: mode must be layered|flat, got {manifest['mode']!r}")
    for entry in manifest.get("slides", []):
        sid = entry.get("slide", "?")
        for field in ("slide", "base", "text"):
            if field not in entry:
                failures.append(f"manifest slide {sid}: missing '{field}'")
        for ov in entry.get("objects", []):
            for field in ("id", "png", "bounds", "z", "sha256"):
                if field not in ov:
                    failures.append(f"manifest slide {sid} overlay {ov.get('id', '?')}: missing '{field}'")


def check_pptx_structure(pptx_path: Path, manifest: dict, tolerance_in: float,
                         failures: list[str], vector_root: Path) -> dict:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(pptx_path))
    slide_w_in = prs.slide_width / EMU_PER_IN
    slide_h_in = prs.slide_height / EMU_PER_IN
    summary = {"slides": []}
    slides = list(prs.slides)
    entries = manifest.get("slides", [])
    if len(slides) != len(entries):
        failures.append(f"pptx: slide count {len(slides)} != manifest {len(entries)}")

    import zipfile

    with zipfile.ZipFile(pptx_path) as zf:
        slide_names = [
            n for n in zf.namelist()
            if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)
        ]
        slide_xml = {
            int(re.search(r"slide(\d+)\.xml$", name).group(1)):
                zf.read(name).decode("utf-8", "ignore")
            for name in slide_names
        }

    for entry, slide in zip(entries, slides):
        sid = entry["slide"]
        cw = float(entry.get("canvasW") or manifest.get("canvasW") or 1920)
        ch = float(entry.get("canvasH") or manifest.get("canvasH") or 1080)
        pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        autoshapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        textboxes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()
                     and s.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE]
        declared_ov = entry.get("objects", [])
        declared_nv = entry.get("natives", [])
        expected_pictures = 1 + len(declared_ov)
        expected_text = len([t for t in entry.get("text", []) if t.get("text", "").strip()])

        if declared_ov and len(pictures) <= 1:
            failures.append(f"slide {sid}: manifest declares {len(declared_ov)} overlays "
                            f"but PPTX holds {len(pictures)} picture(s) — objects glued to background")
        if len(pictures) != expected_pictures:
            failures.append(f"slide {sid}: picture count {len(pictures)} != expected {expected_pictures}")
        if len(autoshapes) != len(declared_nv):
            failures.append(f"slide {sid}: autoshape count {len(autoshapes)} != "
                            f"declared natives {len(declared_nv)}")
        if len(textboxes) != expected_text:
            failures.append(f"slide {sid}: text box count {len(textboxes)} != expected {expected_text}")

        # Names + bounds (tolerance in inches) for overlays AND natives.
        by_name = {s.name: s for s in list(pictures) + list(autoshapes)}
        for kind_label, decl in (("Overlay", declared_ov), ("Native", declared_nv)):
            for item in decl:
                name = f"{kind_label}: {item['id']}"
                shape = by_name.get(name)
                if shape is None:
                    failures.append(f"slide {sid}: missing shape named '{name}'")
                    continue
                b = item["bounds"]
                expect = (b["x"] / cw * slide_w_in, b["y"] / ch * slide_h_in,
                          b["w"] / cw * slide_w_in, b["h"] / ch * slide_h_in)
                actual = (shape.left / EMU_PER_IN, shape.top / EMU_PER_IN,
                          shape.width / EMU_PER_IN, shape.height / EMU_PER_IN)
                for label, e, g in zip(("x", "y", "w", "h"), expect, actual):
                    if abs(e - g) > tolerance_in:
                        failures.append(f"slide {sid} {name}: {label} off by "
                                        f"{abs(e - g):.4f}in (> {tolerance_in}in)")

        # svgBlip only applies when the referenced source is available. The
        # builder otherwise writes the same positioned PNG overlay as fallback.
        expected_svg = expected_svg_blips(declared_ov, vector_root)
        actual_svg = slide_xml.get(sid, "").count("svgBlip")
        if actual_svg < expected_svg:
            failures.append(f"slide {sid}: svgBlip count {actual_svg} < expected "
                            f"{expected_svg} (available vector_source overlays without css effects)")

        # Z-order: shapes after the base must follow the manifest's merged z list.
        merged = sorted(
            [("overlay", f"Overlay: {ov['id']}", int(ov.get("z", 0))) for ov in declared_ov]
            + [("native", f"Native: {nv['id']}", int(nv.get("z", 0))) for nv in declared_nv]
            + [("text", None, int(t.get("z", 0))) for t in entry.get("text", [])
               if t.get("text", "").strip()],
            key=lambda t: t[2])
        actual_kinds = []
        for s in list(slide.shapes)[1:]:  # skip base picture
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                actual_kinds.append(("overlay", s.name))
            elif s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                actual_kinds.append(("native", s.name))
            elif s.has_text_frame and s.text_frame.text.strip():
                actual_kinds.append(("text", None))
        expected_kinds = [(k, n) for k, n, _ in merged]
        if actual_kinds != expected_kinds:
            failures.append(f"slide {sid}: z-order mismatch — expected {expected_kinds}, "
                            f"got {actual_kinds}")

        summary["slides"].append({"slide": sid, "pictures": len(pictures),
                                  "autoshapes": len(autoshapes),
                                  "textboxes": len(textboxes),
                                  "svg_blips": actual_svg})
    return summary


# A non-wrapping box is only reported once its single line is estimated to run
# this far past its own width. PowerPoint does not clip text to its shape, so
# anything beyond the box paints over whatever sits next to it. The slack
# absorbs the glyph-width estimate's error on a heading that legitimately fills
# its box (~10-15% high); the defect this gate exists for overruns by ~95%.
TEXT_OVERFLOW_SLACK = 1.35
# …and by at least one em in absolute terms. A ratio alone misjudges tiny boxes
# shrink-wrapped to one glyph (a numbered badge in a 6pt box): the estimator is
# a whole glyph off there, while an overrun narrower than a single character
# cannot reach a neighbouring shape. Real card-body overruns are many ems wide.
TEXT_OVERFLOW_MIN_EM = 1.0


def check_text_overflow(pptx_path: Path, failures: list[str]) -> dict:
    """Native text rendering outside its own text box — the PPTX-only defect.

    Browser and PDF parity cannot see this: the browser wrapped the same string
    inside the same box, so both look correct while PowerPoint paints one long
    line across its neighbours. Checked here on the real exported shapes.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from build_hybrid_pptx import AVERAGE_GLYPH_WIDTH

    prs = Presentation(str(pptx_path))
    slide_w_in = prs.slide_width / EMU_PER_IN
    slide_h_in = prs.slide_height / EMU_PER_IN
    checked = overflowing = 0

    for index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or not shape.has_text_frame:
                continue
            frame = shape.text_frame
            if not frame.text.strip():
                continue
            checked += 1
            left_in = shape.left / EMU_PER_IN
            top_in = shape.top / EMU_PER_IN
            width_pt = shape.width / EMU_PER_IN * 72.0

            if (left_in + shape.width / EMU_PER_IN > slide_w_in + 0.02
                    or top_in + shape.height / EMU_PER_IN > slide_h_in + 0.02):
                failures.append(
                    f"slide {index} text {shape.name!r}: box extends past the slide "
                    f"({left_in:.2f}+{shape.width / EMU_PER_IN:.2f} x "
                    f"{top_in:.2f}+{shape.height / EMU_PER_IN:.2f}in on "
                    f"{slide_w_in:.2f}x{slide_h_in:.2f}in)")

            if frame.word_wrap:
                continue  # PowerPoint keeps a wrapping box inside its own width
            for para in frame.paragraphs:
                text = "".join(r.text for r in para.runs).strip()
                if not text:
                    continue
                size_pt = max((r.font.size.pt for r in para.runs if r.font.size), default=0)
                estimate_pt = len(text) * size_pt * AVERAGE_GLYPH_WIDTH
                if (estimate_pt > width_pt * TEXT_OVERFLOW_SLACK
                        and estimate_pt - width_pt > size_pt * TEXT_OVERFLOW_MIN_EM):
                    overflowing += 1
                    failures.append(
                        f"slide {index} text {shape.name!r}: non-wrapping line needs "
                        f"~{estimate_pt:.0f}pt in a {width_pt:.0f}pt box — it will render "
                        f"over neighbouring shapes. Enable word wrap for this box "
                        f"(build_hybrid_pptx wrapping contract)")

    return {"text_boxes_checked": checked, "overflowing": overflowing}


def check_parity(parity_dir: Path, thresholds: dict, failures: list[str]) -> list[dict]:
    reports = sorted(parity_dir.rglob("report.json"))
    if not reports:
        failures.append(f"parity: no report.json found under {parity_dir}")
        return []
    seen = []
    for report_path in reports:
        report = load_json(report_path)
        relative_parts = report_path.relative_to(parity_dir).parts
        tier = next((part for part in relative_parts if part in ("tier1", "tier2")), None)
        if tier is None:
            failures.append(f"parity: unrecognized report path {report_path}")
            continue
        slide = next((part for part in relative_parts if part.startswith("slide-")), "?")
        gate = thresholds[tier]
        metrics = report["metrics"]
        ok = (metrics["mean_absolute_error"] <= gate["max_mean_err"]
              and metrics["changed_pixel_ratio"] <= gate["max_changed_ratio"])
        if not ok:
            failures.append(
                f"parity {slide}/{tier}: mean_err "
                f"{metrics['mean_absolute_error']} / changed_ratio "
                f"{metrics['changed_pixel_ratio']} exceeds {gate}")
        seen.append({"report": str(report_path), "tier": tier, "pass": ok, **metrics})
    return seen


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--manifest", required=True)
    parser.add_argument("--parity-dir", help="Directory holding compare_renders report.json files")
    parser.add_argument("--thresholds", default=str(DEFAULT_THRESHOLDS))
    parser.add_argument("--tolerance-in", type=float, default=0.02)
    parser.add_argument("--allow-untagged", action="store_true",
                        help="Accept layered decks whose visuals are untagged "
                             "(they bake into the background)")
    parser.add_argument("--allow-full-bleed", action="store_true",
                        help="Accept overlays that cover (almost) the whole canvas "
                             "(objects inside them stay glued into one picture)")
    args = parser.parse_args()

    failures: list[str] = []
    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        raise SystemExit(f"PPTX not found: {pptx_path}")
    try:
        manifest = load_json(args.manifest)
    except (OSError, json.JSONDecodeError) as error:
        raise SystemExit(f"Manifest unreadable: {error}")

    check_manifest_contract(manifest, failures)

    # Tagging-contract gate: a layered export whose visuals carry no
    # data-export-* tags silently degrades to a baked background — the exact
    # regression this whole pipeline exists to prevent. Fail loudly unless the
    # caller explicitly accepts it.
    if manifest.get("mode") == "layered" and not args.allow_untagged:
        for entry in manifest.get("slides", []):
            untagged = entry.get("untagged_visuals", [])
            if untagged:
                sample = ", ".join(f"<{u['tag']}> {u['w']}x{u['h']}" for u in untagged[:3])
                failures.append(
                    f"slide {entry.get('slide', '?')}: {len(untagged)} untagged visual(s) "
                    f"baked into background ({sample}…) — tag them with "
                    f"data-export-layer/data-export-native in the deck HTML, "
                    f"or pass --allow-untagged to accept a flattened result")

    # Granularity gate: tagging a full-page artwork as ONE overlay passes the
    # untagged check but glues every object into a single picture — a renamed
    # background. Fail any overlay covering >= max_ratio of the canvas.
    if manifest.get("mode") == "layered" and not args.allow_full_bleed:
        max_ratio = load_json(args.thresholds).get(
            "overlay_coverage", {}).get("max_ratio", 0.85)
        for entry in manifest.get("slides", []):
            cw = float(entry.get("canvasW") or manifest.get("canvasW") or 1920)
            ch = float(entry.get("canvasH") or manifest.get("canvasH") or 1080)
            for ov in entry.get("objects", []):
                b = ov.get("visual_bounds") or ov.get("bounds") or {}
                coverage = (b.get("w", 0) * b.get("h", 0)) / (cw * ch)
                if coverage >= max_ratio:
                    failures.append(
                        f"slide {entry.get('slide', '?')} overlay '{ov.get('id', '?')}': "
                        f"covers {coverage:.0%} of the canvas (>= {max_ratio:.0%}) — "
                        f"this is a full-bleed artwork tagged as one overlay, not a "
                        f"separated object. Decompose it into per-object overlays "
                        f"(one tag per card/arrow/icon/illustration), or pass "
                        f"--allow-full-bleed to accept the merged picture")

    structure = check_pptx_structure(pptx_path, manifest, args.tolerance_in, failures,
                                     Path(args.manifest).resolve().parent) \
        if not failures else {}
    if structure:
        structure["text_overflow"] = check_text_overflow(pptx_path, failures)
    parity = []
    if args.parity_dir:
        thresholds = load_json(args.thresholds)
        parity = check_parity(Path(args.parity_dir), thresholds, failures)

    verdict = {
        "pptx": str(pptx_path),
        "mode": manifest.get("mode"),
        "pass": not failures,
        "failures": failures,
        "structure": structure,
        "parity": parity,
    }
    write_json(pptx_path.with_suffix(".validation.json"), verdict)
    print(json.dumps(verdict, indent=2, ensure_ascii=False))
    return 0 if not failures else 1


if __name__ == "__main__":
    raise SystemExit(main())
