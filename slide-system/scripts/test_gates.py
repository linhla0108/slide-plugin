#!/usr/bin/env python3
"""Unit tests for the slide-pipeline gate scripts.

Covers the highest-consequence, previously-untested paths:
  - cleanup_run: deck.html + .pptx survive; intermediates are removed.
  - score_visual_items: published-only reuse/text-only decisions + extraction recommendation.
  - score_visual_items hybrid retrieval: capped secondary lexical credit,
    anti-use-case / count-fit / zero-slot penalties, published-only enrichment.
  - validate_selection_report: equal-score plausibility, provenance, T1 shape-lock.
  - scaffold_slide_from_component: slots preserved, no base64, .bg placeholder.
  - validate_component_fidelity: slot-id coverage pass/fail.
  - read_text_slots: slim projection shape.

Run directly (`python3 test_gates.py`) or under pytest. No network, no install.
"""

from __future__ import annotations

import json
import sys
import tempfile
import importlib.util
import subprocess
from pathlib import Path

SCRIPTS = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPTS))

import cleanup_run
import compare_renders
import export_pptx
import score_visual_items as svi
import validate_selection_report as vsr
import scaffold_slide_from_component as scaffold
import validate_component_fidelity as fidelity
import validate_export_objects as export_objects
import validate_deck_assets as deck_assets
import read_text_slots
import _common
import validate_slot_content_plan as slot_plan
import component_units as units

REGISTRY = SCRIPTS.parent / "registries" / "visual-library.json"

# A real template item with positioned slots (verified to have .slot divs).
ITEM_WITH_SLOTS = "sun.interview-workshop-sunriser.04-mindset"


# --------------------------------------------------------------------------- #
# slot-content-plan — pre-build capacity contract
# --------------------------------------------------------------------------- #
def _slot_contract(slot_id: str = "headline", *, role: str = "heading",
                   width: float = 0.5, height: float = 0.1,
                   font_size: float = 48.0) -> dict:
    return {
        "slots": [{
            "id": slot_id,
            "role": role,
            "bounds": {"x": 0.1, "y": 0.1, "width": width, "height": height},
            "typography": {"font_size": font_size, "line_height": 1.0},
        }]
    }


def _slot_plan(copy: str, *, slot_id: str = "headline") -> dict:
    return {
        "schema_version": 1,
        "slides": [{
            "request_id": "slide-01",
            "item_id": "sun.component.example",
            "slots": [{"slot_id": slot_id, "display_copy": copy}],
        }],
    }


def test_slot_content_plan_requires_one_reuse_entry_per_selected_slide() -> None:
    report = {"slides": [{"request_id": "slide-01", "decision": {
        "action": "reuse", "item_id": "sun.component.example"}}]}
    errors = slot_plan.validate_plan({"schema_version": 1, "slides": []}, report, {})
    assert any("missing a reuse entry" in error for error in errors), errors


def test_slot_content_plan_rejects_copy_that_exceeds_native_slot_capacity() -> None:
    report = {"slides": [{"request_id": "slide-01", "decision": {
        "action": "reuse", "item_id": "sun.component.example"}}]}
    contracts = {"sun.component.example": _slot_contract(width=0.12, height=0.04, font_size=28)}
    errors = slot_plan.validate_plan(_slot_plan("A deliberately long sentence that cannot fit this narrow slot."),
                                     report, contracts)
    assert any("exceeds native capacity" in error for error in errors), errors


def test_slot_content_plan_rejects_slot_below_projection_readability_floor() -> None:
    report = {"slides": [{"request_id": "slide-01", "decision": {
        "action": "reuse", "item_id": "sun.component.example"}}]}
    contracts = {"sun.component.example": _slot_contract(role="body", font_size=14)}
    errors = slot_plan.validate_plan(_slot_plan("Readable body", slot_id="headline"), report, contracts)
    assert any("below the projection floor" in error for error in errors), errors


def test_slot_content_plan_accepts_compact_copy_in_readable_native_slot() -> None:
    report = {"slides": [{"request_id": "slide-01", "decision": {
        "action": "reuse", "item_id": "sun.component.example"}}]}
    contracts = {"sun.component.example": _slot_contract(width=0.5, height=0.1, font_size=48)}
    errors = slot_plan.validate_plan(_slot_plan("Build with AI"), report, contracts)
    assert errors == [], errors


# --------------------------------------------------------------------------- #
# component_units — repeatable visual-unit contract
#
# Fixtures are synthetic geometry only. The gate must work from a component's
# own bounds/typography, so no test may branch on a real item id or deck name.
# --------------------------------------------------------------------------- #
def _slot(slot_id: str, x: float, y: float, w: float, h: float, font: float,
          role: str = "body") -> dict:
    return {
        "id": slot_id, "role": role,
        "bounds": {"x": x, "y": y, "width": w, "height": h},
        "typography": {"font_size": font, "line_height": 1.2},
    }


def _card_row_contract(units: int, *, with_chrome: bool = True, body_lines: int = 1) -> dict:
    """A title + footer chrome + `units` congruent cards laid out in a row.

    Each card is a label over `body_lines` body line(s), separated from its
    neighbours by a real gutter — the same grammar as any published card/column/
    step set. Extraction splits a card's wrapped paragraph into one slot per
    drawn line, so `body_lines > 1` is what a dense published card looks like.
    """
    slots: list[dict] = []
    if with_chrome:
        slots += [
            _slot("deck-title", 0.08, 0.15, 0.60, 0.08, 55.0, role="heading"),
            _slot("footer-legal", 0.04, 0.92, 0.21, 0.02, 20.0, role="footer"),
            _slot("page-number", 0.95, 0.92, 0.01, 0.02, 20.0, role="label"),
        ]
    for i in range(units):
        x = 0.08 + i * 0.30
        slots.append(_slot(f"card-{i + 1}-label", x, 0.35, 0.22, 0.04, 32.0, role="label"))
        for line in range(body_lines):
            suffix = "" if line == 0 else f"-{line + 1}"
            slots.append(_slot(f"card-{i + 1}-body{suffix}", x, 0.41 + line * 0.035, 0.22, 0.03, 20.0))
    return {"source": {"view_box": [0, 0, 1920, 1080]}, "slots": slots}


def _unit_plan(item_id: str, slot_ids: list[str]) -> dict:
    return {"schema_version": 1, "slides": [{
        "request_id": "slide-02", "item_id": item_id,
        "slots": [{"slot_id": sid, "display_copy": "Short label"} for sid in slot_ids],
    }]}


def _unit_report(item_id: str) -> dict:
    return {"slides": [{"request_id": "slide-02",
                        "decision": {"action": "reuse", "item_id": item_id}}]}


def test_visual_units_are_inferred_from_geometry_not_slot_names() -> None:
    model = units.unit_model(_card_row_contract(3))
    assert model["primary_unit_count"] == 3, model
    # Title, footer and page number are chrome/singletons: never repeat units.
    grouped = {s["id"] for group in model["groups"] for unit in group["units"] for s in unit}
    assert grouped == {f"card-{i}-{part}" for i in (1, 2, 3) for part in ("label", "body")}, grouped


def test_visual_units_scale_with_the_number_of_repeats_drawn() -> None:
    for drawn in (2, 4, 5):
        assert units.unit_model(_card_row_contract(drawn))["primary_unit_count"] == drawn


def test_plan_gate_rejects_a_blank_card_in_an_engaged_repeat_set() -> None:
    """The observed defect: every mapped slot fits, but card 3 ships blank."""
    contract = _card_row_contract(3)
    errors = slot_plan.validate_plan(
        _unit_plan("sun.component.cards", ["card-1-label", "card-1-body",
                                           "card-2-label", "card-2-body"]),
        _unit_report("sun.component.cards"), {"sun.component.cards": contract})
    assert len(errors) == 1, errors
    assert "slide 'slide-02'" in errors[0], errors
    assert "2 of 3 native unit(s) carry copy" in errors[0], errors
    assert "1 drawn unit(s) would ship blank" in errors[0], errors
    assert "unit 3 [card-3-label, card-3-body]" in errors[0], errors


def test_plan_gate_accepts_a_fully_filled_repeat_set() -> None:
    contract = _card_row_contract(3)
    errors = slot_plan.validate_plan(
        _unit_plan("sun.component.cards", [f"card-{i}-{p}" for i in (1, 2, 3)
                                           for p in ("label", "body")]),
        _unit_report("sun.component.cards"), {"sun.component.cards": contract})
    assert errors == [], errors


def test_plan_gate_still_allows_empty_titles_footers_and_page_numbers() -> None:
    """Non-unit slots stay deliberately empty; only repeats must be finished."""
    contract = _card_row_contract(3)
    errors = slot_plan.validate_plan(
        _unit_plan("sun.component.cards", [f"card-{i}-label" for i in (1, 2, 3)]),
        _unit_report("sun.component.cards"), {"sun.component.cards": contract})
    assert errors == [], errors


def test_plan_gate_leaves_an_untouched_repeat_group_alone() -> None:
    """A component may carry a second grammar the brief never engages."""
    contract = _card_row_contract(3)
    contract["slots"] += [
        _slot("tab-1", 0.08, 0.62, 0.10, 0.03, 24.0, role="label"),
        _slot("tab-2", 0.40, 0.62, 0.10, 0.03, 24.0, role="label"),
    ]
    errors = slot_plan.validate_plan(
        _unit_plan("sun.component.cards", [f"card-{i}-label" for i in (1, 2, 3)]),
        _unit_report("sun.component.cards"), {"sun.component.cards": contract})
    assert errors == [], errors


# --------------------------------------------------------------------------- #
# Repeated-unit readability budget
#
# Physical capacity is not the ceiling inside a repeated card/strip: copy can
# fit every native line and still project as four ragged narrow columns.
# --------------------------------------------------------------------------- #
def test_plan_gate_rejects_a_repeated_card_that_fits_physically_but_reads_dense() -> None:
    """The observed defect: every slot fits, the card is unreadable anyway."""
    contract = _card_row_contract(4, body_lines=4)
    slot_ids = [f"card-{i}-{p}" for i in (1, 2, 3, 4)
                for p in ("label", "body", "body-2", "body-3", "body-4")]
    errors = slot_plan.validate_plan(
        _unit_plan("sun.component.cards", slot_ids),
        _unit_report("sun.component.cards"), {"sun.component.cards": contract})
    # Every slot is individually within native capacity...
    assert not any("exceeds native capacity" in error for error in errors), errors
    # ...but each card carries 5 display lines against a budget of 3.
    assert len(errors) == 4, errors
    assert "slide 'slide-02'" in errors[0], errors
    assert "unit 1 of 4" in errors[0], errors
    assert "5 display lines" in errors[0] and "budget 3" in errors[0], errors
    assert "'card-1-label'" in errors[0] and "'card-1-body-4'" in errors[0], errors
    assert "speaker notes" in errors[0], errors


def test_plan_gate_rejects_a_repeated_label_that_needs_more_than_one_line() -> None:
    contract = _card_row_contract(3)
    # Tall enough that the label physically holds four lines; the budget still
    # rejects a label that needs more than one.
    for slot in contract["slots"]:
        if slot["id"].endswith("-label"):
            slot["bounds"]["height"] = 0.12
    plan = _unit_plan("sun.component.cards", [f"card-{i}-{p}" for i in (1, 2, 3)
                                              for p in ("label", "body")])
    plan["slides"][0]["slots"][0]["display_copy"] = (
        "A label long enough to wrap across the card")
    errors = slot_plan.validate_plan(plan, _unit_report("sun.component.cards"),
                                     {"sun.component.cards": contract})
    assert len(errors) == 1, errors
    assert "label copy fills" in errors[0] and "budget 1" in errors[0], errors
    assert "'card-1-label'" in errors[0], errors


def test_plan_gate_accepts_a_concise_card_row_and_strip() -> None:
    """A 4-card row and a 5-cell strip, one label plus one compact line each."""
    for count in (4, 5):
        contract = _card_row_contract(count)
        errors = slot_plan.validate_plan(
            _unit_plan("sun.component.cards", [f"card-{i}-{p}" for i in range(1, count + 1)
                                               for p in ("label", "body")]),
            _unit_report("sun.component.cards"), {"sun.component.cards": contract})
        assert errors == [], (count, errors)


def test_readability_budget_leaves_non_repeating_long_form_copy_alone() -> None:
    """A long-form body slot is read, not scanned: no unit budget applies."""
    contract = {"source": {"view_box": [0, 0, 1920, 1080]}, "slots": [
        _slot("section-title", 0.08, 0.15, 0.60, 0.10, 55.0, role="heading"),
        _slot("prose", 0.08, 0.35, 0.80, 0.40, 30.0, role="body"),
    ]}
    plan = _unit_plan("sun.component.essay", ["section-title", "prose"])
    plan["slides"][0]["slots"][1]["display_copy"] = (
        "Long-form slides still carry real paragraphs. This one runs several lines and "
        "explains the argument in full, because nothing here repeats and the audience "
        "reads it rather than scanning a row of cards.")
    errors = slot_plan.validate_plan(plan, _unit_report("sun.component.essay"),
                                     {"sun.component.essay": contract})
    assert errors == [], errors


def _unit_selection_report() -> dict:
    return {"slides": [{"request_id": "slide-07", "decision": {
        "action": "reuse", "item_id": "sun.component.steps"}}]}


def test_selection_gate_rejects_repeat_count_that_cannot_host_the_item_count() -> None:
    """The observed defect: a 4-step flow selected for 3 parallel items."""
    errs = vsr._validate_unit_model(
        _unit_selection_report(), True, {"slide-07": 3},
        {"sun.component.steps": _card_row_contract(4)})
    assert len(errs) == 1, errs
    assert "slide 'slide-07'" in errs[0], errs
    assert "needs 3 parallel item(s)" in errs[0] and "repeats 4 native unit(s)" in errs[0], errs
    assert "1 unit(s) would ship blank" in errs[0], errs


def test_selection_gate_accepts_a_matching_repeat_count() -> None:
    assert vsr._validate_unit_model(
        _unit_selection_report(), True, {"slide-07": 3},
        {"sun.component.steps": _card_row_contract(3)}) == []


def test_selection_gate_ignores_components_with_no_repeat_structure() -> None:
    """Unknown unit model is not a mismatch — same rule as the set-of-N guard."""
    cover = {"source": {"view_box": [0, 0, 1920, 1080]}, "slots": [
        _slot("kicker", 0.08, 0.30, 0.20, 0.03, 24.0, role="label"),
        _slot("headline", 0.08, 0.40, 0.70, 0.12, 88.0, role="heading"),
    ]}
    assert vsr._validate_unit_model(_unit_selection_report(), True,
                                    {"slide-07": 3}, {"sun.component.steps": cover}) == []


def _unit_registry(tmp: Path, spec: dict[str, int]) -> Path:
    """A minimal published registry whose items point at real slot contracts."""
    items = []
    for item_id, unit_count in spec.items():
        contract = tmp / f"{item_id}.json"
        contract.write_text(json.dumps(_card_row_contract(unit_count)), encoding="utf-8")
        items.append({"id": item_id, "status": "published",
                      "intent": ["timeline"], "tags": ["timeline"],
                      "paths": {"text_slots": str(contract)}})
    registry = tmp / "registry.json"
    registry.write_text(json.dumps({"items": items}), encoding="utf-8")
    return registry


def test_scorer_skips_top_ranked_incompatible_unit_count_for_next_candidate() -> None:
    """The P1 fallback: unit fit is decided during selection, not after it.

    The best-scoring candidate repeats four units and cannot host three ideas,
    so the scorer must fall through to the lower-scored three-unit component
    rather than emitting a report a later gate can only reject.
    """
    with tempfile.TemporaryDirectory() as td:
        counts = svi.load_unit_profiles(_unit_registry(Path(td), {
            "sun.set.four-step": 4, "sun.set.three-card": 3}))
        assert {k: v["unit_count"] for k, v in counts.items()} == {
            "sun.set.four-step": 4, "sun.set.three-card": 3}, counts

    best = _item(id="sun.set.four-step", content_structure=["a"])
    worse = _item(id="sun.set.three-card", content_structure=[], density="fixed")
    dec, cands = svi.score_request(dict(_req(), item_count=3), [best, worse],
                                   svi.WEIGHTS, None, unit_profiles=counts)

    assert cands[0]["item_id"] == "sun.set.four-step", "incompatible candidate must stay ranked"
    assert cands[0]["score"] > cands[1]["score"], cands
    assert any("Visual-unit fit" in r and "repeats 4 native unit(s)" in r
               for r in cands[0]["reasons"]), cands[0]["reasons"]
    assert cands[0]["retrieval"]["unit_count"] == 4, cands[0]
    assert dec["action"] == "reuse", dec
    assert dec["item_id"] == "sun.set.three-card", dec


def test_scorer_emits_text_only_when_every_candidate_has_the_wrong_unit_count() -> None:
    """No compatible published component means text-only — never a gate failure."""
    with tempfile.TemporaryDirectory() as td:
        counts = svi.load_unit_profiles(_unit_registry(Path(td), {
            "sun.set.four-step": 4, "sun.set.six-card": 6}))

    dec, cands = svi.score_request(
        dict(_req(), item_count=3),
        [_item(id="sun.set.four-step"), _item(id="sun.set.six-card")],
        svi.WEIGHTS, None, unit_profiles=counts)

    assert dec["action"] == "text-only", dec
    assert dec["item_id"] is None and dec["score"] == 0, dec
    assert dec["extraction_recommended"] is True, "text-only fallback keeps extraction evidence"
    assert "visual-unit" in dec["reason"], dec
    assert {c["item_id"] for c in cands} == {"sun.set.four-step", "sun.set.six-card"}, cands


def test_scorer_keeps_components_without_a_repeat_model_compatible() -> None:
    cover = {"source": {"view_box": [0, 0, 1920, 1080]}, "slots": [
        _slot("headline", 0.08, 0.40, 0.70, 0.12, 88.0, role="heading")]}
    with tempfile.TemporaryDirectory() as td:
        contract = Path(td) / "cover.json"
        contract.write_text(json.dumps(cover), encoding="utf-8")
        registry = Path(td) / "registry.json"
        registry.write_text(json.dumps({"items": [{
            "id": "sun.set.cover", "status": "published",
            "paths": {"text_slots": str(contract)}}]}), encoding="utf-8")
        counts = svi.load_unit_profiles(registry)

    assert counts["sun.set.cover"]["unit_count"] is None, counts
    dec, _ = svi.score_request(dict(_req(), item_count=3), [_item(id="sun.set.cover")],
                               svi.WEIGHTS, None, unit_profiles=counts)
    assert dec["action"] == "reuse", dec


def test_scorer_unit_fit_is_inert_below_two_parallel_items() -> None:
    with tempfile.TemporaryDirectory() as td:
        counts = svi.load_unit_profiles(_unit_registry(Path(td), {"sun.set.six-card": 6}))
    for count in (None, 1):
        request = dict(_req())
        if count is not None:
            request["item_count"] = count
        dec, _ = svi.score_request(request, [_item(id="sun.set.six-card")],
                                   svi.WEIGHTS, None, unit_profiles=counts)
        assert dec["action"] == "reuse", (count, dec)


def _quote_panel_contract() -> dict:
    """A two-panel editorial layout: display/quote panel + dense working list.

    Structurally the grammar that hosts one statement, not N parallel items —
    no repeat group, a short slot set in display type, and a much denser
    surface beside it.
    """
    slots = [_slot("deck-title", 0.08, 0.14, 0.50, 0.07, 55.0, role="heading"),
             _slot("quote-line-1", 0.09, 0.50, 0.34, 0.06, 48.0, role="heading"),
             _slot("quote-line-2", 0.09, 0.57, 0.30, 0.06, 48.0, role="heading")]
    for i in range(12):
        slots.append(_slot(f"list-{i + 1}", 0.56, 0.34 + i * 0.04, 0.34, 0.03, 28.0,
                           role="list-item"))
    return {"source": {"view_box": [0, 0, 1920, 1080]}, "slots": slots}


def test_display_surface_detects_a_quote_panel_beside_a_dense_surface() -> None:
    panel = units.display_surface(_quote_panel_contract())
    assert panel is not None
    assert panel["slot_ids"] == ["quote-line-1", "quote-line-2"], panel
    assert panel["font_px"] == 48.0 and panel["body_slot_count"] == 12, panel


def test_display_surface_ignores_an_evenly_repeating_card_layout() -> None:
    """The signal must not fire on a component whose grammar IS parallel."""
    assert units.display_surface(_card_row_contract(4)) is None
    assert units.display_surface(_card_row_contract(2)) is None


def test_display_surface_ignores_headers_that_are_not_display_type() -> None:
    """Column headers are short and non-repeating but stay near body size."""
    contract = _card_row_contract(3)
    contract["slots"].append(_slot("section-note", 0.08, 0.62, 0.20, 0.03, 26.0,
                                   role="label"))
    assert units.display_surface(contract) is None


def test_quote_layout_loses_to_a_matching_multi_item_layout() -> None:
    """The reported defect: a quote-heavy panel chosen for a 4-principle brief."""
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        quote_contract = tmp / "quote.json"
        quote_contract.write_text(json.dumps(_quote_panel_contract()), encoding="utf-8")
        cards_contract = tmp / "cards.json"
        cards_contract.write_text(json.dumps(_card_row_contract(4)), encoding="utf-8")
        registry = tmp / "registry.json"
        registry.write_text(json.dumps({"items": [
            {"id": "sun.set.quote-panel", "status": "published",
             "paths": {"text_slots": str(quote_contract)}},
            {"id": "sun.set.four-card", "status": "published",
             "paths": {"text_slots": str(cards_contract)}},
        ]}), encoding="utf-8")
        profiles = svi.load_unit_profiles(registry)

    # The quote panel narrowly outranks the card set on metadata alone, the way
    # two similarly-tagged published components realistically differ.
    quote = _item(id="sun.set.quote-panel", content_structure=["a"])
    cards = _item(id="sun.set.four-card", content_structure=["a"], density="fixed")
    request = dict(_req(), item_count=4)

    baseline, _ = svi.score_request(request, [quote, cards], svi.WEIGHTS, None)
    assert baseline["item_id"] == "sun.set.quote-panel", "fixture must start mismatched"

    dec, cands = svi.score_request(request, [quote, cards], svi.WEIGHTS, None,
                                   unit_profiles=profiles)
    assert dec["item_id"] == "sun.set.four-card", dec
    penalised = [c for c in cands if c["item_id"] == "sun.set.quote-panel"][0]
    assert any("Layout-grammar fit" in r and "display/quote surface" in r
               for r in penalised["reasons"]), penalised["reasons"]
    assert penalised["retrieval"]["display_surface"]["font_px"] == 48.0, penalised
    assert penalised["score"] == baseline["score"] - svi.DISPLAY_SURFACE_PENALTY, penalised


def test_quote_layout_still_wins_when_nothing_better_is_published() -> None:
    """Penalty, not eligibility: the existing fallback must stay intact."""
    with tempfile.TemporaryDirectory() as td:
        contract = Path(td) / "quote.json"
        contract.write_text(json.dumps(_quote_panel_contract()), encoding="utf-8")
        registry = Path(td) / "registry.json"
        registry.write_text(json.dumps({"items": [{
            "id": "sun.set.quote-panel", "status": "published",
            "paths": {"text_slots": str(contract)}}]}), encoding="utf-8")
        profiles = svi.load_unit_profiles(registry)

    dec, _ = svi.score_request(dict(_req(), item_count=4),
                               [_item(id="sun.set.quote-panel")],
                               svi.WEIGHTS, None, unit_profiles=profiles)
    assert dec["action"] == "reuse", dec
    assert dec["item_id"] == "sun.set.quote-panel", dec
    assert any("layout-grammar mismatch" in w for w in dec["warnings"]), dec["warnings"]


def test_quote_layout_is_untouched_by_a_single_statement_request() -> None:
    """A cover/quote brief wants exactly this grammar — never penalise it."""
    with tempfile.TemporaryDirectory() as td:
        contract = Path(td) / "quote.json"
        contract.write_text(json.dumps(_quote_panel_contract()), encoding="utf-8")
        registry = Path(td) / "registry.json"
        registry.write_text(json.dumps({"items": [{
            "id": "sun.set.quote-panel", "status": "published",
            "paths": {"text_slots": str(contract)}}]}), encoding="utf-8")
        profiles = svi.load_unit_profiles(registry)

    for request in (_req(), dict(_req(), item_count=1)):
        dec, cands = svi.score_request(request, [_item(id="sun.set.quote-panel")],
                                       svi.WEIGHTS, None, unit_profiles=profiles)
        assert dec["action"] == "reuse", dec
        assert dec["warnings"] == [] or not any(
            "layout-grammar" in w for w in dec["warnings"]), dec
        assert not any("Layout-grammar fit" in r for r in cands[0]["reasons"]), cands


def test_scorer_produced_report_passes_the_selection_validator() -> None:
    """End-to-end: the scorer's own fallback output must clear the pre-build gate."""
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        analysis = tmp / "analysis"
        analysis.mkdir()
        unit_registry = _unit_registry(tmp, {"sun.set.four-step": 4, "sun.set.three-card": 3})
        compact = tmp / "compact.json"
        compact.write_text(json.dumps({"items": [
            dict(_item(id="sun.set.four-step"), tags=["timeline"]),
            dict(_item(id="sun.set.three-card"), tags=["timeline"]),
        ]}), encoding="utf-8")
        requests = analysis / "visual-requests.json"
        requests.write_text(json.dumps({"job_id": "unit-fit", "slides": [{
            "request_id": "slide-01", "intent": ["timeline"], "tags": [],
            "content_structure": ["a"], "content_shape": "timeline",
            "density": "medium", "brand": "sun", "item_count": 3,
        }]}), encoding="utf-8")
        report = analysis / "selection-report.json"

        scored = subprocess.run(
            [sys.executable, str(SCRIPTS / "score_visual_items.py"),
             "--batch-request", str(requests), "--registry", str(compact),
             "--unit-registry", str(unit_registry), "--retrieval-index", "none",
             "--output", str(report)],
            capture_output=True, text=True)
        assert scored.returncode == 0, scored.stderr

        decision = json.loads(report.read_text(encoding="utf-8"))["slides"][0]["decision"]
        assert decision["action"] == "reuse", decision
        assert decision["item_id"] == "sun.set.three-card", decision

        validated = subprocess.run(
            [sys.executable, str(SCRIPTS / "validate_selection_report.py"),
             "--selection-report", str(report), "--visual-requests", str(requests),
             "--registry", str(unit_registry)],
            capture_output=True, text=True)
        assert validated.returncode == 0, validated.stdout + validated.stderr
        # Defense in depth actually ran and agreed, rather than being skipped.
        result = json.loads((analysis / "selection-validation.json").read_text(encoding="utf-8"))
        assert result["valid"] is True, result
        lock = [c for c in result["checks"] if c["name"] == "visual_unit_lock"]
        assert lock and lock[0]["pass"] is True, result["checks"]


def test_selection_gate_is_inert_without_a_declared_item_count() -> None:
    for count in (None, 1, "three"):
        assert vsr._validate_unit_model(
            _unit_selection_report(), True, {"slide-07": count},
            {"sun.component.steps": _card_row_contract(6)}) == []


# --------------------------------------------------------------------------- #
# cleanup_run
# --------------------------------------------------------------------------- #
def test_cleanup_keeps_deck_and_pptx() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        run = Path(tmp)
        (run / "deck.html").write_text("<html></html>")
        (run / "out.pptx").write_bytes(b"PK\x03\x04fake")
        (run / "analysis").mkdir()
        (run / "analysis" / "selection-report.json").write_text("{}")
        (run / "slide-1-bg.png").write_bytes(b"png")
        (run / "parity").mkdir()
        (run / "parity" / "x.png").write_bytes(b"png")

        removed = cleanup_run.cleanup(run, dry_run=False)

        assert (run / "deck.html").exists(), "deck.html must survive cleanup"
        assert (run / "out.pptx").exists(), ".pptx must survive cleanup"
        assert (run / "analysis" / "selection-report.json").exists(), "analysis/ must survive"
        assert not (run / "slide-1-bg.png").exists(), "intermediate png must be removed"
        assert not (run / "parity").exists(), "parity/ dir must be removed"
        assert any("parity" in r for r in removed)


# --------------------------------------------------------------------------- #
# export_pptx / compare_renders — cache consistency and AA-aware parity
# --------------------------------------------------------------------------- #
def test_export_invalidates_stale_verdict_and_parity_artifacts() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp) / "export"
        out.mkdir()
        output = out / "deck.pptx"
        output.write_bytes(b"PK\x03\x04current")
        (out / "export-result.json").write_text("{}", encoding="utf-8")
        (out / ".capture-fingerprint.json").write_text("{}", encoding="utf-8")
        (out / ".parity-fingerprint.json").write_text("{}", encoding="utf-8")
        (out / "export-manifest.json").write_text("{}", encoding="utf-8")
        output.with_suffix(".validation.json").write_text("{}", encoding="utf-8")
        report = out / "parity" / "slide-01" / "tier2" / "report.json"
        report.parent.mkdir(parents=True)
        report.write_text("{}", encoding="utf-8")

        export_pptx.invalidate_stale_artifacts(out, output, capture_stale=True)

        assert output.exists(), "invalidation must not delete the built PPTX"
        assert not (out / "export-result.json").exists()
        assert not output.with_suffix(".validation.json").exists()
        assert not (out / ".capture-fingerprint.json").exists()
        assert not (out / ".parity-fingerprint.json").exists()
        assert not (out / "export-manifest.json").exists()
        assert not (out / "parity").exists()


def test_export_reuses_only_complete_fingerprint_bound_parity_reports() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        parity = Path(tmp) / "parity"
        manifest = {"slides": [{"slide": 1}, {"slide": 2}]}
        fingerprint = {"html_sha": "current", "compare_script_sha": "metric-v2"}
        for report in export_pptx.expected_parity_reports(parity, manifest):
            report.parent.mkdir(parents=True, exist_ok=True)
            report.write_text('{"metrics":{"changed_pixel_ratio":0}}', encoding="utf-8")

        export_pptx.write_parity_fingerprint(parity, manifest, fingerprint)

        assert export_pptx.parity_cache_valid(parity, manifest, fingerprint)
        missing = parity / "slide-02" / "tier2" / "report.json"
        missing.unlink()
        assert not export_pptx.parity_cache_valid(parity, manifest, fingerprint)


def test_export_generation_gates_apply_only_to_normal_slide_jobs() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        standalone = root / "standalone.html"
        standalone.write_text("<html/>", encoding="utf-8")
        assert export_pptx.generation_gate_commands(standalone) == []

        run = root / "job" / "run-01"
        analysis = run / "analysis"
        analysis.mkdir(parents=True)
        deck = run / "deck.html"
        deck.write_text("<html/>", encoding="utf-8")
        (analysis / "selection-report.json").write_text("{}", encoding="utf-8")
        commands = export_pptx.generation_gate_commands(deck)

        assert [label for _command, label in commands] == [
            "selection report", "deck asset resolution", "deck stage runtime", "brand compliance", "component fidelity",
        ]
        assert "--visual-requests" not in commands[0][0]
        (analysis / "visual-requests.json").write_text("[]", encoding="utf-8")
        assert "--visual-requests" in export_pptx.generation_gate_commands(deck)[0][0]


def test_deck_stage_starter_is_present_and_recognized() -> None:
    starter = SCRIPTS.parent / "boilerplates" / "deck_stage.js"
    source = starter.read_text(encoding="utf-8")
    assert "customElements.define(TAG, DeckStage)" in source
    assert "goTo(index)" in source
    assert 'observedAttributes()' in source

    html = '<script src="deck_stage.js"></script><deck-stage width="1920" height="1080"></deck-stage>'
    check = importlib.import_module("validate_deck_stage_runtime").check_deck_stage(html)
    assert check["pass"] is True, check


def test_project_python_path_uses_windows_virtualenv_layout() -> None:
    root = Path("C:/slide-plugin")
    assert _common.project_python_path(root, os_name="nt") == (
        root / ".venv" / "Scripts" / "python.exe"
    )


def test_project_python_path_uses_posix_virtualenv_layout() -> None:
    root = Path("/workspace/slide-plugin")
    assert _common.project_python_path(root, os_name="posix") == (
        root / ".venv" / "bin" / "python3"
    )


def test_missing_project_python_has_platform_install_hint() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        try:
            _common.require_project_python(Path(tmp), os_name="nt")
        except _common.ProjectPythonError as exc:
            message = str(exc)
            assert ".venv\\Scripts\\python.exe" in message
            assert "setup.ps1" in message
        else:
            raise AssertionError("missing project virtualenv must fail")


def test_invalid_project_python_has_platform_install_hint() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        python = Path(tmp) / ".venv" / "Scripts" / "python.exe"
        python.parent.mkdir(parents=True)
        python.write_text("not an executable", encoding="utf-8")
        try:
            _common.require_project_python(Path(tmp), os_name="nt")
        except _common.ProjectPythonError as exc:
            message = str(exc)
            assert "not usable" in message
            assert "setup.ps1" in message
        else:
            raise AssertionError("invalid project virtualenv must fail")


def test_preflight_export_and_smoke_select_same_project_python() -> None:
    import check_base_requirements
    import test_export_stack
    catalog_path = SCRIPTS.parents[1] / "slide-system" / "catalog" / "catalog_server.py"
    spec = importlib.util.spec_from_file_location("catalog_server_python_test", catalog_path)
    assert spec and spec.loader
    catalog_server = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(catalog_server)

    expected = _common.require_project_python(SCRIPTS.parents[1])
    assert check_base_requirements.selected_python() == expected
    assert export_pptx.selected_python() == expected
    assert test_export_stack.selected_python() == expected
    assert catalog_server.selected_python() == expected


def test_distribution_surfaces_parse_and_expose_entrypoints() -> None:
    root = SCRIPTS.parents[1]
    plugin = json.loads(
        (root / ".agents/.claude-plugin/plugin.json").read_text(encoding="utf-8")
    )
    marketplace = json.loads(
        (root / ".claude-plugin/marketplace.json").read_text(encoding="utf-8")
    )
    assert plugin["name"] == "sun-riser"
    assert marketplace["plugins"][0]["source"] == "./.agents"

    skill_names = {
        path.parent.name
        for path in (root / ".agents/skills").glob("*/SKILL.md")
    }
    assert {"slide-generator", "component-extractor", "extract-preflight"} <= skill_names

    component_command = (
        root / ".opencode/commands/component.md"
    ).read_text(encoding="utf-8")
    assert "component-extractor" in component_command
    assert "$ARGUMENTS" in component_command


def test_pdf_component_entrypoint_runs_preflight_before_analysis_and_staging() -> None:
    import extract_pdf_components

    calls: list[list[str]] = []

    def runner(cmd: list[str], **_kwargs) -> subprocess.CompletedProcess:
        calls.append([str(value) for value in cmd])
        return subprocess.CompletedProcess(cmd, 0, stdout='{"status":"ok"}', stderr="")

    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        pdf = root / "sample.pdf"
        pdf.write_bytes(b"%PDF-1.4\n%%EOF\n")
        extract_pdf_components.run_workflow(
            pdf=pdf,
            extraction_id="sample-components",
            output_root=root / "extractions",
            history=root / "history.json",
            registry=root / "registry.json",
            catalog_output=root / "catalog-data.json",
            marker=root / "extract-readiness.json",
            python=Path(sys.executable),
            runner=runner,
        )

    scripts = [Path(cmd[1]).name for cmd in calls]
    assert scripts == [
        "check_base_requirements.py",
        "analyze_with_docling.py",
        "auto_stage_candidates.py",
        "build_component_catalog.py",
    ]
    assert calls[0][2:] == ["--input", "pdf", "--json", "--marker", str(root / "extract-readiness.json")]
    assert "publish_extraction.py" not in " ".join(" ".join(cmd) for cmd in calls)


def test_export_smoke_uses_external_markitdown_when_project_module_is_missing() -> None:
    import test_export_stack

    command = test_export_stack.markitdown_command(
        Path("C:/repo/.venv/Scripts/python.exe"),
        module_available=False,
        executable="C:/Tools/markitdown.exe",
        pptx_path=Path("C:/temp/deck.pptx"),
    )
    assert command == ["C:/Tools/markitdown.exe", str(Path("C:/temp/deck.pptx"))]


def test_compare_renders_ignores_small_delta_aa_edges() -> None:
    try:
        from PIL import Image
    except ImportError:
        return

    reference = Image.new("RGB", (400, 400), (120, 120, 120))
    candidate = reference.copy()
    for index in range(5000):
        x = index % 400
        y = index // 400
        candidate.putpixel((x, y), (152, 120, 120))

    metrics = compare_renders.compute_metrics(reference, candidate)

    assert metrics["mean_absolute_error"] < 1.0
    assert metrics["changed_pixel_ratio"] == 0.0


def test_compare_renders_still_fails_shifted_solid_block() -> None:
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        return

    reference = Image.new("RGB", (200, 200), "white")
    candidate = reference.copy()
    ImageDraw.Draw(reference).rectangle((20, 60, 69, 109), fill="black")
    ImageDraw.Draw(candidate).rectangle((100, 60, 149, 109), fill="black")

    metrics = compare_renders.compute_metrics(reference, candidate)

    assert metrics["mean_absolute_error"] > 1.0
    assert metrics["changed_pixel_ratio"] > 0.01


# --------------------------------------------------------------------------- #
# score_visual_items — decision thresholds
# --------------------------------------------------------------------------- #
def _req() -> dict:
    return {"intent": ["timeline"], "tags": [], "content_structure": ["a"],
            "density": "medium", "brand": "sun", "required_exports": []}


def _item(**over) -> dict:
    base = {"id": "sun.set.x", "status": "published", "intent": ["timeline"],
            "tags": [], "content_structure": ["a"], "density": "any",
            "brand": None, "limitations": []}
    base.update(over)
    return base


def test_score_perfect_match_is_reuse() -> None:
    dec, _ = svi.score_request(_req(), [_item()], svi.WEIGHTS, None)
    assert dec["action"] == "reuse", dec
    assert dec["extraction_recommended"] is False


def test_score_safe_mid_rank_is_reused() -> None:
    # semantic 35 + structure 0 + density 4 + brand 10 + export 15 + access 10 = 74
    item = _item(content_structure=[], density="fixed")
    dec, _ = svi.score_request(_req(), [item], svi.WEIGHTS, None)
    assert 65 <= dec["score"] < 75, dec
    assert dec["action"] == "reuse", dec


def test_score_without_buildable_candidate_is_text_only_with_extraction() -> None:
    # Topic/semantic evidence ranks published items, but an incompatible content
    # shape is still a hard physical constraint.
    item = _item(intent=["timeline"], content_structure=["a"])
    dec, _ = svi.score_request(dict(_req(), content_shape="profile"), [item], svi.WEIGHTS, None)
    assert dec["item_id"] is None, dec
    assert dec["score"] == 0, dec
    assert dec["action"] == "text-only", dec
    assert dec["extraction_recommended"] is True, "missing component must recommend extraction"


def test_draft_and_staging_items_never_enter_selection() -> None:
    """Only `published` is selectable — Draft/staging/qa never reach a deck.

    `build_enrichment` already drops non-published records, but that only
    governs the retrieval-index side. Selection eligibility is decided
    independently in `score_request` off the registry item's own `status`, so
    it needs its own proof: an item that would otherwise be a PERFECT match
    (same intent, structure, density, brand as the request) must still fall
    through to text-only purely because it is not published.
    """
    for status in ("staging", "qa", "draft"):
        item = _item(id="sun.component.card-set-grid", status=status)
        dec, cands = svi.score_request(_req(), [item], svi.WEIGHTS, None)
        assert dec["action"] == "text-only", (status, dec)
        assert dec["item_id"] is None, (status, dec)
        assert dec["score"] == 0, (status, dec)
        # It stays visible as ranked evidence for the user, but is not eligible.
        assert cands and cands[0]["eligible"] is False, (status, cands)

    # Control: the same item published IS selected, proving the exclusion above
    # is the status check and not some other mismatch in the fixture.
    published = _item(id="sun.component.card-set-grid", status="published")
    assert svi.score_request(_req(), [published], svi.WEIGHTS, None)[0]["action"] == "reuse"


def test_score_single_generic_semantic_signal_is_reused_when_buildable() -> None:
    item = _item(intent=["layout"], content_structure=["a"])
    dec, _ = svi.score_request(_req(), [item], svi.WEIGHTS, None)
    assert dec["action"] == "reuse", dec
    assert dec["item_id"] == item["id"], dec


# --------------------------------------------------------------------------- #
# score_visual_items — type-intent bias (component vs template in all-types)
# --------------------------------------------------------------------------- #
def _typed_req(query: str, **over) -> dict:
    base = {"query": query, "intent": ["team", "profile", "roster"],
            "tags": ["contributors"], "content_structure": ["heading", "label"],
            "density": "any", "brand": None, "required_exports": []}
    base.update(over)
    return base


def _template_item() -> dict:
    # Matches every request term -> out-scores the component when there is no bias.
    return _item(id="sun.deck.04-contributors", type="template",
                 intent=["team", "profile", "roster"], tags=["contributors"],
                 content_structure=["heading", "label"])


def _component_item() -> dict:
    # One fewer intent match -> lower raw score than the template above.
    return _item(id="sun.component.profile-circles", type="component",
                 intent=["team", "profile"], tags=["contributors"],
                 content_structure=["heading", "label"])


def test_request_type_intent_detection() -> None:
    assert svi.request_type_intent({"prefer_type": "component"}) == "component"
    assert svi.request_type_intent({"prefer_type": "template"}) == "template"
    assert svi.request_type_intent({"query": "reusable component for KPI strip"}) == "component"
    assert svi.request_type_intent({"query": "full slide template for cover"}) == "template"
    # template intent wins ties (explicit whole-slide ask beats an incidental word)
    assert svi.request_type_intent({"query": "component card set full slide template"}) == "template"
    # markers can arrive via intent/tags, not just free text
    assert svi.request_type_intent({"tags": ["icon-reference"]}) == "component"
    assert svi.request_type_intent({"intent": ["team", "profile"]}) is None


def test_type_intent_component_query_ranks_component_over_template() -> None:
    items = [_template_item(), _component_item()]
    # Neutral phrasing: the template out-scores the component (baseline behavior).
    dec_n, _ = svi.score_request(_typed_req("team roster"), items, svi.WEIGHTS, None)
    assert dec_n["item_id"] == "sun.deck.04-contributors", dec_n
    # Explicit component intent: the template is demoted; the component wins.
    dec_c, cands = svi.score_request(
        _typed_req("reusable component team roster"), items, svi.WEIGHTS, None)
    assert dec_c["item_id"] == "sun.component.profile-circles", dec_c
    tmpl = next(c for c in cands if c["item_id"] == "sun.deck.04-contributors")
    assert tmpl["retrieval"]["type_bias"] == "template-demoted"
    assert any("template demoted" in r for r in tmpl["reasons"])


def test_type_intent_template_query_lets_template_win() -> None:
    items = [_template_item(), _component_item()]
    dec, cands = svi.score_request(
        _typed_req("full slide template for the team page"), items, svi.WEIGHTS, None)
    assert dec["item_id"] == "sun.deck.04-contributors", dec
    tmpl = next(c for c in cands if c["item_id"] == "sun.deck.04-contributors")
    assert "type_bias" not in tmpl.get("retrieval", {}), "template intent must not demote templates"


def test_type_intent_neutral_query_applies_no_bias() -> None:
    items = [_template_item(), _component_item()]
    dec, cands = svi.score_request(_typed_req("team roster"), items, svi.WEIGHTS, None)
    # Same winner and score as a run with the demotion path never triggered.
    assert dec["item_id"] == "sun.deck.04-contributors", dec
    for c in cands:
        assert "type_bias" not in c.get("retrieval", {})
        assert not any("template demoted" in r for r in c["reasons"])


def test_type_intent_leaves_components_unchanged() -> None:
    # Component-only scoring never sees a template, so the demotion cannot fire:
    # a component's score is identical with or without component intent.
    comp = [_component_item()]
    _, neutral = svi.score_request(_typed_req("team roster"), comp, svi.WEIGHTS, None)
    _, biased = svi.score_request(_typed_req("reusable component team roster"), comp, svi.WEIGHTS, None)
    assert neutral[0]["score"] == biased[0]["score"]
    assert "type_bias" not in biased[0].get("retrieval", {})


def test_type_intent_no_component_false_positive_when_nothing_fits() -> None:
    # A component-intent query may still reuse a semantically matching published
    # template; it must not be forced onto the unrelated standalone component.
    template = _item(id="sun.deck.10-chart", type="template",
                     intent=["chart", "statistics"], tags=["pie"],
                     content_structure=["metric", "label"])
    unrelated = _item(id="sun.component.timeline", type="component",
                      intent=["timeline"], tags=[], content_structure=["heading"])
    req = {"query": "reusable component financial pie chart",
           "intent": ["chart", "statistics"], "tags": ["pie-chart", "financial"],
           "content_structure": ["metric", "label"], "density": "any",
           "brand": None, "required_exports": []}
    dec, _ = svi.score_request(req, [template, unrelated], svi.WEIGHTS, None)
    assert dec["action"] == "reuse", dec
    assert dec["item_id"] == "sun.deck.10-chart", dec


# --------------------------------------------------------------------------- #
# score_visual_items — hybrid retrieval (v3.2)
# --------------------------------------------------------------------------- #
def _rreq(**over) -> dict:
    base = {"intent": [], "tags": [], "content_structure": [], "density": "any",
            "brand": None, "required_exports": []}
    base.update(over)
    return base


def test_retrieval_secondary_match_lifts_prose_metadata_item() -> None:
    # metric/KPI strip: docling-style prose intent is invisible to primary
    # matching; index keywords must lift it via capped secondary credit.
    strip = _item(id="sun.component.metric-strip",
                  intent=["revenue team size metric strip"], tags=["strip"],
                  content_structure=[])
    req = _rreq(intent=["statistics", "kpi"], tags=["strip"])
    _, plain = svi.score_request(req, [strip], svi.WEIGHTS, None)
    enrichment = svi.build_enrichment([{
        "id": "sun.component.metric-strip", "status": "published",
        "keywords": ["revenue", "team", "metric", "strip"],
        "component_type": "strip", "slot_count": 5,
    }])
    _, enriched = svi.score_request(req, [strip], svi.WEIGHTS, None, enrichment=enrichment)
    assert enriched[0]["score"] > plain[0]["score"], "index keywords must lift the strip"
    assert enriched[0]["retrieval"]["secondary_matches"] == ["statistics"]
    assert enriched[0]["retrieval"]["slot_count"] == 5


def test_retrieval_tier_strip_trap_stays_below_genuine_component() -> None:
    # level/tier strip: an OCR-named trap with a "levels" keyword must not
    # outrank the component that declares tiers/levels as canonical intent.
    genuine = _item(id="sun.component.tier-set",
                    intent=["ranking", "levels", "tiers"], tags=["set-of-3"],
                    content_structure=["heading", "label"])
    trap = _item(id="sun.component.trap-strip",
                 intent=["spicy autocomplete autonomous levels strip"],
                 tags=["strip"], content_structure=[])
    enrichment = svi.build_enrichment([
        {"id": "sun.component.trap-strip", "status": "published",
         "keywords": ["levels", "autonomous"], "slot_count": 16},
    ])
    req = _rreq(intent=["levels", "tiers", "ranking"],
                content_structure=["heading", "label"])
    dec, cands = svi.score_request(req, [trap, genuine], svi.WEIGHTS, None,
                                   enrichment=enrichment)
    assert dec["item_id"] == "sun.component.tier-set", dec
    scores = {c["item_id"]: c["score"] for c in cands}
    assert scores["sun.component.trap-strip"] < scores["sun.component.tier-set"]


def test_retrieval_secondary_only_match_is_reused_when_buildable() -> None:
    # A published component that clears capacity/shape gates may be reused even
    # when its ranking evidence comes only from retrieval metadata.
    lure = _item(id="sun.component.strip-block", intent=["ai team visual"], tags=[],
                 content_structure=["a"])
    enrichment = svi.build_enrichment([{
        "id": "sun.component.strip-block", "status": "published",
        "keywords": ["timeline", "roadmap"], "slot_count": 4,
    }])
    dec, cands = svi.score_request(_req(), [lure], svi.WEIGHTS, None,
                                   enrichment=enrichment)
    cap_points = svi.SECONDARY_CAP * svi.WEIGHTS["semantic_intent"]
    assert cands[0]["criteria"]["semantic_intent"] <= cap_points + 1e-9
    assert dec["action"] == "reuse", dec
    assert dec["item_id"] == "sun.component.strip-block"
    assert dec["extraction_recommended"] is False


def test_retrieval_top_ranked_buildable_candidate_wins() -> None:
    # The scorer ranks all published candidates. Once physical buildability is
    # established, a secondary-only match can win over a lower-scored runner-up.
    lure = _item(
        id="sun.component.strip-block",
        intent=["decorative visual"],
        tags=[],
        content_structure=["heading", "label"],
    )
    good = _item(
        id="sun.component.timeline-row-set",
        intent=["timeline", "roadmap"],
        tags=[],
        content_structure=["heading"],
    )
    enrichment = svi.build_enrichment([{
        "id": "sun.component.strip-block", "status": "published",
        "keywords": ["timeline", "roadmap", "milestones", "schedule"],
        "slot_count": 4,
    }])
    req = _rreq(
        intent=["timeline", "roadmap", "milestones", "schedule"],
        content_structure=["heading", "label"],
    )
    dec, cands = svi.score_request(req, [lure, good], svi.WEIGHTS, None,
                                   enrichment=enrichment)
    assert cands[0]["item_id"] == "sun.component.strip-block"
    assert cands[0]["criteria"]["semantic_intent"] < svi.WEIGHTS["semantic_intent"] * 0.3
    assert dec["item_id"] == "sun.component.strip-block", dec
    assert dec["action"] == "reuse", dec


def test_retrieval_selected_candidate_stays_in_reported_candidates() -> None:
    # The selected top-ranked candidate stays in the report even when top_n
    # truncates other candidate evidence.
    lures = [
        _item(
            id=f"sun.component.strip-block-{idx}",
            intent=["decorative visual"],
            tags=[],
            content_structure=["heading", "label"],
        )
        for idx in range(6)
    ]
    good = _item(
        id="sun.component.timeline-row-set",
        intent=["timeline", "roadmap"],
        tags=[],
        content_structure=["heading"],
    )
    enrichment = svi.build_enrichment([
        {
            "id": lure["id"], "status": "published",
            "keywords": ["timeline", "roadmap", "milestones", "schedule"],
            "slot_count": 4,
        }
        for lure in lures
    ])
    req = _rreq(
        intent=["timeline", "roadmap", "milestones", "schedule"],
        content_structure=["heading", "label"],
    )
    dec, cands = svi.score_request(req, lures + [good], svi.WEIGHTS, None,
                                   top_n=5, enrichment=enrichment)
    assert dec["item_id"].startswith("sun.component.strip-block-"), dec
    assert cands[0]["item_id"].startswith("sun.component.strip-block-"), cands
    assert len(cands) == 5, cands
    assert any(c["item_id"] == dec["item_id"] for c in cands), cands


def test_retrieval_prose_component_outranks_unrelated_item() -> None:
    # team/contributor/profile: prose-only metadata gains capped rank credit,
    # so the right component surfaces above unrelated ones in candidates.
    team = _item(id="sun.component.profile-circles",
                 intent=["team contributor profile circles layout"], tags=[],
                 content_structure=[])
    other = _item(id="sun.component.faq", intent=["faq"], tags=[],
                  content_structure=[])
    enrichment = svi.build_enrichment([
        {"id": "sun.component.profile-circles", "status": "published",
         "name": "Team Contributor Circles",
         "intent": ["team contributor profile circles layout"], "slot_count": 0},
        {"id": "sun.component.faq", "status": "published",
         "keywords": ["faq"], "slot_count": 4},
    ])
    req = _rreq(intent=["team", "profile"], tags=["circles"])
    _, cands = svi.score_request(req, [other, team], svi.WEIGHTS, None,
                                 enrichment=enrichment)
    assert cands[0]["item_id"] == "sun.component.profile-circles", cands
    assert cands[0]["retrieval"]["secondary_matches"], "must explain the lexical match"


def test_retrieval_anti_use_case_penalty_for_undeclared_domain() -> None:
    badge = _item(id="sun.component.badge-grid", intent=["numbered", "grid"],
                  tags=["cards"])
    req = _rreq(intent=["statistics"], tags=["numbered"], content_structure=["a"])
    base_enr = {"id": "sun.component.badge-grid", "status": "published",
                "keywords": ["badge"], "slot_count": 6}
    _, plain = svi.score_request(req, [badge], svi.WEIGHTS, None,
                                 enrichment=svi.build_enrichment([base_enr]))
    anti_enr = dict(base_enr, anti_use_cases=[
        "Do not use for data-driven charts or metrics; placeholder diagram."])
    _, hit = svi.score_request(req, [badge], svi.WEIGHTS, None,
                               enrichment=svi.build_enrichment([anti_enr]))
    assert plain[0]["score"] - hit[0]["score"] == svi.ANTI_USE_CASE_PENALTY
    assert hit[0]["retrieval"]["anti_hits"] == ["statistics"]
    assert any("Anti-use-case" in r for r in hit[0]["reasons"])


def test_retrieval_anti_hit_on_declared_intent_is_caveat_not_exclusion() -> None:
    # The item declares statistics as honest intent; its anti text mentioning
    # "metrics" is an editing caveat and must NOT be penalized.
    circles = _item(id="sun.component.circle-panel", intent=["statistics", "ranking"],
                    tags=[])
    enrichment = svi.build_enrichment([{
        "id": "sun.component.circle-panel", "status": "published",
        "anti_use_cases": ["Do not reuse the baked metrics without editing text slots."],
        "slot_count": 13,
    }])
    req = _rreq(intent=["statistics"], content_structure=["a"])
    _, cands = svi.score_request(req, [circles], svi.WEIGHTS, None,
                                 enrichment=enrichment)
    assert "anti_hits" not in cands[0].get("retrieval", {}), cands[0]
    assert not any("Anti-use-case" in r for r in cands[0]["reasons"])


def test_retrieval_count_fit_penalty_prefers_matching_set_size() -> None:
    # buildability: wrong declared set size must not beat a better-fit item.
    three = _item(id="sun.component.three", intent=["roles"],
                  tags=["cards", "set-of-3"])
    four = _item(id="sun.component.quad-card-set", intent=["roles"],
                 tags=["cards", "set-of-4"])
    req = _rreq(intent=["roles"], tags=["cards"], content_structure=["a"],
                item_count=4)
    dec, cands = svi.score_request(req, [three, four], svi.WEIGHTS, None)
    assert dec["item_id"] == "sun.component.quad-card-set", dec
    three_cand = next(c for c in cands if c["item_id"] == "sun.component.three")
    assert three_cand["retrieval"]["set_sizes"] == [3]
    assert any("Count fit" in r for r in three_cand["reasons"])


def test_retrieval_zero_slot_component_penalized_when_text_needed() -> None:
    deco = _item(id="sun.component.ring-panel", intent=["team", "profile"], tags=[])
    slotted = _item(id="sun.component.slot-panel", intent=["team", "profile"], tags=[])
    enrichment = svi.build_enrichment([
        {"id": "sun.component.ring-panel", "status": "published", "slot_count": 0},
        {"id": "sun.component.slot-panel", "status": "published", "slot_count": 6},
    ])
    req = _rreq(intent=["team", "profile"], content_structure=["a"])
    dec, cands = svi.score_request(req, [deco, slotted], svi.WEIGHTS, None,
                                   enrichment=enrichment)
    assert dec["item_id"] == "sun.component.slot-panel", dec
    deco_cand = next(c for c in cands if c["item_id"] == "sun.component.ring-panel")
    assert deco_cand["retrieval"]["slot_count"] == 0
    assert any("no editable text slots" in r for r in deco_cand["reasons"])
    # decoration-only requests (no text content) are NOT penalized
    _, cands2 = svi.score_request(_rreq(intent=["team", "profile"]),
                                  [deco, slotted], svi.WEIGHTS, None,
                                  enrichment=enrichment)
    deco2 = next(c for c in cands2 if c["item_id"] == "sun.component.ring-panel")
    assert not any("no editable text slots" in r for r in deco2["reasons"])


def test_retrieval_zero_slot_only_candidate_becomes_text_only() -> None:
    deco = _item(id="sun.component.ring-panel", intent=["team", "profile"], tags=[])
    enrichment = svi.build_enrichment([
        {"id": "sun.component.ring-panel", "status": "published", "slot_count": 0},
    ])
    req = _rreq(intent=["team", "profile"], content_structure=["heading", "label"])

    dec, cands = svi.score_request(req, [deco], svi.WEIGHTS, None, enrichment=enrichment)

    assert dec["action"] == "text-only", dec
    assert dec["item_id"] is None, dec
    assert dec["extraction_recommended"] is True, dec
    assert cands[0]["item_id"] == "sun.component.ring-panel", cands
    assert any("no editable text slots" in reason for reason in cands[0]["reasons"])


def test_scorer_uses_slot_contract_when_retrieval_index_omits_slot_count() -> None:
    """The published text-slot contract, not index completeness, decides capacity."""
    deco = _item(id="sun.component.ring-panel", intent=["team", "profile"], tags=[])
    req = _rreq(intent=["team", "profile"], content_structure=["heading", "label"])

    dec, cands = svi.score_request(
        req,
        [deco],
        svi.WEIGHTS,
        None,
        enrichment={},
        unit_profiles={"sun.component.ring-panel": {"editable_slot_count": 0}},
    )

    assert dec["action"] == "text-only", dec
    assert dec["item_id"] is None, dec
    assert cands[0]["retrieval"]["slot_count"] == 0, cands[0]
    assert any("no editable text slots" in reason for reason in cands[0]["reasons"])


def test_canonicalize_drops_filler_symmetrically_with_index_tokens() -> None:
    """Both sides of the semantic comparison must use one filler rule.

    `_field_tokens` has always dropped STOPWORDS; `_canonicalize` did not, so a
    request carrying prose connectors was scored against a larger denominator
    than the index side ever contributed to.
    """
    assert svi._semantic_terms(["timeline", "of", "a", "the", "with"]) == {"timeline"}
    assert svi._semantic_terms(["of"]) == set()
    # Real vocabulary is still folded, not dropped.
    assert svi._semantic_terms(["roadmap"]) == {"timeline"}
    # content_structure keeps using the unfiltered path: slot names are not prose.
    assert svi._canonicalize(["a"]) == {"a"}


def test_count_mismatch_is_not_buildable_even_when_top_ranked() -> None:
    """A set-of-3 cannot host 4 items; that is capacity, not preference."""
    three = _item(id="sun.component.card-set-grid", intent=["timeline"],
                  tags=["set-of-3"], content_structure=["a"])
    enrichment = svi.build_enrichment([
        {"id": "sun.component.card-set-grid", "status": "published", "slot_count": 6},
    ])
    dec, cands = svi.score_request(dict(_req(), item_count=4), [three], svi.WEIGHTS,
                                   None, enrichment=enrichment)
    assert dec["action"] == "text-only", dec
    assert cands[0]["item_id"] == "sun.component.card-set-grid", cands
    assert dec["extraction_recommended"] is True, dec

    fits, _ = svi.score_request(dict(_req(), item_count=3), [three], svi.WEIGHTS,
                                None, enrichment=enrichment)
    assert fits["action"] == "reuse", fits


def test_noisy_intent_scores_identically_to_canonical_intent() -> None:
    """The headline contract: prose in `intent` must not cost coverage.

    `semantic` is a ratio, so every extra word divides the score. Dropping
    STOPWORDS was not enough — plain nouns like "real"/"case" are not filler
    but are still not vocabulary, and they were halving a perfect match.
    """
    item = _item(id="sun.set.timeline", intent=["timeline"], tags=[],
                 content_structure=["a"])
    canonical = _req()
    noisy = dict(canonical, intent=["timeline", "of", "a", "real", "case"])

    dec_c, cands_c = svi.score_request(canonical, [item], svi.WEIGHTS, None)
    dec_n, cands_n = svi.score_request(noisy, [item], svi.WEIGHTS, None)

    assert cands_n[0]["criteria"]["semantic_intent"] == cands_c[0]["criteria"]["semantic_intent"]
    assert dec_n["score"] == dec_c["score"], (dec_n, dec_c)
    assert dec_n["item_id"] == dec_c["item_id"] == "sun.set.timeline"
    # Dropped prose stays observable rather than vanishing.
    assert dec_n["evidence"]["dropped_intent_terms"] == ["a", "case", "of", "real"]
    assert dec_n["evidence"]["scored_intent_terms"] == ["timeline"]
    assert any("non-canonical prose" in w for w in dec_n["warnings"])


def test_normalize_intent_keeps_only_canonical_vocabulary() -> None:
    """`intent` is canonical-only; literal tags belong in `tags`.

    Corpus membership is not an escape hatch — a junk word survives it whenever
    any single published item happens to carry that string, which would make
    the no-dilution guarantee probabilistic instead of exact.
    """
    kept, dropped = svi.normalize_intent(["roadmap", "set-of-3", "kombucha", "what"])
    assert kept == {"timeline"}, kept
    assert dropped == ["kombucha", "set-of-3", "what"], dropped
    # `tags` is still matched literally, so a real tag is never lost.
    assert "set-of-3" in svi._semantic_terms(["set-of-3"])


def test_content_structure_still_matches_literally() -> None:
    """Intent normalization must not leak into slot-name matching."""
    assert svi.overlap_score(["a", "of"], ["a", "of"]) == 1.0
    assert svi._canonicalize(["a"]) == {"a"}


def test_source_specific_artwork_is_selectable_with_topic_warning() -> None:
    """Topic mismatch is disclosed, but does not block published reuse."""
    themed = _item(id="sun.goal-setting-2026.01-cover", intent=["timeline"], tags=[])
    dec, cands = svi.score_request(_req(), [themed], svi.WEIGHTS, None)
    assert dec["action"] == "reuse", dec
    assert dec["item_id"] == "sun.goal-setting-2026.01-cover"
    assert dec["evidence"]["subject_warnings"] == ["sun.goal-setting-2026.01-cover"]
    assert any("Subject mismatch" in r for r in cands[0]["reasons"]), cands[0]
    assert any("Subject mismatch" in warning for warning in dec["warnings"]), dec
    assert dec["extraction_recommended"] is False


def test_generic_shell_and_on_topic_request_stay_selectable() -> None:
    """Negative cases: the guard must not block legitimate reuse."""
    # A generic shell has no topic tokens at all.
    shell = _item(id="sun.component.card-set-grid", intent=["timeline"], tags=[])
    dec, _ = svi.score_request(_req(), [shell], svi.WEIGHTS, None)
    assert dec["action"] == "reuse", dec

    # Placeholder-named artwork is still a generic shell.
    lorem = _item(id="sun.component.lorem-ipsum-badge-set", intent=["timeline"], tags=[])
    assert svi.subject_tokens("sun.component.lorem-ipsum-badge-set") == set()
    assert svi.score_request(_req(), [lorem], svi.WEIGHTS, None)[0]["action"] == "reuse"

    # A themed item IS reusable when the deck is about that topic. The support
    # goes in `query`: adding it to `tags` would dilute the semantic ratio and
    # test the floor instead of the subject guard.
    themed = _item(id="sun.goal-setting-2026.01-cover", intent=["timeline"], tags=[])
    on_topic = dict(_req(), query="goal setting review deck")
    assert svi.score_request(on_topic, [themed], svi.WEIGHTS, None)[0]["action"] == "reuse"


def test_subject_tokens_read_the_index_name_not_only_the_id() -> None:
    """The index `name` is where a themed capture announces its source."""
    record = {"name": "01 - Cover: Goal Setting 2026"}
    assert svi.subject_tokens("sun.deck.01-cover", record) >= {"goal", "setting"}
    assert svi.subject_tokens("sun.deck.01-cover", {"name": "Cover Slide"}) == set()


def test_scorer_keeps_higher_ranked_subject_mismatch_with_warning() -> None:
    """The topic signal informs review but does not suppress reuse."""
    themed = _item(id="sun.goal-setting-2026.01-cover", intent=["timeline"],
                   tags=["roadmap"], content_structure=["a"])
    shell = _item(id="sun.component.card-set-grid", intent=["timeline"], tags=[])
    dec, _ = svi.score_request(_req(), [themed, shell], svi.WEIGHTS, None)
    assert dec["action"] == "reuse", dec
    assert dec["item_id"] == "sun.goal-setting-2026.01-cover", dec
    assert dec["evidence"]["subject_warnings"] == ["sun.goal-setting-2026.01-cover"]


def test_rejected_items_are_persisted_for_reproducible_reruns() -> None:
    good = _item(id="sun.component.card-set-grid", intent=["timeline"], tags=[])
    bad = _item(id="sun.component.grid-set-card", intent=["timeline"], tags=["roadmap"])
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        (root / "reg.json").write_text(json.dumps({"items": [bad, good]}), encoding="utf-8")
        (root / "req.json").write_text(json.dumps(_req()), encoding="utf-8")
        out = root / "out.json"
        rc = subprocess.run(
            [sys.executable, str(SCRIPTS / "score_visual_items.py"),
             "--request", str(root / "req.json"), "--registry", str(root / "reg.json"),
             "--retrieval-index", "none", "--output", str(out),
             "--reject-item", "sun.component.grid-set-card"],
            capture_output=True, text=True)
        assert rc.returncode == 0, rc.stderr
        report = json.loads(out.read_text(encoding="utf-8"))
        assert report["rejected_items"] == ["sun.component.grid-set-card"], report
        assert report["decision"]["item_id"] == "sun.component.card-set-grid"

        plain = json.loads((root / "req.json").read_text(encoding="utf-8"))
        assert plain  # sanity
        out2 = root / "out2.json"
        subprocess.run(
            [sys.executable, str(SCRIPTS / "score_visual_items.py"),
             "--request", str(root / "req.json"), "--registry", str(root / "reg.json"),
             "--retrieval-index", "none", "--output", str(out2)],
            capture_output=True, text=True, check=True)
        assert json.loads(out2.read_text(encoding="utf-8"))["rejected_items"] == []


def test_superseded_pdfs_quarantined_but_foreign_ones_fail_delivery() -> None:
    """One canonical PDF per job — without deleting files we did not create."""
    with tempfile.TemporaryDirectory() as tmp:
        run_dir = Path(tmp)
        out_dir = run_dir / "_export"
        out_dir.mkdir()
        canonical = run_dir / "deck.pdf"
        canonical.write_bytes(b"%PDF-1.4 canonical")
        old = run_dir / "deck-v1.pdf"
        old.write_bytes(b"%PDF-1.4 old")
        (out_dir / export_pptx.PDF_HISTORY).write_text(json.dumps(["deck-v1.pdf"]),
                                                       encoding="utf-8")
        stranger = run_dir / "someone-elses.pdf"
        stranger.write_bytes(b"%PDF-1.4 foreign")

        quarantined, foreign = export_pptx.quarantine_superseded_pdfs(
            run_dir, canonical, out_dir)

        assert len(quarantined) == 1 and quarantined[0].endswith("deck-v1.pdf")
        assert not old.exists(), "our own superseded artifact is moved aside"
        assert (run_dir / export_pptx.SUPERSEDED_DIR / "deck-v1.pdf").exists()
        assert foreign == ["someone-elses.pdf"], foreign
        assert stranger.exists(), "a file we did not create must never be touched"
        assert canonical.exists()


def test_every_shape_lock_key_has_scorer_vocabulary() -> None:
    """SHAPE_TYPE_MAP and SYNONYMS must not drift apart.

    A content_shape the validator accepts but the scorer has no canonical token
    for can never clear the semantic floor, which looks like a library gap and
    is really a vocabulary gap.
    """
    smap = svi._build_synonym_map()
    missing = sorted(shape for shape in vsr.SHAPE_TYPE_MAP if shape not in smap)
    assert not missing, f"content_shape(s) with no canonical scorer vocabulary: {missing}"


def test_semantic_score_ranks_candidates_without_a_reuse_floor() -> None:
    """Semantic evidence orders reusable published candidates; it does not gate them."""
    lure = _item(id="sun.component.low-semantic", intent=["unrelated"], tags=[])
    dec, _ = svi.score_request(_req(), [lure], svi.WEIGHTS, None)
    assert dec["action"] == "reuse", dec


def test_pdf_geometry_reads_pages_and_landscape_orientation() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        pdf = Path(tmp) / "d.pdf"
        pdf.write_bytes(
            b"%PDF-1.4\n/Type /Pages\n"
            b"/Type /Page /MediaBox [ 0 0 1440 810 ]\n"
            b"/Type /Page /MediaBox [ 0 0 1440 810 ]\n")
        pages, (w, h) = export_pptx.pdf_geometry(pdf)
        assert pages == 2, pages
        assert (w, h) == (1440.0, 810.0), (w, h)
        assert w > h, "16:9 deck pages must be landscape"


def test_export_pdf_does_not_combine_landscape_with_explicit_paper() -> None:
    """Chromium swaps the paper box when both are given — that was the P0."""
    source = (SCRIPTS / "export-pdf.js").read_text(encoding="utf-8")
    pdf_call = source[source.index("await page.pdf({"):]
    pdf_call = pdf_call[:pdf_call.index("});")]
    assert "width:" in pdf_call and "height:" in pdf_call, pdf_call
    assert "landscape" not in pdf_call, "explicit width/height already fix orientation"


def test_export_pdf_locks_the_stage_transform_for_print() -> None:
    source = (SCRIPTS / "export-pdf.js").read_text(encoding="utf-8")
    assert 'setProperty("transform", "none", "important")' in source, (
        "deck-stage fit() re-runs on the print resize and must not win")


# --------------------------------------------------------------------------- #
# build_hybrid_pptx wrapping contract — PPTX-only card-text overlap
# --------------------------------------------------------------------------- #
CARD_CANVAS_W, CARD_CANVAS_H = 1920.0, 1080.0
# Two 380px-wide cards 40px apart: the real narrow-card geometry from the
# readability-budget run, where a wrapped 3-line body was exported as one line.
CARD_BODY = "Claude Code: viết và review code, debug, tự động hoá dự án."
CARD_LEFT = {"x": 200.0, "y": 500.0, "w": 380.0, "h": 99.0, "z": 10,
             "text": CARD_BODY, "fontSize": "25px", "lineHeight": "33px",
             "color": "rgb(23,23,23)", "align": "start"}
CARD_RIGHT = dict(CARD_LEFT, x=620.0, text="Draft nội dung, brief slide, chỉnh tone bài viết.")
CARD_HEADING = {"x": 200.0, "y": 420.0, "w": 380.0, "h": 48.0, "z": 9,
                "text": "Chọn công cụ", "fontSize": "40px", "lineHeight": "48px",
                "color": "rgb(23,23,23)", "align": "start"}
RICH_LEAD = {
    "x": 96.0, "y": 240.0, "w": 1560.0, "h": 81.0, "z": 3,
    "text": "Claude Desktop app phù hợp mọi đối tượng — Chat · Cowork · Code.",
    "fontSize": "30px", "lineHeight": "40.5px", "color": "rgb(23,23,23)",
    "fontWeight": "400", "align": "start",
    "runs": [
        {"text": "Claude Desktop app phù hợp mọi đối tượng", "fontSize": "30px",
         "fontWeight": "700", "color": "rgb(255,85,51)"},
        {"text": " — Chat · Cowork · Code.", "fontSize": "30px",
         "fontWeight": "400", "color": "rgb(23,23,23)"},
    ],
}


def _export_text_items(items: list[dict], tmp: Path, name: str = "cards.pptx"):
    """Run the real builder + a save/reload round-trip.

    Returns (pptx_path, shapes) read back from the saved file, so assertions see
    the actual exported PPTX objects rather than in-memory python-pptx state.
    """
    import build_hybrid_pptx as hybrid
    from pptx import Presentation

    prs = Presentation()
    prs.slide_width = _pptx_inches(13.333333)
    prs.slide_height = _pptx_inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for item in items:
        hybrid.add_text_box_v2(slide, item, "Arial", CARD_CANVAS_W, CARD_CANVAS_H)
    out = tmp / name
    prs.save(str(out))
    return out, list(Presentation(str(out)).slides)[0].shapes


def _pptx_inches(value: float):
    from pptx.util import Inches
    return Inches(value)


def _rendered_width_pt(shape) -> float:
    """Widest single line this shape would paint if it never wrapped."""
    import build_hybrid_pptx as hybrid
    widest = 0.0
    for para in shape.text_frame.paragraphs:
        text = "".join(r.text for r in para.runs).strip()
        size = max((r.font.size.pt for r in para.runs if r.font.size), default=0)
        widest = max(widest, len(text) * size * hybrid.AVERAGE_GLYPH_WIDTH)
    return widest


def test_pptx_card_body_wraps_instead_of_overrunning_its_neighbour() -> None:
    """Long card copy must stay inside its own box, not paint over the next card.

    The browser wrapped this string to 3 lines inside a 380px card; exported
    without word wrap PowerPoint lays it on ONE line and, because native PPTX
    text is not clipped to its shape, it runs straight across the adjacent
    card's text box. Browser/PDF parity cannot see that.
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        _, shapes = _export_text_items([CARD_HEADING, CARD_LEFT, CARD_RIGHT], Path(tmpdir))
        left = next(s for s in shapes if CARD_BODY[:20] in s.text_frame.text)
        right = next(s for s in shapes if "Draft nội dung" in s.text_frame.text)

        assert left.text_frame.word_wrap is True, "browser-wrapped body must wrap in PPTX"
        # Geometry is untouched: the box still is the card, and stays clear of
        # its neighbour, so a wrapping box can no longer reach it.
        assert left.left < right.left, (left.left, right.left)
        assert left.left + left.width <= right.left, "card boxes must not overlap"
        expected_w = CARD_LEFT["w"] / CARD_CANVAS_W * 13.333333
        assert abs(left.width / 914400 - expected_w) < 0.01, left.width
        assert abs(left.left / 914400 - CARD_LEFT["x"] / CARD_CANVAS_W * 13.333333) < 0.01
        assert abs(left.top / 914400 - CARD_LEFT["y"] / CARD_CANVAS_H * 7.5) < 0.01

        # Height budget: room for the wrapped lines, still on the canvas.
        assert left.height >= right.height > 0
        assert left.height / 914400 >= 2 * 33.0 / CARD_CANVAS_H * 7.5, left.height
        assert left.top + left.height <= _pptx_inches(7.5), "box must stay on the slide"


def test_pptx_inline_rich_text_exports_as_one_editable_box() -> None:
    """A lead with a bold coloured prefix must not become overlapping boxes."""
    with tempfile.TemporaryDirectory() as tmpdir:
        _, shapes = _export_text_items([RICH_LEAD], Path(tmpdir), "rich-lead.pptx")
        assert len(shapes) == 1, "all inline runs belong to the same PowerPoint textbox"
        paragraph = shapes[0].text_frame.paragraphs[0]
        assert [run.text for run in paragraph.runs] == [
            "Claude Desktop app phù hợp mọi đối tượng", " — Chat · Cowork · Code."
        ]
        assert str(paragraph.runs[0].font.color.rgb) == "FF5533"
        assert paragraph.runs[0].font.bold is True
        assert str(paragraph.runs[1].font.color.rgb) == "171717"
        assert paragraph.runs[1].font.bold is False


def test_pptx_short_heading_stays_one_unwrapped_line() -> None:
    """A one-line heading must not be re-wrapped by PowerPoint's own metrics."""
    with tempfile.TemporaryDirectory() as tmpdir:
        _, shapes = _export_text_items([CARD_HEADING, CARD_LEFT], Path(tmpdir))
        heading = next(s for s in shapes if "Chọn công cụ" in s.text_frame.text)

        assert heading.text_frame.word_wrap is False, (
            "a heading the browser fitted on one line must not be allowed to re-wrap")
        assert len(heading.text_frame.paragraphs) == 1
        # Not wrapping is only safe because the line genuinely fits its box.
        assert _rendered_width_pt(heading) <= heading.width / 914400 * 72.0, (
            "a non-wrapping heading must fit its own box width")


def test_export_gate_catches_a_non_wrapping_line_wider_than_its_box() -> None:
    """The post-export geometry gate must fail overflowing native text."""
    with tempfile.TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir)
        good, _ = _export_text_items([CARD_HEADING, CARD_LEFT, CARD_RIGHT], tmp)
        failures: list[str] = []
        summary = export_objects.check_text_overflow(good, failures)
        assert failures == [], failures
        assert summary["text_boxes_checked"] == 3 and summary["overflowing"] == 0, summary

        # Same copy, but declared as a single captured line: the builder keeps
        # word_wrap off, so the gate must report the overrun.
        broken, _ = _export_text_items([dict(CARD_LEFT, h=33.0)], tmp, "overflow.pptx")
        failures = []
        summary = export_objects.check_text_overflow(broken, failures)
        assert summary["overflowing"] == 1, summary
        assert failures and "render over neighbouring shapes" in failures[0], failures


def test_export_result_declares_the_editability_tier() -> None:
    """`editable PPTX` must not be read as `every shape is editable`."""
    layered = export_pptx.EDITABILITY_TIERS["layered"]
    assert layered["tier"] == "text-editable", layered
    assert "background" in layered["graphics"], layered
    assert layered["limitation"], "the limitation must be spelled out for the reader"


def test_fidelity_report_publishes_coverage_ratio() -> None:
    report = {"slides": [
        {"request_id": "s1", "decision": {"action": "reuse", "item_id": "sun.set.a"}},
        {"request_id": "s2", "decision": {"action": "text-only", "item_id": None}},
        {"request_id": "s3", "decision": {"action": "text-only", "item_id": None}},
    ]}
    checked = fidelity.check_fidelity("<html></html>", report, {"items": []})
    assert len(checked) == 1, checked
    assert len(list(fidelity._decisions(report))) == 3


def test_fidelity_accepts_a_validated_planned_subset_of_native_slots() -> None:
    """Unused component slots must stay empty when the run plan does not need them."""
    with tempfile.TemporaryDirectory() as tmp:
        preview = Path(tmp) / "preview.html"
        preview.write_text(
            '<div class="slot" data-slot-id="headline"></div>'
            '<div class="slot" data-slot-id="unused-detail"></div>',
            encoding="utf-8",
        )
        report = {"slides": [{"request_id": "slide-01", "decision": {
            "action": "reuse", "item_id": "sun.component.example"}}]}
        registry = {"items": [{"id": "sun.component.example", "paths": {"preview": str(preview)}}]}
        plan = {"slides": [{"request_id": "slide-01", "item_id": "sun.component.example",
                             "slots": [{"slot_id": "headline", "display_copy": "Short copy"}]}]}
        checked = fidelity.check_fidelity(
            '<div class="bg"></div><div data-slot-id="headline">Short copy</div>',
            report, registry, plan,
        )
        assert checked[0]["pass_"] is True, checked
        assert checked[0]["coverage"] == 1.0, checked
        assert checked[0]["coverage_scope"] == "planned-slots", checked


def test_export_validator_requires_svg_blip_only_when_vector_source_exists() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        (root / "present.svg").write_text("<svg/>", encoding="utf-8")
        overlays = [
            {"id": "present", "vector_source": "present.svg"},
            {"id": "missing", "vector_source": "missing.svg"},
            {"id": "effects", "vector_source": "present.svg", "css_effects": True},
        ]
        assert export_objects.expected_svg_blips(overlays, root) == 1


def test_deck_asset_gate_rejects_missing_relative_visuals_but_ignores_data_urls() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        html = root / "deck.html"
        html.write_text(
            '<img src="assets/missing.svg"><div style="background:url(data:image/png;base64,abc)"></div>',
            encoding="utf-8",
        )
        errors = deck_assets.missing_local_assets(html)
        assert errors == ["assets/missing.svg"], errors
        (root / "assets").mkdir()
        (root / "assets" / "missing.svg").write_text("<svg/>", encoding="utf-8")
        assert deck_assets.missing_local_assets(html) == []


# --------------------------------------------------------------------------- #
# Render legibility: reused components must render readably, not just match ids
# --------------------------------------------------------------------------- #
def _txt(x, y, w, h, text="copy", color="rgb(23, 23, 23)"):
    return {"tag": "h3", "text": text, "x": x, "y": y, "w": w, "h": h, "color": color}


def test_wrapped_text_colliding_with_the_next_slot_fails() -> None:
    # Real slide 6 geometry: 28px copy that wrapped to 2 lines (h 56) inside a
    # slot pitched 36px apart, so it lands on top of the bullet below it.
    collisions = fidelity.find_text_collisions([
        _txt(1107.7, 576.7, 663.6, 56.0, "Thông tin khách hàng, hợp đồng"),
        _txt(1106.4, 612.7, 342.0, 28.0, "Phân vân thì hỏi lead trước."),
    ])
    assert collisions, "overflowing copy overlapping the next slot must be reported"
    assert collisions[0]["ratio"] >= fidelity.TEXT_OVERLAP_MIN_RATIO


def test_nested_and_neighbouring_text_boxes_are_not_collisions() -> None:
    # A parent box enclosing its own child (DOM nesting) and two stacked
    # single-line slots that merely touch are both legitimate layouts.
    assert fidelity.find_text_collisions([
        _txt(72, 353, 683, 549, "BẠN ĐANG TỐN THỜI GIAN"),
        _txt(72, 512, 650, 274, "VÀO VIỆC GÌ?"),
    ]) == []
    assert fidelity.find_text_collisions([
        _txt(403.1, 500.5, 345.7, 19.0),
        _txt(421.1, 521.5, 367.1, 19.0),
    ]) == []


def test_white_copy_on_warm_paper_fails_contrast() -> None:
    # Slides 4/5: the component artwork never rendered, leaving white text on
    # the #FFFDF8 paper background — a ~1.03:1 ratio.
    assert fidelity.contrast_ratio((255, 255, 255), (255, 253, 248)) < fidelity.CONTRAST_MIN
    assert fidelity.contrast_ratio((23, 23, 23), (255, 253, 248)) >= fidelity.CONTRAST_MIN
    assert fidelity.parse_css_color("rgb(255, 253, 248)") == (255, 253, 248)
    assert fidelity.parse_css_color("#FFFDF8") == (255, 253, 248)


def test_object_outside_the_canvas_is_reported() -> None:
    slide = {"slide": 4, "objects": [
        {"id": "page-04-obj-01", "bounds": {"x": 60, "y": -1491, "w": 593, "h": 593}},
        {"id": "page-04-obj-04", "bounds": {"x": 224, "y": 40, "w": 1525, "h": 550}},
    ]}
    off = fidelity.off_canvas_objects(slide, 1920, 1080)
    assert [o["id"] for o in off] == ["page-04-obj-01"], off


def test_render_legibility_rejects_the_reproduced_deck_defects() -> None:
    manifest = {"canvasW": 1920, "canvasH": 1080, "slides": [
        {"slide": 4, "objects": [
            {"id": "page-04-obj-01", "bounds": {"x": 60, "y": -1491, "w": 593, "h": 593}}],
         "text": [_txt(293.8, 483.0, 241.0, 61.1, "Hỏi đáp,", "rgb(255, 255, 255)")]},
        {"slide": 6, "objects": [], "text": [
            _txt(1107.7, 576.7, 663.6, 56.0, "Thông tin khách hàng"),
            _txt(1106.4, 612.7, 342.0, 28.0, "Phân vân thì hỏi lead"),
        ]},
    ]}
    report = fidelity.check_render_legibility(manifest, renders_dir=None)
    assert report["valid"] is False
    kinds = {f["check"] for f in report["failures"]}
    assert "off_canvas_object" in kinds, report["failures"]
    assert "text_collision" in kinds, report["failures"]


def test_render_legibility_passes_a_clean_slide() -> None:
    # The headline band clears the artwork band: no text-vs-text overlap and no
    # text standing on the component's illustration.
    manifest = {"canvasW": 1920, "canvasH": 1080, "slides": [
        {"slide": 1, "objects": [
            {"id": "obj-01", "bounds": {"x": 100, "y": 300, "w": 400, "h": 300}}],
         "text": [_txt(72, 129, 430, 55, "AI LÀM ĐƯỢC GÌ"),
                  _txt(72, 189, 481, 55, "CHO TỪNG NGƯỜI")]},
    ]}
    assert fidelity.check_render_legibility(manifest, renders_dir=None)["valid"] is True


# --------------------------------------------------------------------------- #
# Placement contract: text vs component ARTWORK (not just text vs text).
#
# `find_text_collisions` compares text against text, so a generated caption
# dropped over a component's illustration passed every legibility check while
# being unreadable on the slide. Bounds here are synthetic but proportioned
# like the real defect: a caption band directly above a circle-badge set whose
# artwork reaches up into it.
# --------------------------------------------------------------------------- #
def _circle_band(prefix: str = "page-02") -> list[dict]:
    return [{"id": f"{prefix}-obj-0{i + 1}",
             "bounds": {"x": 57 + i * 615, "y": 433, "w": 571, "h": 570}}
            for i in range(3)]


def test_caption_intersecting_component_artwork_fails() -> None:
    hits = fidelity.find_text_over_artwork(
        {"slide": 4, "objects": _circle_band(),
         "text": [_txt(168, 348, 380, 132, "Hỏi đáp, brainstorm, soạn thảo — việc nhanh")]})
    assert hits, "a caption reaching into the component's artwork must be reported"
    assert hits[0]["placement"] == "external"
    assert hits[0]["object_id"] == "page-02-obj-01"
    assert hits[0]["ratio"] >= fidelity.TEXT_ARTWORK_MAX_RATIO


def test_chrome_title_clear_of_the_artwork_passes() -> None:
    # Declared slide chrome is not exempt — it passes by clearing the artwork.
    chrome = _txt(96, 130, 1728, 83, "CHỌN ĐÚNG APP CHO ĐÚNG VIỆC")
    chrome["placement"] = "chrome"
    assert fidelity.find_text_over_artwork(
        {"slide": 4, "objects": _circle_band(), "text": [chrome]}) == []


def test_native_slot_inside_its_own_component_artwork_passes() -> None:
    # The component drew this box for this copy; the slot plan governs it.
    slot = _txt(197, 699, 109, 42, "CHAT", "rgb(255, 255, 255)")
    slot["slotId"] = "lorem-ipsum"
    assert fidelity.find_text_over_artwork(
        {"slide": 4, "objects": _circle_band(), "text": [slot]}) == []


def test_undeclared_text_is_treated_as_external_not_exempt() -> None:
    # A missing attribute must never buy an exemption from the artwork check.
    assert fidelity.text_placement({"text": "x"}) == "external"
    assert fidelity.text_placement({"text": "x", "placement": "nonsense"}) == "external"
    assert fidelity.text_placement({"text": "x", "slotId": "01"}) == "slot"


def test_render_legibility_rejects_the_reproduced_caption_over_circles() -> None:
    # Slide 4 as reported: three captions above a reused circle-badge set whose
    # artwork reaches up into them, with the component's own slots filled.
    slots = []
    for x, text in ((197, "CHAT"), (813, "COWORK"), (1428, "CODE")):
        slot = _txt(x, 699, 109, 42, text, "rgb(255, 255, 255)")
        slot["slotId"] = text.lower()
        slots.append(slot)
    captions = [_txt(168 + i * 616, 348, 380, 132, f"caption {i + 1}") for i in range(3)]
    manifest = {"canvasW": 1920, "canvasH": 1080, "slides": [
        {"slide": 4, "objects": _circle_band(), "text": slots + captions}]}
    report = fidelity.check_render_legibility(manifest, renders_dir=None)
    assert report["valid"] is False
    over = [f for f in report["failures"] if f["check"] == "text_over_artwork"]
    assert len(over) == 3, report["failures"]
    assert {f["id"] for f in over} == {"page-02-obj-01", "page-02-obj-02", "page-02-obj-03"}
    # Evidence names the slide, the text and the overlay, and rules out z-index.
    assert all(f["slide"] == 4 and "caption" in f["detail"] for f in over)
    assert all("z-index" in f["detail"] for f in over)


def test_export_runs_render_legibility_on_the_capture_manifest() -> None:
    source = (SCRIPTS / "export_pptx.py").read_text(encoding="utf-8")
    assert "--export-manifest" in source, \
        "capture output must be gated for legibility before the PPTX is built"
    assert "VALIDATE_FIDELITY" in source


def test_retrieval_enrichment_published_only_and_missing_index() -> None:
    # Draft/staging records never enrich scoring, even from a stale file, and
    # a missing index file degrades to plain primary-only scoring.
    assert svi.build_enrichment([
        {"id": "sun.component.x", "status": "staging", "keywords": ["timeline"]},
        {"status": "published", "keywords": ["timeline"]},
    ]) == {}
    with tempfile.TemporaryDirectory() as tmp:
        assert svi.load_retrieval_index(Path(tmp) / "missing.jsonl") == {}


def test_retrieval_corrupt_index_degrades_to_empty_enrichment() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        index = Path(tmp) / "component-retrieval-index.jsonl"
        index.write_text(
            '{"id":"sun.component.valid","status":"published","keywords":["kpi"]}\n'
            '{"id":"sun.component.truncated",',
            encoding="utf-8",
        )
        assert svi.load_retrieval_index(index) == {}
        index.write_bytes(b"\xff\xfe")
        assert svi.load_retrieval_index(index) == {}


def test_retrieval_index_projects_slot_count() -> None:
    import build_component_retrieval_index as bri
    registry = {"items": [{
        "id": "sun.component.slots", "status": "published", "type": "component",
        "intent": ["grid"], "text_contract": {"slot_count": 7},
    }]}
    records = bri.build_records(registry)
    assert records[0]["slot_count"] == 7
    assert records[0]["schema_version"] == 2


# --------------------------------------------------------------------------- #
# validate_selection_report
# --------------------------------------------------------------------------- #
def _single_report(item_id="sun.interview-workshop-sunriser.02-timeline", score=80.0,
                   action="reuse", scores=(80.0, 80.0)) -> dict:
    crit = {k: 1.0 for k in vsr.REQUIRED_CRITERIA}
    return {
        "request_id": "s1",
        "generated_at": "x",
        "generated_by": "score_visual_items.py",
        "decision": {
            "action": action,
            "item_id": item_id,
            "score": score,
            "reason": "r",
            "extraction_recommended": False,
        },
        "candidates": [{"item_id": f"c{i}", "eligible": True, "score": s, "criteria": crit}
                       for i, s in enumerate(scores)],
    }


def test_equal_scores_are_plausible() -> None:
    checks, errors, warnings = [], [], []
    vsr._validate_single(_single_report(scores=(48.0, 48.0)), checks, errors, warnings, errors.append)
    plaus = next(c for c in checks if c["name"] == "score_plausibility")
    assert plaus["pass"] is True, "two equal scores must NOT fail plausibility"


def test_eligible_all_zero_still_fails() -> None:
    checks, errors, warnings = [], [], []
    vsr._validate_single(_single_report(scores=(0.0, 0.0)), checks, errors, warnings, errors.append)
    plaus = next(c for c in checks if c["name"] == "score_plausibility")
    assert plaus["pass"] is False, "eligible items all scoring 0 must fail"


def test_selection_report_accepts_text_only_with_suggestions() -> None:
    report = _single_report(item_id=None, score=0.0, action="text-only", scores=(32.0, 21.0))
    report["decision"]["extraction_recommended"] = True
    checks, errors, warnings = [], [], []
    vsr._validate_single(report, checks, errors, warnings, errors.append)
    assert not errors, errors
    assert next(check for check in checks if check["name"] == "adoption_compliance")["pass"] is True


def test_selection_report_rejects_text_only_with_component_or_score() -> None:
    report = _single_report(item_id="sun.component.timeline", score=25.0, action="text-only")
    checks, errors, warnings = [], [], []
    vsr._validate_single(report, checks, errors, warnings, errors.append)
    assert any("text-only action must have item_id null" in error for error in errors), errors
    assert any("text-only action must have score 0" in error for error in errors), errors


# --------------------------------------------------------------------------- #
# T1 parallel-set allowance — shape lock reads declared repeated-item evidence
# --------------------------------------------------------------------------- #
def _set_request(shape: str, n: int, *, declare: bool = True, count: bool = True) -> dict:
    request = {"content_shape": shape, "intent": [shape],
               "content_structure": ["label", "heading", "body"]}
    if declare:
        request["content_structure"].append(f"repeatable-set-of-{n}")
    if count:
        request["item_count"] = n
    return request


def _set_item(n: int, *, terms: set[str] | None = None) -> tuple[set[str], set[int]]:
    """A component that names its grammar, not the request's shape label."""
    return (terms if terms is not None else {"role-cards", "cards", "personas"},
            {n})


def test_parallel_set_allowance_accepts_a_matching_repeated_component() -> None:
    """The reported root cause: a 4-card set excluded from a 4-item checklist."""
    terms, sizes = _set_item(4)
    assert svi.shape_lock_ok("checklist", _set_request("checklist", 4), terms, sizes)


def test_parallel_set_allowance_still_accepts_a_literal_checklist() -> None:
    """The base rule is untouched: a real checklist matches on its own tokens."""
    assert svi.shape_lock_ok("checklist", _set_request("checklist", 4),
                             {"checklist", "preparation"}, set())


def test_parallel_set_allowance_rejects_a_count_mismatch() -> None:
    terms, sizes = _set_item(5)
    assert not svi.shape_lock_ok("checklist", _set_request("checklist", 4), terms, sizes)


def test_parallel_set_allowance_needs_declared_repeat_evidence() -> None:
    """item_count alone must not reclassify an ordinary request."""
    terms, sizes = _set_item(4)
    assert not svi.shape_lock_ok("checklist", _set_request("checklist", 4, declare=False),
                                 terms, sizes)
    # ...and a stray tag without item_count is not evidence either.
    assert not svi.shape_lock_ok("checklist", _set_request("checklist", 4, count=False),
                                 terms, sizes)


def test_parallel_set_allowance_never_applies_to_single_statement_shapes() -> None:
    """A repeated set does not make a card set a cover or a closing."""
    for shape in ("cover", "closing", "two-column"):
        terms, sizes = _set_item(4)
        assert not svi.shape_lock_ok(shape, _set_request(shape, 4), terms, sizes), shape


def test_parallel_set_allowance_rejects_a_component_declaring_no_set() -> None:
    assert not svi.shape_lock_ok("checklist", _set_request("checklist", 4),
                                 {"role-cards"}, set())


def test_real_ai_workflow_checklist_slides_select_count_compatible_components() -> None:
    """End-to-end on the real library with the real AI-workflow requests.

    Slides 6 and 8 must land on published components that genuinely repeat 4 and
    5 units, rather than the quote-heavy prep template or a sparse CTA.
    """
    registry = SCRIPTS.parent / "registries" / "visual-library-compact.json"
    if not registry.is_file() or not REGISTRY.is_file():
        return
    items = [i for i in read_text_slots.load_json(registry).get("items", [])
             if i.get("status") == "published"]
    profiles = svi.load_unit_profiles(REGISTRY)
    enrichment = svi.load_retrieval_index(
        SCRIPTS.parent / "registries" / "component-retrieval-index.jsonl")

    real_requests = [
        {"request_id": "slide-06-four-principles",
         "query": "reusable component: four numbered usage principles checklist",
         "prefer_type": "component", "intent": ["checklist"], "content_shape": "checklist",
         "tags": ["numbered", "principles", "checklist", "set-of-4", "rules"],
         "content_structure": ["label", "heading", "body", "repeatable-set-of-4"],
         "item_count": 4, "density": "medium", "brand": "sun-studio"},
        {"request_id": "slide-08-pro-tips",
         "query": "reusable component: five numbered pro tips list with small tags",
         "prefer_type": "component", "intent": ["checklist"], "content_shape": "checklist",
         "tags": ["tips", "numbered", "list", "checklist", "set-of-5"],
         "content_structure": ["label", "heading", "body", "repeatable-set-of-5"],
         "item_count": 5, "density": "high", "brand": "sun-studio"},
    ]
    for request in real_requests:
        decision, _ = svi.score_request(request, items, svi.WEIGHTS, None,
                                        enrichment=enrichment, unit_profiles=profiles)
        wanted = request["item_count"]
        assert decision["action"] == "reuse", (request["request_id"], decision)
        chosen = decision["item_id"]
        assert profiles.get(chosen, {}).get("unit_count") == wanted, (
            f"{request['request_id']}: {chosen} repeats "
            f"{profiles.get(chosen, {}).get('unit_count')} units, needs {wanted}")
        # The previously-selected mismatches must not come back.
        assert chosen not in {"sun.interview-workshop-sunriser.05-prep",
                              "sun.sun-presentation.08-next-steps-cta"}, chosen


def test_shape_lock_matches_and_mismatches() -> None:
    reg_tokens = vsr._registry_tokens(read_text_slots.load_json(REGISTRY))
    rep = _single_report(item_id="sun.interview-workshop-sunriser.02-timeline")
    # match: timeline shape -> timeline item
    errs, _ = vsr._validate_shape_lock(rep, False, {"s1": {"content_shape": "timeline"}}, reg_tokens, strict_shape=False)
    assert not errs, f"timeline->timeline should pass: {errs}"
    # mismatch: cover shape -> timeline item
    errs, _ = vsr._validate_shape_lock(rep, False, {"s1": {"content_shape": "cover"}}, reg_tokens, strict_shape=False)
    assert errs, "cover shape locked to a timeline item must fail"


def test_missing_shape_warns_unless_strict() -> None:
    reg_tokens = vsr._registry_tokens(read_text_slots.load_json(REGISTRY))
    rep = _single_report(item_id="sun.interview-workshop-sunriser.02-timeline")
    errs, warns = vsr._validate_shape_lock(rep, False, {"s1": {"content_shape": None}}, reg_tokens, strict_shape=False)
    assert not errs and warns, "missing shape is a warning by default"
    errs, warns = vsr._validate_shape_lock(rep, False, {"s1": {"content_shape": None}}, reg_tokens, strict_shape=True)
    assert errs and not warns, "missing shape is an error under --strict-shape"


def test_shape_lock_covers_component_first_shapes() -> None:
    reg_tokens = vsr._registry_tokens(read_text_slots.load_json(REGISTRY))
    cases = {
        "profile": "sun.component.team-contributor-circles.g01",
        "tiers": "sun.component.spicy-autocomplete-autonomous-levels-strip",
        "icons": "sun.component.brand-icon-reference-sheet",
        "review": "sun.goal-setting-2026.07-quarterly-check-in",
    }
    for shape, item_id in cases.items():
        rep = _single_report(item_id=item_id)
        errs, _ = vsr._validate_shape_lock(rep, False, {"s1": {"content_shape": shape}}, reg_tokens, strict_shape=True)
        assert not errs, f"{shape} -> {item_id} should pass: {errs}"
    rep = _single_report(item_id="sun.salary-benefits-2026.01-cover")
    errs, _ = vsr._validate_shape_lock(rep, False, {"s1": {"content_shape": "tiers"}}, reg_tokens, strict_shape=True)
    assert errs, "tiers shape locked to a cover template must fail"


def test_selection_report_rejects_non_contract_custom_action() -> None:
    """A scorer result cannot be relabeled as a local custom build."""
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        report_path = root / "analysis" / "selection-report.json"
        report_path.parent.mkdir()
        report_path.write_text(json.dumps({
            "job_id": "test-job",
            "generated_at": "2026-07-13T00:00:00+00:00",
            "generated_by": "score_visual_items.py",
            "slides": [{
                "request_id": "s1",
                "decision": {
                    "action": "custom-local",
                    "item_id": None,
                    "score": 80.0,
                    "reason": "Manually rejected despite a strong match.",
                },
                "candidates": [{
                    "item_id": "sun.component.timeline",
                    "eligible": True,
                    "score": 80.0,
                    "criteria": {**{key: 1.0 for key in vsr.REQUIRED_CRITERIA},
                                 "semantic_intent": 20.0},
                }],
            }],
        }), encoding="utf-8")
        original_argv = sys.argv
        try:
            sys.argv = ["validate_selection_report.py", "--selection-report", str(report_path)]
            assert vsr.main() == 1, "manual curation must fail the blocking selection gate"
        finally:
            sys.argv = original_argv


def test_slide_generator_requires_fresh_selection_for_new_jobs() -> None:
    skill = (SCRIPTS.parent.parent / ".agents" / "skills" / "slide-generator" / "SKILL.md")
    text = skill.read_text(encoding="utf-8")
    assert "Do not read `docs/logs/`" in text
    assert "outputs/slide-jobs/` to judge current library fit" in text
    assert "never edit `selection-report.json`" in text


# --------------------------------------------------------------------------- #
# scaffold_slide_from_component
# --------------------------------------------------------------------------- #
def test_scaffold_preserves_slots_no_base64() -> None:
    preview = scaffold._preview_path(ITEM_WITH_SLOTS, str(REGISTRY))
    slots = scaffold._extract_slots(preview.read_text(encoding="utf-8", errors="replace"))
    assert slots, "expected positioned slots in this component"
    frag = scaffold.build_scaffold(ITEM_WITH_SLOTS, slots)
    assert frag.count("data-slot-id=") >= len(slots), "every slot id must survive"
    assert "base64" not in frag, "scaffold must NOT embed the raster SVG"
    assert 'class="bg"' in frag, "scaffold must include a .bg placeholder"
    assert ".slide-scaffold .slot { z-index: 20; }" in frag, \
        "editable slot copy must remain above decomposed artwork overlays"


def test_scaffold_falls_back_to_text_slot_contract_when_preview_has_no_slots() -> None:
    contract = {
        "source": {"view_box": [0, 0, 1000, 500]},
        "slots": [{
            "id": "headline",
            "html_tag": "h2",
            "bounds": {"x": 0.1, "y": 0.2, "width": 0.5, "height": 0.1},
            "typography": {
                "font_family": "Arial",
                "font_size": 30,
                "font_weight": "bold",
                "font_style": "normal",
                "line_height": 1.1,
                "color": "#ffffff",
            },
            "horizontal_align": "left",
        }],
    }

    slots = scaffold._slots_from_contract(contract)
    fragment = scaffold.build_scaffold("sun.component.example", slots)

    assert len(slots) == 1
    assert 'data-slot-id="headline"' in fragment
    assert "left:10.0000%" in fragment
    assert "font-size:64.8000px" in fragment


def test_scaffold_rejects_compact_registry() -> None:
    # A registry whose items lack `paths` (the compact shape) must be rejected
    # with a clear error rather than silently producing nothing.
    import json
    with tempfile.TemporaryDirectory() as tmp:
        compact = Path(tmp) / "compact.json"
        compact.write_text(json.dumps({"items": [{"id": ITEM_WITH_SLOTS}]}))
        try:
            scaffold._preview_path(ITEM_WITH_SLOTS, str(compact))
        except SystemExit as exc:
            assert "paths.preview" in str(exc) or "full" in str(exc).lower()
        else:
            raise AssertionError("registry without paths must be rejected")


# --------------------------------------------------------------------------- #
# validate_component_fidelity
# --------------------------------------------------------------------------- #
def _registry_dict() -> dict:
    return read_text_slots.load_json(REGISTRY)


def test_fidelity_pass_and_fail() -> None:
    reg = _registry_dict()
    preview = scaffold._preview_path(ITEM_WITH_SLOTS, str(REGISTRY))
    slots = scaffold._extract_slots(preview.read_text(encoding="utf-8", errors="replace"))
    deck_ok = scaffold.build_scaffold(ITEM_WITH_SLOTS, slots)  # full coverage + .bg
    report = {"slides": [{"request_id": "s1",
                          "decision": {"action": "reuse", "item_id": ITEM_WITH_SLOTS}}]}

    res = fidelity.check_fidelity(deck_ok, report, reg)
    assert res and res[0]["pass_"] is True, res

    res = fidelity.check_fidelity("<html>nothing</html>", report, reg)
    assert res and res[0]["pass_"] is False, "deck with no slots/bg must fail fidelity"


# --------------------------------------------------------------------------- #
# read_text_slots
# --------------------------------------------------------------------------- #
def test_read_text_slots_projection() -> None:
    # Resolve the slots fixture from the registry (not a hardcoded path) so a
    # pruned/renamed item can never leave this test bound to a ghost folder —
    # the exact failure mode that made guideline-board-layouts break it.
    registry = read_text_slots.load_json(REGISTRY)
    entry = next(i for i in registry["items"] if i["id"] == ITEM_WITH_SLOTS)
    slots_path = SCRIPTS.parents[1] / entry["paths"]["text_slots"]
    data = read_text_slots.load_json(slots_path)
    slim = read_text_slots.project(data["slots"], with_typography=False)
    assert len(slim) == len(data["slots"])
    assert set(slim[0].keys()) == set(read_text_slots.SLIM_FIELDS)


# --------------------------------------------------------------------------- #
# crop_svg_region
# --------------------------------------------------------------------------- #
import json as _json
import crop_svg_region as crop
import validate_text_slots as vts


def _crop_fixture(tmp: Path, region: dict) -> Path:
    item = tmp / "items" / "metric-card"
    (item / "artifact").mkdir(parents=True)
    (item / "artifact" / "visual.svg").write_text(
        '<?xml version="1.0"?>\n<svg xmlns="http://www.w3.org/2000/svg" '
        'viewBox="0 0 1000 600" width="1000" height="600">'
        '<defs><clipPath id="c1"><rect width="1000" height="600"/></clipPath></defs>'
        '<rect x="500" y="60" width="400" height="240" fill="#f60"/></svg>',
        encoding="utf-8",
    )
    (item / "artifact" / "text-slots.json").write_text(_json.dumps({
        "slots": [
            {"id": "in", "bounds": {"x": 0.55, "y": 0.12, "width": 0.30, "height": 0.10}, "z_order": 1},
            {"id": "out", "bounds": {"x": 0.05, "y": 0.80, "width": 0.10, "height": 0.05}, "z_order": 2},
        ],
        "source": {"view_box": [0, 0, 1000, 600], "canvas_width": 1000, "canvas_height": 600},
    }), encoding="utf-8")
    (item / "mapping.json").write_text(_json.dumps(
        {"item_id": "metric-card", "type": "component", "source": {"region": region}}), encoding="utf-8")
    return item


def test_crop_region_rewrites_viewbox_and_slots() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        item = _crop_fixture(Path(tmp), {"x": 0.5, "y": 0.1, "width": 0.4, "height": 0.4, "unit": "normalized"})
        res = crop.crop_item(item)
        assert res["status"] == "cropped" and res["slots_dropped"] == 1, res
        svg = (item / "artifact" / "visual.svg").read_text()
        assert 'viewBox="0 0 400 240"' in svg and "translate(-500.0 -60.0)" in svg, svg
        slots = crop.load_json(item / "artifact" / "text-slots.json")
        assert [s["id"] for s in slots["slots"]] == ["in"], "out-of-region slot must drop"
        b = slots["slots"][0]["bounds"]
        assert abs(b["x"] - 0.125) < 1e-6 and abs(b["width"] - 0.75) < 1e-6, b
        assert slots["source"]["region_crop"]["crop_window"] == [500.0, 60.0, 400.0, 240.0]


def test_crop_region_idempotent_and_full_page_noop() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        item = _crop_fixture(Path(tmp), {"x": 0.5, "y": 0.1, "width": 0.4, "height": 0.4, "unit": "normalized"})
        crop.crop_item(item)
        assert crop.crop_item(item)["status"] == "already-cropped"
    with tempfile.TemporaryDirectory() as tmp:
        item = _crop_fixture(Path(tmp), {"x": 0, "y": 0, "width": 1, "height": 1, "unit": "normalized"})
        assert crop.crop_item(item)["status"] == "full-page-noop"


def test_crop_region_honors_absolute_units() -> None:
    # On the 1000x600 fixture page, a pt region must crop identically to the
    # equivalent normalized region (regression: pt was treated as a 0-1 fraction
    # and silently produced a ~1000x-too-large viewBox with every slot dropped).
    with tempfile.TemporaryDirectory() as tmp:
        item = _crop_fixture(Path(tmp), {"x": 500, "y": 60, "width": 400, "height": 240, "unit": "pt"})
        res = crop.crop_item(item)
        assert res["status"] == "cropped" and res["slots_dropped"] == 1, res
        svg = (item / "artifact" / "visual.svg").read_text()
        assert 'viewBox="0 0 400 240"' in svg, svg
        slots = crop.load_json(item / "artifact" / "text-slots.json")
        assert slots["source"]["region_crop"]["crop_window"] == [500.0, 60.0, 400.0, 240.0]
    # an unsupported unit must fail loud, never silently mis-scale.
    with tempfile.TemporaryDirectory() as tmp:
        item = _crop_fixture(Path(tmp), {"x": 1, "y": 1, "width": 1, "height": 1, "unit": "furlongs"})
        try:
            crop.crop_item(item)
        except SystemExit:
            pass
        else:
            raise AssertionError("unknown region unit must raise SystemExit")


def test_validate_excludes_cropped_out_source_text() -> None:
    # After a region crop, source text outside the region has no slot. The
    # full-page source-with-text.svg must NOT report it as unmapped, because
    # crop_svg_region.py recorded it in source.region_crop.dropped_source_refs.
    def _build(item: Path, with_marker: bool) -> None:
        (item / "artifact").mkdir(parents=True)
        (item / "evidence").mkdir(parents=True)
        (item / "artifact" / "visual.svg").write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
            '<rect width="100" height="100"/></svg>', encoding="utf-8")
        (item / "evidence" / "source-with-text.svg").write_text(
            '<svg xmlns="http://www.w3.org/2000/svg"><text>IN</text>'
            '<text>OUT</text></svg>', encoding="utf-8")
        source = {}
        if with_marker:
            source["region_crop"] = {"dropped_source_refs": [
                {"text_index": 1, "tspan_index": 0, "character_range": [0, 3]}]}
        (item / "artifact" / "text-slots.json").write_text(_json.dumps({
            "schema_version": 1,
            "slots": [{"id": "s1", "editable": True, "allow_empty": True,
                       "bounds": {"x": 0.1, "y": 0.1, "width": 0.2, "height": 0.1},
                       "source_refs": [{"text_index": 0, "tspan_index": 0,
                                        "character_range": [0, 2]}]}],
            "source": source,
        }), encoding="utf-8")

    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "card"
        _build(item, with_marker=True)
        assert vts.validate(item) == [], "cropped-out text must be excluded"
    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "card"
        _build(item, with_marker=False)
        errs = vts.validate(item)
        assert any("Unmapped source text" in e for e in errs), errs


# --------------------------------------------------------------------------- #
# decompose_svg_objects — ancestor transforms and off-canvas geometry
# --------------------------------------------------------------------------- #
import decompose_svg_objects as decompose
import xml.etree.ElementTree as _ET2

_SVG_NS = "{http://www.w3.org/2000/svg}"
_INK_NS = "{http://www.inkscape.org/namespaces/inkscape}"


def _cropped_component_svg() -> "_ET2.Element":
    # Shape of a real cropped extraction: the drawing sits under a layer that a
    # crop offset translated into view. Measuring in the browser sees the
    # translated position; copying the group alone does not.
    return _ET2.fromstring(
        '<svg xmlns="http://www.w3.org/2000/svg" width="1999" height="620" '
        'viewBox="0 0 1999 620">'
        '<defs><clipPath id="c"><path d="M0 0H2938V2623H0Z"/></clipPath></defs>'
        '<g transform="translate(-499.89 -1859.3)">'
        '  <g xmlns:ns="http://www.inkscape.org/namespaces/inkscape" '
        '     ns:groupmode="layer">'
        '    <g id="circle-a"><circle cx="998" cy="1900" r="271"/></g>'
        '    <g id="circle-b"><circle cx="1400" cy="1900" r="271"/></g>'
        '  </g>'
        '</g></svg>')


def test_decompose_carries_ancestor_transforms_into_fragments() -> None:
    root = _cropped_component_svg()
    groups = decompose.document_groups(root)
    assert [g.get("id") for g in groups] == ["circle-a", "circle-b"], groups
    parent_map = {c: p for p in root.iter() for c in p}
    for group in groups:
        assert decompose.ancestor_transform(root, parent_map, group) == \
            "translate(-499.89 -1859.3)", (
                "a fragment that drops the crop offset paints outside its own "
                "viewBox and renders blank")


def test_decompose_reads_the_viewbox_and_rejects_cropped_away_geometry() -> None:
    root = _cropped_component_svg()
    assert decompose.viewbox(root) == (0.0, 0.0, 1999.0, 620.0)
    # y = -1491 is the real measurement for the circle set that rendered blank.
    assert not decompose.intersects_canvas(
        {"x": 60, "y": -1491, "w": 593, "h": 593}, 1999.0, 620.0)
    assert decompose.intersects_canvas(
        {"x": 224, "y": 40, "w": 1525, "h": 550}, 1999.0, 620.0)
    # Straddling the top edge is still partly visible — keep it.
    assert decompose.intersects_canvas(
        {"x": 10, "y": -100, "w": 200, "h": 200}, 1999.0, 620.0)


# --------------------------------------------------------------------------- #
# crop_svg_region — off-canvas <image> pruning
# --------------------------------------------------------------------------- #
import xml.etree.ElementTree as _ET

_SVG = "{http://www.w3.org/2000/svg}"


def _crop_image_fixture(tmp: Path, body: str, defs: str = "") -> Path:
    # 1000x600 page; region -> crop window page-space [500, 60, 400, 240],
    # i.e. the page-space rectangle x:500..900, y:60..300.
    item = tmp / "items" / "img-card"
    (item / "artifact").mkdir(parents=True)
    (item / "artifact" / "visual.svg").write_text(
        '<?xml version="1.0"?>\n<svg xmlns="http://www.w3.org/2000/svg" '
        'viewBox="0 0 1000 600" width="1000" height="600">'
        f'<defs>{defs}</defs>{body}</svg>', encoding="utf-8")
    (item / "artifact" / "text-slots.json").write_text(_json.dumps({
        "slots": [],
        "source": {"view_box": [0, 0, 1000, 600], "canvas_width": 1000, "canvas_height": 600},
    }), encoding="utf-8")
    (item / "mapping.json").write_text(_json.dumps(
        {"item_id": "img-card", "type": "component",
         "source": {"region": {"x": 0.5, "y": 0.1, "width": 0.4, "height": 0.4, "unit": "normalized"}}}),
        encoding="utf-8")
    return item


def _img_ids(item: Path) -> list[str]:
    root = _ET.parse(item / "artifact" / "visual.svg").getroot()
    return sorted(im.get("id") for im in root.iter(_SVG + "image"))


def test_crop_prunes_offcanvas_body_images() -> None:
    # inside the window -> keep; straddling the edge (partial overlap) -> keep;
    # wholly below the window -> drop.
    body = (
        '<image id="inside" x="550" y="80" width="100" height="50" href="#a"/>'
        '<image id="straddle" x="480" y="80" width="60" height="50" href="#b"/>'
        '<image id="outside" x="50" y="400" width="80" height="40" href="#c"/>'
    )
    with tempfile.TemporaryDirectory() as tmp:
        item = _crop_image_fixture(Path(tmp), body)
        res = crop.crop_item(item)
        assert res["images_pruned"] == 1, res
        assert _img_ids(item) == ["inside", "straddle"], _img_ids(item)


def test_crop_keeps_defs_images() -> None:
    # an off-canvas image painted indirectly via <defs> is never pruned.
    defs = '<image id="d1" x="50" y="400" width="80" height="40" href="#d"/>'
    with tempfile.TemporaryDirectory() as tmp:
        item = _crop_image_fixture(Path(tmp), "", defs)
        res = crop.crop_item(item)
        assert res["images_pruned"] == 0, res
        assert _img_ids(item) == ["d1"], _img_ids(item)


def test_crop_failsafe_unparseable_transform() -> None:
    # an off-canvas body image under a non-affine transform (rotate) is KEPT —
    # we never drop an element we cannot fully reason about.
    body = ('<g transform="rotate(45)">'
            '<image id="rot" x="50" y="400" width="80" height="40" href="#e"/></g>')
    with tempfile.TemporaryDirectory() as tmp:
        item = _crop_image_fixture(Path(tmp), body)
        res = crop.crop_item(item)
        assert res["images_pruned"] == 0, res
        assert "rot" in _img_ids(item), _img_ids(item)


def test_crop_affine_transform_honored() -> None:
    # 'moved-in' is off-canvas by raw coords but a translate brings it into the
    # window -> kept. 'moved-out' is on-canvas raw but a translate pushes it
    # wholly below the window -> dropped.
    body = (
        '<g transform="translate(600 0)">'
        '<image id="moved-in" x="0" y="80" width="100" height="50" href="#f"/></g>'
        '<g transform="translate(0 400)">'
        '<image id="moved-out" x="550" y="80" width="100" height="50" href="#g"/></g>'
    )
    with tempfile.TemporaryDirectory() as tmp:
        item = _crop_image_fixture(Path(tmp), body)
        res = crop.crop_item(item)
        assert res["images_pruned"] == 1, res
        ids = _img_ids(item)
        assert "moved-in" in ids and "moved-out" not in ids, ids


def test_crop_also_crops_evidence_svg() -> None:
    # the full-page evidence SVG is cropped to the same window so it stops
    # referencing off-canvas images; its <text> is preserved (validate relies on
    # the text enumeration), and the off-canvas image ref is dropped.
    with tempfile.TemporaryDirectory() as tmp:
        item = _crop_image_fixture(
            Path(tmp), '<image id="vis" x="550" y="80" width="100" height="50" href="#a"/>')
        ev = item / "evidence"
        ev.mkdir(parents=True)
        (ev / "source-with-text.svg").write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1000 600">'
            '<image id="ev-in" x="550" y="80" width="100" height="50" href="#a"/>'
            '<image id="ev-out" x="50" y="400" width="80" height="40" href="#c"/>'
            '<text>LABEL</text></svg>', encoding="utf-8")
        res = crop.crop_item(item)
        assert res["evidence_images_pruned"] == 1, res
        root = _ET.parse(ev / "source-with-text.svg").getroot()
        assert root.get("viewBox") == "0 0 400 240", root.get("viewBox")
        ids = sorted(im.get("id") for im in root.iter(_SVG + "image"))
        assert ids == ["ev-in"], ids
        assert len(list(root.iter(_SVG + "text"))) == 1, "evidence text must survive"


def test_gc_removes_unreferenced_assets() -> None:
    import externalize_svg_images as ext
    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "card"
        assets = item / "artifact" / "assets"
        assets.mkdir(parents=True)
        (item / "artifact" / "visual.svg").write_text(
            '<svg xmlns="http://www.w3.org/2000/svg">'
            '<image href="assets/keep.png"/></svg>', encoding="utf-8")
        (assets / "keep.png").write_bytes(b"k")
        (assets / "orphan.png").write_bytes(b"o")
        removed = ext.gc_unreferenced_assets(ext.item_svg_specs(item), assets)
        assert removed == 1, removed
        assert sorted(p.name for p in assets.iterdir()) == ["keep.png"]


def test_catalog_preview_skips_fullpage_reference_when_cropped() -> None:
    # A cropped component must NOT surface the whole-page reference.png as a
    # preview (that is what made the Draft show the full slide). The cropped
    # source-with-text.svg is the preview instead.
    import build_component_catalog as bcc
    def _build(item: Path, cropped: bool) -> None:
        (item / "artifact").mkdir(parents=True)
        (item / "evidence").mkdir(parents=True)
        src = {"region_crop": {"crop_window": [0, 0, 100, 50]}} if cropped else {}
        (item / "artifact" / "text-slots.json").write_text(
            _json.dumps({"slots": [], "source": src}), encoding="utf-8")
        (item / "evidence" / "source-with-text.svg").write_text("<svg/>", encoding="utf-8")
        (item / "evidence" / "reference.png").write_bytes(b"png")
    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "c"
        _build(item, cropped=True)
        labels = [im["label"] for im in bcc.collect_images(item)]
        assert "Reference" not in labels and "Source with text" in labels, labels
    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "c"
        _build(item, cropped=False)
        labels = [im["label"] for im in bcc.collect_images(item)]
        assert "Reference" in labels, "full-page items still surface reference.png"


def test_catalog_surfaces_text_free_variant_for_cropped_draft() -> None:
    import build_component_catalog as bcc
    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "c"
        (item / "artifact").mkdir(parents=True)
        (item / "evidence").mkdir(parents=True)
        (item / "artifact" / "text-slots.json").write_text(
            _json.dumps({"slots": [], "source": {"region_crop": {"crop_window": [0, 0, 100, 50]}}}),
            encoding="utf-8",
        )
        (item / "artifact" / "visual.svg").write_text("<svg id='text-free'/>", encoding="utf-8")
        (item / "evidence" / "source-with-text.svg").write_text("<svg id='with-text'/>", encoding="utf-8")

        labels = [im["label"] for im in bcc.collect_images(item)]

        assert labels[:2] == ["Source with text", "Text-free visual"], labels


def test_catalog_pairs_classifier_cards_with_text_free_variants() -> None:
    import build_component_catalog as bcc
    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "ai-coding-maturity-levels-strip"
        components = item / "artifact" / "components"
        evidence = item / "evidence"
        components.mkdir(parents=True)
        evidence.mkdir(parents=True)
        (item / "artifact" / "visual.svg").write_text("<svg id='full-text-free'/>", encoding="utf-8")
        (evidence / "source-with-text.svg").write_text("<svg id='full-source'/>", encoding="utf-8")
        for name in [
            "ai-coding-maturity-levels-strip-group-01.svg",
            "ai-coding-maturity-levels-strip-group-01-card-01-source.svg",
            "ai-coding-maturity-levels-strip-group-01-card-01.svg",
            "ai-coding-maturity-levels-strip-group-01-card-02-source.svg",
            "ai-coding-maturity-levels-strip-group-01-card-02.svg",
        ]:
            (components / name).write_text(f"<svg id='{name}'/>", encoding="utf-8")
        (components / "components-manifest.json").write_text(_json.dumps({
            "groups": [{
                "group_id": "ai-coding-maturity-levels-strip-group-01",
                "file": "components/ai-coding-maturity-levels-strip-group-01.svg",
                "shape_class": 1,
                "title": "Level Cards",
                "member_count": 5,
                "distinct_card_count": 5,
                "cards": [
                    {
                        "card_id": "ai-coding-maturity-levels-strip-group-01-card-01",
                        "title": "Level 1 Spicy Autocomplete",
                        "source_file": "components/ai-coding-maturity-levels-strip-group-01-card-01-source.svg",
                        "file": "components/ai-coding-maturity-levels-strip-group-01-card-01.svg",
                        "duplicate_count": 1,
                    },
                    {
                        "card_id": "ai-coding-maturity-levels-strip-group-01-card-02",
                        "title": "Level 2 AI Coding Assistants",
                        "source_file": "components/ai-coding-maturity-levels-strip-group-01-card-02-source.svg",
                        "file": "components/ai-coding-maturity-levels-strip-group-01-card-02.svg",
                        "duplicate_count": 1,
                    },
                ],
            }],
        }), encoding="utf-8")

        labels = [im["label"] for im in bcc.collect_images(item)]

        assert labels == [
            "Full component",
            "Full component (Text-free)",
            "Level 1 Spicy Autocomplete",
            "Level 1 Spicy Autocomplete (Text-free)",
            "Level 2 AI Coding Assistants",
            "Level 2 AI Coding Assistants (Text-free)",
        ], labels


def test_catalog_pairs_single_layout_row_with_text_free_variant() -> None:
    import build_component_catalog as bcc
    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "goal-card"
        components = item / "artifact" / "components"
        evidence = item / "evidence"
        components.mkdir(parents=True)
        evidence.mkdir(parents=True)
        (item / "artifact" / "visual.svg").write_text("<svg id='full-text-free'/>", encoding="utf-8")
        (evidence / "source-with-text.svg").write_text("<svg id='full-source'/>", encoding="utf-8")
        (components / "goal-card-row-01.svg").write_text("<svg id='row-free'/>", encoding="utf-8")
        (components / "goal-card-row-01-source.svg").write_text("<svg id='row-source'/>", encoding="utf-8")
        (components / "components-manifest.json").write_text(_json.dumps({
            "groups": [{
                "group_id": "goal-card-row-01",
                "file": "components/goal-card-row-01.svg",
                "shape_class": 1001,
                "layout_group": "row",
                "title": "Goal Key Result Task",
                "member_count": 5,
                "distinct_card_count": 1,
                "cards": [{
                    "card_id": "goal-card-row-01",
                    "title": "Goal Key Result Task",
                    "source_file": "components/goal-card-row-01-source.svg",
                    "file": "components/goal-card-row-01.svg",
                    "duplicate_count": 1,
                }],
            }],
        }), encoding="utf-8")

        labels = [im["label"] for im in bcc.collect_images(item)]

        assert labels == [
            "Full component",
            "Full component (Text-free)",
            "Goal Key Result Task",
            "Goal Key Result Task (Text-free)",
        ], labels


def test_catalog_rel_uses_web_safe_posix_paths() -> None:
    import build_component_catalog as bcc
    path = bcc.PROJECT_ROOT / "slide-system" / "library" / "x" / "visual.svg"
    assert bcc.rel(path) == "slide-system/library/x/visual.svg"


# --------------------------------------------------------------------------- #
# build_registry
# --------------------------------------------------------------------------- #
import build_registry as breg


def test_build_registry_projection_and_compact_keys() -> None:
    items = [{"id": "sun.x.y", "type": "card", "intent": ["a"], "tags": ["t"],
              "status": "published", "name": "drop me", "paths": {"x": 1}}]
    compact = breg.project_compact(items)
    row = compact["items"][0]
    # only the compact keys survive; heavy/identifying extras are dropped.
    assert set(row.keys()) == set(breg.COMPACT_KEYS)
    assert "name" not in row and "paths" not in row
    assert row["id"] == "sun.x.y" and row["intent"] == ["a"]


def test_build_registry_live_is_clean() -> None:
    # the real registry must have no dangling entries (every artifact exists).
    reg = breg.load_json(breg.REGISTRY)
    dangling = [i["id"] for i in reg["items"]
                if i.get("paths", {}).get("artifact")
                and not breg.resolve_repo_path(i["paths"]["artifact"]).exists()]
    assert dangling == [], f"dangling registry entries: {dangling}"


def _breg_env(tmp: Path, items: list[dict]):
    """Point build_registry's module paths at a temp registry with no library
    (so no orphans) and no history (so no zombies). Returns (registry, compact,
    retrieval) paths. Caller restores globals in a finally block."""
    reg = tmp / "visual-library.json"
    reg.write_text(json.dumps({"items": items}), encoding="utf-8")
    (tmp / "library").mkdir()
    breg.REGISTRY = reg
    breg.COMPACT = tmp / "visual-library-compact.json"
    breg.RETRIEVAL = tmp / "component-retrieval-index.jsonl"
    breg.HISTORY = tmp / "extraction-history.json"  # absent -> no zombies
    breg.LIBRARY = tmp / "library"                  # empty -> no orphans
    return reg, breg.COMPACT, breg.RETRIEVAL


def _run_breg(*flags: str) -> int:
    old = sys.argv[:]
    sys.argv = ["build_registry.py", *flags]
    try:
        return breg.main()
    finally:
        sys.argv = old


def test_build_registry_check_detects_stale_compact() -> None:
    saved = (breg.REGISTRY, breg.COMPACT, breg.RETRIEVAL, breg.HISTORY, breg.LIBRARY)
    items = [{"id": "sun.component.x", "type": "component", "status": "published",
              "intent": ["kpi"], "tags": ["t"], "content_structure": ["metric"],
              "brand": "sun-studio", "density": "any", "limitations": []}]
    with tempfile.TemporaryDirectory() as tmp:
        try:
            reg, compact, retrieval = _breg_env(Path(tmp), items)
            # Correct retrieval index so ONLY compact drift is under test.
            retrieval.write_text(breg.retrieval_jsonl(items), encoding="utf-8", newline="\n")
            # Stale compact: content that does not match the projection.
            compact.write_text('{"items": []}', encoding="utf-8")
            assert _run_breg("--check") == 1, "stale compact must fail --check"
            # A fresh, correct compact passes.
            compact.write_text(breg.compact_text(items), encoding="utf-8")
            assert _run_breg("--check") == 0, "matching compact must pass --check"
            # A missing compact also fails.
            compact.unlink()
            assert _run_breg("--check") == 1, "missing compact must fail --check"
        finally:
            breg.REGISTRY, breg.COMPACT, breg.RETRIEVAL, breg.HISTORY, breg.LIBRARY = saved


def test_build_registry_write_regenerates_stale_compact() -> None:
    saved = (breg.REGISTRY, breg.COMPACT, breg.RETRIEVAL, breg.HISTORY, breg.LIBRARY)
    items = [{"id": "sun.component.y", "type": "component", "status": "published",
              "intent": ["checklist"], "tags": ["t"], "content_structure": ["heading"],
              "brand": "sun-studio", "density": "any", "limitations": []}]
    with tempfile.TemporaryDirectory() as tmp:
        try:
            reg, compact, retrieval = _breg_env(Path(tmp), items)
            compact.write_text('{"items": []}', encoding="utf-8")  # stale
            assert _run_breg("--check") == 1
            assert _run_breg("--write") == 0
            # --write cleared the drift and produced the exact projection.
            assert compact.read_text(encoding="utf-8") == breg.compact_text(items)
            assert _run_breg("--check") == 0
        finally:
            breg.REGISTRY, breg.COMPACT, breg.RETRIEVAL, breg.HISTORY, breg.LIBRARY = saved


def test_build_registry_live_compact_projection_is_clean() -> None:
    # The committed compact must equal the deterministic projection of the live
    # full registry (guards against the stale-compact class of bug this fixes).
    reg = breg.load_json(breg.REGISTRY)
    assert breg.COMPACT.read_text(encoding="utf-8") == breg.compact_text(reg["items"]), \
        "visual-library-compact.json is stale — run build_registry.py --write"


# --------------------------------------------------------------------------- #
# classify_page_components (pure-logic paths — no browser/Chromium needed)
# --------------------------------------------------------------------------- #
import classify_page_components as cpc
import extract_editable_text_slots as eets


def test_row_title_reads_heading_columns_left_to_right() -> None:
    slots = [
        {"text": "Kết quả muốn đạt được", "role": "subheading", "x": 0.10, "y": 0.60,
         "w": 0.20, "h": 0.03, "size": 30},
        {"text": "GOAL", "role": "heading", "x": 0.10, "y": 0.45,
         "w": 0.08, "h": 0.04, "size": 53},
        {"text": "KEY", "role": "heading", "x": 0.48, "y": 0.44,
         "w": 0.05, "h": 0.03, "size": 53},
        {"text": "RESULT", "role": "heading", "x": 0.45, "y": 0.47,
         "w": 0.10, "h": 0.03, "size": 53},
        {"text": "TASK", "role": "heading", "x": 0.82, "y": 0.45,
         "w": 0.07, "h": 0.04, "size": 53},
    ]
    assert cpc._row_title(slots, "Row") == "GOAL KEY RESULT TASK"


def test_split_runs_breaks_on_large_forward_gap() -> None:
    # Three column headings concatenated in one tspan, separated by large x-gaps
    # (no space chars) must split into three runs — the bug that merged
    # STRATEGIST/DRIVER/COACH into one slot.
    text = "ABCDEFGHI"
    xs = [0, 10, 20, 200, 210, 220, 400, 410, 420]   # 3 clusters, gap 180 >> advance 10
    runs = eets.split_runs(text, xs, font_size=20)
    assert [r[0] for r in runs] == ["ABC", "DEF", "GHI"], runs


def test_split_runs_keeps_tight_text_together() -> None:
    # Ordinary evenly-spaced glyphs (a single word) are NOT split.
    text = "ABCDEF"
    xs = [0, 10, 20, 30, 40, 50]
    runs = eets.split_runs(text, xs, font_size=20)
    assert [r[0] for r in runs] == ["ABCDEF"], runs


def test_split_runs_still_breaks_on_line_wrap() -> None:
    # A backward x-jump (x resets left) is a line wrap and still splits.
    text = "ABCDE"
    xs = [0, 10, 20, 2, 12]
    runs = eets.split_runs(text, xs, font_size=20)
    assert [r[0] for r in runs] == ["ABC", "DE"], runs

_SVGNS = "http://www.w3.org/2000/svg"


def _measured(width: float, height: float, groups: list[dict]) -> dict:
    out = []
    for i, g in enumerate(groups):
        out.append({"index": i, "x": g["x"], "y": g["y"], "w": g["w"], "h": g["h"],
                    "children": g.get("children", [])})
    return {"width": width, "height": height, "groups": out}


def test_classify_drops_offcanvas_leaves() -> None:
    # A crop leaves vector junk below the viewBox; those leaves must be dropped.
    m = _measured(1000, 400, [
        {"x": 10, "y": 10, "w": 100, "h": 100},      # on-canvas
        {"x": 10, "y": 900, "w": 100, "h": 100},     # off-canvas (y >> 400)
    ])
    leaves = cpc._leaf_boxes(m, window_pad=20)
    assert len(leaves) == 1 and leaves[0]["y"] == 10, leaves


def test_classify_tall_card_absorbs_icon() -> None:
    # A portrait card filling 96% of canvas height (but narrow) must NOT be
    # treated as a bridging "bar": it has to absorb its own icon into one
    # cluster. This is the bug that left icons orphaned.
    H = 500
    m = _measured(1000, H, [
        {"x": 8, "y": 10, "w": 380, "h": 480},        # tall narrow card
        {"x": 300, "y": 400, "w": 80, "h": 80},       # icon inside the card
    ])
    leaves = cpc._leaf_boxes(m, window_pad=20)
    clusters = cpc._cluster_spatial(leaves, merge_gap=6, canvas_w=1000, canvas_h=H)
    assert len(clusters) == 1, [len(c) for c in clusters]


def test_classify_divider_does_not_bridge() -> None:
    # A thin full-width rule (extreme aspect ratio) must stay its own singleton
    # so it does not glue two separated cards into one blob.
    m = _measured(1000, 500, [
        {"x": 8, "y": 10, "w": 380, "h": 480},        # card A
        {"x": 612, "y": 10, "w": 380, "h": 480},      # card B (x-disjoint)
        {"x": 5, "y": 240, "w": 990, "h": 8},         # full-width divider
    ])
    leaves = cpc._leaf_boxes(m, window_pad=20)
    clusters = cpc._cluster_spatial(leaves, merge_gap=6, canvas_w=1000, canvas_h=500)
    # 2 cards + 1 divider singleton == 3 clusters; cards never merge.
    assert len(clusters) == 3, [cpc._bounds(c) for c in clusters]


def test_classify_dedups_same_shape_different_color() -> None:
    # Three congruent instances => one shape class (identical AND same-shape-
    # different-color both collapse). Slightly different size => its own class.
    inst = [{"w": 380, "h": 480}, {"w": 378, "h": 476}, {"w": 381, "h": 479},
            {"w": 120, "h": 120}]
    classes = cpc._shape_classes(inst, tol=0.14)
    sizes = sorted(len(c) for c in classes)
    assert sizes == [1, 3], sizes


def test_classify_groups_adjacent_same_shape_run() -> None:
    # A row of 3 same-shape instances with small gutters is ONE proximity group
    # (the whole run, rendered with each member's variant preserved).
    inst = [{"x": 0, "y": 0, "w": 100, "h": 200},
            {"x": 120, "y": 0, "w": 100, "h": 200},   # gutter 20 < 0.6*100
            {"x": 240, "y": 0, "w": 100, "h": 200}]
    groups = cpc._proximity_groups(inst, [0, 1, 2], gap_frac=0.6)
    assert len(groups) == 1 and sorted(groups[0]) == [0, 1, 2], groups


def test_classify_splits_distant_same_shape() -> None:
    # Same-shape instances sitting far apart are NOT one group — each is its own
    # standalone item.
    inst = [{"x": 0, "y": 0, "w": 100, "h": 100},
            {"x": 900, "y": 0, "w": 100, "h": 100}]   # gap 800 >> 0.6*100
    groups = cpc._proximity_groups(inst, [0, 1], gap_frac=0.6)
    assert len(groups) == 2, groups


def test_classify_keeps_different_shapes_separate() -> None:
    # Two adjacent boxes of clearly different shape land in different shape-
    # classes, so grouping (which is within-class) keeps them as 2 groups even
    # though they sit side by side.
    inst = [{"x": 0, "y": 0, "w": 100, "h": 200},     # tall
            {"x": 110, "y": 0, "w": 300, "h": 80}]    # wide — different shape
    classes = cpc._shape_classes(inst, tol=0.14)
    groups: list = []
    for class_idxs in classes:
        groups += cpc._proximity_groups(inst, class_idxs, gap_frac=0.6)
    assert len(groups) == 2, (classes, groups)


def test_child_count_mismatch_detects() -> None:
    # A group whose ElementTree child count differs from the measured child
    # count is flagged (its measured indices would copy wrong nodes).
    svg = (f'<svg xmlns="{_SVGNS}"><g><rect/><rect/></g><g><rect/></g></svg>')
    root = _ET.fromstring(svg)
    groups = list(root)  # two <g>
    measured = [{"children": [{}, {}]}, {"children": [{}, {}]}]  # 2nd says 2, ET has 1
    bad = cpc._child_count_mismatch(groups, measured)
    assert bad == [(1, 1, 2)], bad


def test_child_count_mismatch_clean() -> None:
    svg = (f'<svg xmlns="{_SVGNS}"><g><rect/><rect/></g><g><rect/></g></svg>')
    root = _ET.fromstring(svg)
    groups = list(root)
    measured = [{"children": [{}, {}]}, {"children": [{}]}]
    assert cpc._child_count_mismatch(groups, measured) == []


# Exact-equality distance for the dedup tests: 0 when equal (<= any threshold),
# large otherwise. The real call injects perceptual-signature distance instead.
_EQ_DIST = lambda a, b: 0.0 if a == b else 999.0


def test_collapse_duplicates_keeps_first_and_counts() -> None:
    # Identical items collapse into the first; counts track how many.
    kept, counts = cpc._collapse_duplicates(["a", "b", "a", "a", "c"], _EQ_DIST, 0.0)
    assert kept == [0, 1, 4] and counts == [3, 1, 1], (kept, counts)


def test_collapse_duplicates_none_never_merges() -> None:
    # A failed render (None) is always kept on its own — never silently
    # dropped or merged with another None.
    kept, counts = cpc._collapse_duplicates([None, "a", None, "a"], _EQ_DIST, 0.0)
    assert kept == [0, 1, 2] and counts == [1, 2, 1], (kept, counts)


def test_collapse_duplicates_merges_within_threshold() -> None:
    # Distance-based: items within `threshold` collapse (near-identical →
    # "tương tự"), items beyond it stay distinct (different color/icon).
    dist = lambda a, b: abs(a - b)
    # 10 and 11 are within 3 of each other; 50 is far → kept separate.
    kept, counts = cpc._collapse_duplicates([10, 11, 50, 10], dist, 3.0)
    assert kept == [0, 2] and counts == [3, 1], (kept, counts)


def test_split_on_gutter_separates_bridged_components() -> None:
    # Two big leaves with a 30px gutter, bridged by one tiny leaf sitting in the
    # gap (the card↔photo failure). The split ignores the tiny leaf when finding
    # the gutter and assigns it to the nearer side → two components.
    members = [
        {"x": 0, "y": 0, "w": 100, "h": 100},     # big left
        {"x": 130, "y": 0, "w": 100, "h": 100},   # big right (30px gutter)
        {"x": 110, "y": 45, "w": 12, "h": 12},    # tiny bridge in the gutter
    ]
    parts = cpc._split_on_gutter(members, min_gutter_px=16.0)
    assert len(parts) == 2, parts
    assert {len(p) for p in parts} == {1, 2}, parts  # tiny joins one side


def test_split_on_gutter_keeps_single_component_intact() -> None:
    # A genuine component (parts within a small gap, < threshold) is NOT split.
    members = [
        {"x": 0, "y": 0, "w": 100, "h": 100},
        {"x": 108, "y": 0, "w": 100, "h": 100},   # 8px gap < 16px threshold
    ]
    assert cpc._split_on_gutter(members, min_gutter_px=16.0) == [members]


def test_heading_picks_largest_font_with_subtitle() -> None:
    # Heading = largest font tier; a short second tier (subtitle) is appended,
    # read top-to-bottom.
    slots = [{"text": "Level 1", "x": 0.1, "y": 0.10, "w": 0.1, "h": 0.02, "size": 53.0},
             {"text": "Spicy", "x": 0.1, "y": 0.14, "w": 0.1, "h": 0.02, "size": 42.0},
             {"text": "Autocomplete", "x": 0.1, "y": 0.16, "w": 0.1, "h": 0.02, "size": 42.0},
             {"text": "a long body copy line here", "x": 0.1, "y": 0.30, "w": 0.1, "h": 0.02, "size": 18.0}]
    assert cpc._heading(slots) == "Level 1 Spicy Autocomplete", cpc._heading(slots)


def test_heading_drops_paragraph_tier() -> None:
    # When the second font tier is a multi-slot paragraph (>3 slots), it is NOT
    # appended — only the heading survives.
    slots = [{"text": "TRANSLATOR", "x": 0.1, "y": 0.10, "w": 0.1, "h": 0.02, "size": 38.0}]
    slots += [{"text": f"w{i}", "x": 0.1, "y": 0.2 + i * 0.01, "w": 0.05, "h": 0.01, "size": 20.0}
              for i in range(5)]
    assert cpc._heading(slots) == "TRANSLATOR", cpc._heading(slots)


def test_group_title_common_prefix_and_join() -> None:
    assert cpc._group_title(["Level 1 X", "Level 2 Y", "Level 3 Z"]) == "Level cards"
    # No shared prefix → join the heading-like (short, capitalized) titles only.
    assert cpc._group_title(["TRANSLATOR", "a long body copy fallback here", "DRIVER"]) \
        == "TRANSLATOR / DRIVER"


def test_layout_cells_split_single_row_cards() -> None:
    instances = [
        {"x": 10, "y": 20, "w": 120, "h": 240, "members": [{"group": 0, "child": None}]},
        {"x": 150, "y": 20, "w": 120, "h": 240, "members": [{"group": 1, "child": None}]},
        {"x": 290, "y": 20, "w": 120, "h": 240, "members": [{"group": 2, "child": None}]},
    ]
    small = [
        {"x": 330, "y": 40, "w": 20, "h": 20, "members": [{"group": 3, "child": None}]},
    ]

    assert cpc._cluster_layout_rows(instances, small) == []
    cells = cpc._cluster_layout_cells(instances, small)

    assert len(cells) == 3
    assert [cell["col_index"] for cell in cells] == [1, 2, 3]
    assert len(cells[2]["elements"]) == 2, cells


def test_tags_from_dedups_and_skips_stopwords() -> None:
    tags = cpc._tags_from(["Level 1 Spicy", "Level 2 Coding", "the and Spicy"])
    assert tags == ["Level", "1", "Spicy", "2", "Coding"], tags


def test_slots_in_uses_center_point() -> None:
    slots = [{"text": "in", "x": 0.10, "y": 0.10, "w": 0.05, "h": 0.05, "size": 10},
             {"text": "out", "x": 0.80, "y": 0.80, "w": 0.05, "h": 0.05, "size": 10}]
    got = cpc._slots_in(slots, 0, 0, 500, 500, 1000, 1000)  # px box covers left-top quadrant
    assert [s["text"] for s in got] == ["in"], got


def test_classify_records_dropped_small_with_bounds() -> None:
    # A cluster below the area floor is recorded with its bounds (not just
    # counted), so a genuine small component stays inspectable.
    canvas_w, canvas_h = 1000.0, 1000.0
    area = canvas_w * canvas_h
    clusters = [
        [{"x": 10, "y": 10, "w": 20, "h": 20, "group": 0, "child": None}],   # tiny < 1.5%
        [{"x": 100, "y": 100, "w": 300, "h": 300, "group": 1, "child": None}],
    ]
    dropped = []
    for members in clusters:
        x0, y0, w, h = cpc._bounds(members)
        a = w * h
        if a < 0.015 * area:
            dropped.append({"x": round(x0, 1), "y": round(y0, 1),
                            "w": round(w, 1), "h": round(h, 1),
                            "area_frac": round(a / area, 4)})
    assert len(dropped) == 1
    assert set(dropped[0]) == {"x", "y", "w", "h", "area_frac"}, dropped[0]
    assert dropped[0]["w"] == 20 and dropped[0]["area_frac"] == 0.0004, dropped[0]


def test_classify_excludes_fullbleed_background() -> None:
    # On a full page the background cluster (≈ canvas size) must be routed to
    # background_candidates, NOT emitted as a component class.
    canvas_w, canvas_h = 1000.0, 800.0
    clusters = [
        [{"x": 2, "y": 2, "w": 996, "h": 796, "group": 0, "child": None}],   # bg
        [{"x": 100, "y": 100, "w": 200, "h": 200, "group": 1, "child": None}],
        [{"x": 400, "y": 100, "w": 200, "h": 200, "group": 2, "child": None}],
    ]
    area = canvas_w * canvas_h
    bg, inst, dropped = [], [], 0
    for members in clusters:
        x0, y0, w, h = cpc._bounds(members)
        a = w * h
        if a < 0.015 * area:
            dropped += 1
        elif a >= 0.7 * area:
            bg.append(members)
        else:
            inst.append(members)
    assert len(bg) == 1 and len(inst) == 2, (len(bg), len(inst))


def test_classify_ancestor_transform_and_fragment() -> None:
    # The crop wrapper transform must be captured and re-applied in fragments,
    # else geometry lands off-canvas and the fragment renders blank.
    svg = (f'<svg xmlns="{_SVGNS}" xmlns:inkscape="http://www.inkscape.org/namespaces/inkscape" '
           'viewBox="0 0 100 100">'
           '<g transform="translate(-440 -440)"><g>'
           '<g inkscape:groupmode="layer">'
           '<g><rect x="450" y="450" width="20" height="20"/></g>'
           '</g></g></g></svg>')
    root = _ET.fromstring(svg)
    groups = cpc.document_groups(root)
    parent_map = {c: p for p in root.iter() for c in p}
    chain = cpc._ancestor_transform(root, parent_map, groups[0])
    assert chain == "translate(-440 -440)", repr(chain)
    members = [{"x": 10, "y": 10, "w": 20, "h": 20, "group": 0, "child": None}]
    frag = cpc._build_fragment(members, groups, [], margin=3,
                               source_dir=Path("."), ancestor_transform=chain)
    blob = _ET.tostring(frag, encoding="unicode")
    assert "translate(-440 -440)" in blob and "rect" in blob, blob


# --------------------------------------------------------------------------- #
# split_icon_sheet — icon-sheet decomposition helpers
# --------------------------------------------------------------------------- #
def test_split_cluster_1d_groups_within_tol() -> None:
    import split_icon_sheet as sis
    # two tight groups (~5 and ~105) separated by a wide gap -> two lines
    lines = sis._cluster_1d([4, 5, 6, 104, 105, 106], tol=20)
    assert len(lines) == 2, lines
    assert abs(lines[0] - 5) < 1 and abs(lines[1] - 105) < 1, lines


def test_split_merge_within_fuses_overlap_keeps_distant() -> None:
    import split_icon_sheet as sis
    # A and B overlap (gap 0); C is far away on x -> {A,B} merged, C separate.
    clusters = [
        {"x": 0, "y": 0, "w": 20, "h": 20, "members": [{"x": 0, "y": 0, "w": 20, "h": 20}]},
        {"x": 18, "y": 2, "w": 20, "h": 20, "members": [{"x": 18, "y": 2, "w": 20, "h": 20}]},
        {"x": 200, "y": 0, "w": 20, "h": 20, "members": [{"x": 200, "y": 0, "w": 20, "h": 20}]},
    ]
    out = sis._merge_within(clusters, gap=10)
    assert len(out) == 2, [(o["x"], o["w"]) for o in out]
    big = max(out, key=lambda o: o["w"])
    assert len(big["members"]) == 2, big


def test_split_per_row_gap_separates_neighbours_fuses_fragments() -> None:
    """The core grid rule: cells in one row split into icons by the gap valley —
    fragments (gap < col_tol) fuse, neighbours (gap >> col_tol) stay separate."""
    import split_icon_sheet as sis
    # row of cells (same y): two fragments of icon A (x=0..20, 25..45, gap 5),
    # then icon B far right (x=120..145, gap 75). col_tol=35 sits in the valley.
    row = [
        {"x": 0, "y": 0, "w": 20, "h": 30},
        {"x": 25, "y": 0, "w": 20, "h": 30},
        {"x": 120, "y": 0, "w": 25, "h": 30},
    ]
    by_cell: dict = {}
    row.sort(key=lambda c: c["x"])
    col, right = 0, None
    for c in row:
        if right is not None and (c["x"] - right) >= 35:
            col += 1
        by_cell.setdefault(col, []).append(c)
        right = max(right if right is not None else c["x"] + c["w"], c["x"] + c["w"])
    assert len(by_cell) == 2, by_cell
    assert len(by_cell[0]) == 2 and len(by_cell[1]) == 1, by_cell


def test_build_catalog_collect_icon_set_parses_and_absent() -> None:
    import build_component_catalog as bcc
    import json as _json
    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp)
        # no manifest -> None
        assert bcc.collect_icon_set(item) is None
        icons_dir = item / "artifact" / "icons"
        icons_dir.mkdir(parents=True)
        (icons_dir / "icon-000.svg").write_text("<svg/>")
        (icons_dir / "icon-001.svg").write_text("<svg/>")
        (icons_dir / "icons-manifest.json").write_text(_json.dumps({"icons": [
            {"index": 0, "file": "icon-000.svg", "slug": "bod", "name": "BOD",
             "region": "frequently-used", "row": -1, "col": -1},
            {"index": 1, "file": "icon-001.svg", "slug": "wifi", "name": "wifi",
             "region": "grid", "row": 1, "col": 2},
            {"index": 2, "file": "missing.svg", "slug": "x", "name": "x",
             "region": "grid", "row": 1, "col": 3},
        ]}))
        iset = bcc.collect_icon_set(item)
        assert iset and iset["count"] == 2, iset  # missing file dropped
        slugs = [i["slug"] for i in iset["icons"]]
        assert slugs == ["bod", "wifi"], slugs
        assert iset["icons"][0]["path"].endswith("icon-000.svg")


def test_build_catalog_skips_blank_text_free_visual_marked_by_quality_gate() -> None:
    import build_component_catalog as bcc

    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "demo"
        comps = item / "artifact" / "components"
        evidence = item / "evidence"
        comps.mkdir(parents=True)
        evidence.mkdir(parents=True)
        (item / "mapping.json").write_text(json.dumps({
            "item_id": "demo",
            "status": "staging",
            "quality_gate": {"blank_item_visual": True},
        }), encoding="utf-8")
        (evidence / "source-with-text.svg").write_text(
            '<svg xmlns="http://www.w3.org/2000/svg">'
            '<rect width="100" height="80" fill="#3333FF"/></svg>',
            encoding="utf-8",
        )
        (item / "artifact" / "visual.svg").write_text(
            '<svg xmlns="http://www.w3.org/2000/svg">'
            '<rect width="100" height="80" fill="#3333FF"/></svg>',
            encoding="utf-8",
        )
        (comps / "components-manifest.json").write_text('{"groups":[]}', encoding="utf-8")

        labels = [image["label"] for image in bcc.collect_images(item)]

        assert labels == ["Full component"], labels


def test_build_catalog_skips_standalone_blank_visual_drafts() -> None:
    import build_component_catalog as bcc

    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp) / "component-extractions"
        item = root / "demo" / "items" / "tiny-label"
        (item / "artifact").mkdir(parents=True)
        (item / "evidence").mkdir(parents=True)
        (item / "mapping.json").write_text(json.dumps({
            "item_id": "tiny-label",
            "candidate_stable_id": "sun.component.tiny-label",
            "name": "Tiny Label",
            "status": "staging",
            "type": "component",
            "category": "component",
            "source": {"path": "source.pdf", "slide_or_page": 1},
            "quality_gate": {"blank_item_visual": True},
        }), encoding="utf-8")
        (item / "artifact" / "visual.svg").write_text("<svg/>", encoding="utf-8")
        (item / "evidence" / "source-with-text.svg").write_text(
            '<svg xmlns="http://www.w3.org/2000/svg"><text>04</text></svg>',
            encoding="utf-8",
        )
        registry = Path(tmp) / "registry.json"
        registry.write_text('{"items":[]}', encoding="utf-8")
        output = Path(tmp) / "catalog-data.json"
        old_argv = sys.argv[:]
        sys.argv = [
            "build_component_catalog.py",
            "--registry", str(registry),
            "--extractions", str(root),
            "--output", str(output),
        ]
        try:
            assert bcc.main() == 0
        finally:
            sys.argv = old_argv

        catalog = json.loads(output.read_text(encoding="utf-8"))

        assert [item["id"] for item in catalog["items"]] == []


def test_quality_gate_prunes_blank_refs_and_empty_manifests() -> None:
    import quality_gate as qg

    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "demo"
        comps = item / "artifact" / "components"
        comps.mkdir(parents=True)
        (item / "mapping.json").write_text(json.dumps({
            "item_id": "demo",
            "status": "staging",
        }), encoding="utf-8")
        (comps / "blank.svg").write_text(
            '<svg xmlns="http://www.w3.org/2000/svg"></svg>',
            encoding="utf-8",
        )
        (comps / "card.svg").write_text(
            '<svg xmlns="http://www.w3.org/2000/svg">'
            '<rect width="100" height="80" fill="#3333FF"/></svg>',
            encoding="utf-8",
        )
        manifest = comps / "components-manifest.json"
        manifest.write_text(json.dumps({
            "groups": [
                {"group_id": "blank", "file": "components/blank.svg", "cards": []},
                {"group_id": "card", "file": "components/card.svg", "cards": [
                    {"card_id": "card", "file": "components/card.svg"},
                    {"card_id": "blank-card", "file": "components/blank.svg"},
                ]},
            ],
        }), encoding="utf-8")

        summary = qg.sanitize_item(item)
        cleaned = json.loads(manifest.read_text(encoding="utf-8"))
        mapping = json.loads((item / "mapping.json").read_text(encoding="utf-8"))

        assert summary["blank_refs_pruned"] == 2, summary
        assert [g["group_id"] for g in cleaned["groups"]] == ["card"]
        assert cleaned["groups"][0]["cards"] == [{"card_id": "card", "file": "components/card.svg"}]
        assert mapping["quality_gate"]["status"] == "reviewable"

    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "empty"
        comps = item / "artifact" / "components"
        comps.mkdir(parents=True)
        (item / "mapping.json").write_text(json.dumps({
            "item_id": "empty",
            "status": "staging",
        }), encoding="utf-8")
        manifest = comps / "components-manifest.json"
        manifest.write_text('{"groups":[]}', encoding="utf-8")

        summary = qg.sanitize_item(item)
        mapping = json.loads((item / "mapping.json").read_text(encoding="utf-8"))

        assert summary["empty_manifests_removed"] == 1, summary
        assert not manifest.exists()
        assert mapping["quality_gate"]["status"] == "needs_review"


def test_quality_gate_prunes_render_blank_refs_and_marks_base_visual() -> None:
    import quality_gate as qg

    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "rendered"
        artifact = item / "artifact"
        comps = artifact / "components"
        comps.mkdir(parents=True)
        (item / "mapping.json").write_text(json.dumps({
            "item_id": "rendered",
            "status": "staging",
        }), encoding="utf-8")
        visual = artifact / "visual.svg"
        blank = comps / "render-blank.svg"
        source = comps / "source.svg"
        for path in (visual, blank, source):
            path.write_text(
                '<svg xmlns="http://www.w3.org/2000/svg">'
                '<rect width="100" height="80" fill="#3333FF"/></svg>',
                encoding="utf-8",
            )
        manifest = comps / "components-manifest.json"
        manifest.write_text(json.dumps({
            "groups": [{
                "group_id": "rendered",
                "file": "components/render-blank.svg",
                "cards": [
                    {"card_id": "blank", "file": "components/render-blank.svg"},
                    {"card_id": "source", "source_file": "components/source.svg"},
                ],
            }],
        }), encoding="utf-8")
        render_results = {
            str(visual.resolve()).lower(): {"nonwhite_ratio": 0.0},
            str(blank.resolve()).lower(): {"nonwhite_ratio": 0.0},
            str(source.resolve()).lower(): {"nonwhite_ratio": 0.05},
        }

        summary = qg.sanitize_items([item], render_results=render_results)[0]
        cleaned = json.loads(manifest.read_text(encoding="utf-8"))
        mapping = json.loads((item / "mapping.json").read_text(encoding="utf-8"))

        assert summary["blank_item_visual"], summary
        assert summary["render_blank_refs_pruned"] == 2, summary
        assert summary["status"] == "needs_review"
        assert "file" not in cleaned["groups"][0]
        assert cleaned["groups"][0]["cards"] == [
            {"card_id": "source", "source_file": "components/source.svg"},
        ]
        assert mapping["quality_gate"]["blank_item_visual"] is True
        assert mapping["quality_gate"]["item_visual_nonwhite_ratio"] == 0.0


def test_quality_gate_ignores_white_defs_and_masks() -> None:
    import quality_gate as qg

    with tempfile.TemporaryDirectory() as tmp:
        svg = Path(tmp) / "masked.svg"
        svg.write_text(
            '<svg xmlns="http://www.w3.org/2000/svg">'
            '<defs><rect width="100" height="100" fill="#3333FF"/></defs>'
            '<mask id="m"><rect width="100" height="100" fill="white"/></mask>'
            '<rect width="100" height="100" fill="white"/></svg>',
            encoding="utf-8",
        )

        assert not qg.svg_has_visible_content(svg)


def test_retrieval_index_builds_published_search_records() -> None:
    import build_component_retrieval_index as bri

    registry = {
        "items": [
            {
                "id": "sun.component.metric-strip",
                "name": "Metric Strip",
                "status": "published",
                "type": "component",
                "component_type": "strip",
                "layout_role": "metric comparison strip",
                "intent": ["revenue growth"],
                "tags": ["metric", "growth"],
                "keywords": ["revenue", "team-size"],
                "content_structure": ["label", "percentage"],
                "use_cases": ["Show KPI change"],
                "anti_use_cases": ["Do not use for pie charts"],
                "visual_summary": "Two horizontal metric bars.",
                "retrieval_notes": "Use when user asks for KPI cards.",
                "source": {"kind": "extraction", "path": "/Users/home/private/source.pdf"},
                "paths": {"artifact": "slide-system/library/components/metric-strip"},
            },
            {
                "id": "sun.component.draft-only",
                "name": "Draft Only",
                "status": "staging",
                "intent": ["draft"],
            },
        ],
    }

    records = bri.build_records(registry)

    assert [record["id"] for record in records] == ["sun.component.metric-strip"]
    assert records[0]["retrieval_mode"] == "lexical-ready"
    assert "revenue" in records[0]["retrieval_terms"]
    assert "users" not in records[0]["retrieval_terms"]
    assert "pie" in records[0]["search_text"]
    assert records[0]["paths"]["artifact"] == "slide-system/library/components/metric-strip"


def test_publish_preserves_retrieval_metadata_and_index() -> None:
    import importlib

    publish = importlib.import_module("publish_extraction")
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        extraction_dir = root / "extract"
        item_dir = extraction_dir / "items" / "metric-strip"
        (item_dir / "artifact").mkdir(parents=True)
        (item_dir / "preview").mkdir()
        (item_dir / "evidence").mkdir()
        (item_dir / "artifact" / "visual.svg").write_text(
            '<svg xmlns="http://www.w3.org/2000/svg"><rect width="10" height="10"/></svg>',
            encoding="utf-8",
        )
        (item_dir / "preview" / "thumbnail.png").write_bytes(b"not-a-real-png")
        (item_dir / "evidence" / "source-with-text.svg").write_text("<svg/>", encoding="utf-8")
        (item_dir / "mapping.json").write_text(json.dumps({
            "extraction_id": "publish-meta-demo",
            "item_id": "metric-strip",
            "candidate_stable_id": "sun.component.metric-strip",
            "name": "Metric Strip",
            "status": "staging",
            "type": "component",
            "category": "metrics",
            "brand": "sun-studio",
            "semantic_intent": ["revenue growth"],
            "tags": ["metric", "growth"],
            "content_structure": ["label", "percentage"],
            "content_fields": {},
            "density": "any",
            "component_type": "strip",
            "layout_role": "metric comparison strip",
            "visual_summary": "Two horizontal KPI bars.",
            "keywords": ["revenue", "team-size"],
            "use_cases": ["Show KPI change"],
            "anti_use_cases": ["Do not use for pie charts"],
            "quality_notes": "Reviewed in Draft.",
            "retrieval_notes": "Use when user asks for KPI cards.",
            "artifact_status": "ready",
            "approval": {"status": "approved"},
            "source": {
                "path": str(root / "source.pdf"),
                "slide_or_page": 1,
                "region": {"x": 0, "y": 0, "width": 1, "height": 1, "unit": "normalized"},
                "sha256": "source-hash",
            },
            "fingerprints": {
                "region_identity_sha256": "region-hash",
                "semantic_signature_sha256": "semantic-hash",
            },
        }), encoding="utf-8")

        registry = root / "visual-library.json"
        registry.write_text('{"items":[]}', encoding="utf-8")
        history = root / "history.json"
        history.write_text('{"attempts":[]}', encoding="utf-8")
        library = root / "library"
        old_argv = sys.argv[:]
        sys.argv = [
            "publish_extraction.py",
            "--extraction-dir", str(extraction_dir),
            "--item-id", "metric-strip",
            "--registry", str(registry),
            "--history", str(history),
            "--library-root", str(library),
        ]
        try:
            assert publish.main() == 0
        finally:
            sys.argv = old_argv

        item = read_text_slots.load_json(registry)["items"][0]
        assert item["component_type"] == "strip"
        assert item["keywords"] == ["revenue", "team-size"]
        assert item["use_cases"] == ["Show KPI change"]
        assert item["retrieval_notes"] == "Use when user asks for KPI cards."

        index = (root / "component-retrieval-index.jsonl").read_text(encoding="utf-8")
        record = json.loads(index.strip())
        assert record["id"] == "sun.component.metric-strip"
        assert "revenue" in record["retrieval_terms"]
        assert "kpi" in record["retrieval_terms"]


def test_publish_rejects_failed_auto_stage_artifacts() -> None:
    import importlib

    publish = importlib.import_module("publish_extraction")
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        extraction_dir = root / "extract"
        item_dir = extraction_dir / "items" / "broken-strip"
        (item_dir / "artifact").mkdir(parents=True)
        (item_dir / "preview").mkdir()
        (item_dir / "evidence").mkdir()
        (item_dir / "artifact" / "visual.svg").write_text("<svg/>", encoding="utf-8")
        (item_dir / "preview" / "thumbnail.png").write_bytes(b"not-a-real-png")
        (item_dir / "evidence" / "source-with-text.svg").write_text("<svg/>", encoding="utf-8")
        (item_dir / "mapping.json").write_text(json.dumps({
            "extraction_id": "publish-failed-demo",
            "item_id": "broken-strip",
            "candidate_stable_id": "sun.component.broken-strip",
            "name": "Broken Strip",
            "status": "staging",
            "type": "component",
            "category": "metrics",
            "brand": "sun-studio",
            "semantic_intent": ["broken metric"],
            "tags": [],
            "content_structure": [],
            "content_fields": {},
            "artifact_status": "failed",
            "artifact_log": "validate_text_slots.py: failed",
            "approval": {"status": "approved"},
            "source": {
                "path": str(root / "source.pdf"),
                "slide_or_page": 1,
                "region": {"x": 0, "y": 0, "width": 1, "height": 1, "unit": "normalized"},
                "sha256": "source-hash",
            },
            "fingerprints": {
                "region_identity_sha256": "region-hash",
                "semantic_signature_sha256": "semantic-hash",
            },
        }), encoding="utf-8")
        registry = root / "visual-library.json"
        registry.write_text('{"items":[]}', encoding="utf-8")
        history = root / "history.json"
        history.write_text('{"attempts":[]}', encoding="utf-8")
        old_argv = sys.argv[:]
        sys.argv = [
            "publish_extraction.py",
            "--extraction-dir", str(extraction_dir),
            "--item-id", "broken-strip",
            "--registry", str(registry),
            "--history", str(history),
            "--library-root", str(root / "library"),
        ]
        try:
            try:
                publish.main()
            except SystemExit as exc:
                assert "Artifact build status is failed" in str(exc)
            else:
                raise AssertionError("failed auto-stage artifacts must not publish")
        finally:
            sys.argv = old_argv
        assert read_text_slots.load_json(registry)["items"] == []


# --------------------------------------------------------------------------- #
# validate_component_metadata — retrieval metadata quality gate
# --------------------------------------------------------------------------- #
def _meta_component(**over) -> dict:
    """A fully retrieval-ready component item (passes the metadata gate)."""
    base = {
        "id": "sun.component.demo-strip",
        "type": "component",
        "category": "component",
        "name": "Demo Metric Strip",
        "intent": ["statistics", "metrics"],
        "tags": ["strip", "kpi", "set-of-3"],
        "content_structure": ["label", "metric"],
        "component_type": "strip",
        "layout_role": "horizontal metric strip",
        "visual_summary": "Three KPI figures side by side with labels.",
        "keywords": ["revenue", "growth", "kpi"],
        "use_cases": ["Show three headline KPIs on one row"],
        "anti_use_cases": ["Do not use for narrative body text"],
        "retrieval_notes": "Select when the slide needs a compact KPI row.",
        "quality_notes": "Manually reviewed against source render.",
        "text_contract": {"slot_count": 6},
    }
    base.update(over)
    return base


def test_component_metadata_valid_passes() -> None:
    import validate_component_metadata as vcm
    assert vcm.validate_item(_meta_component()) == []
    assert vcm.validate_registry({"items": [_meta_component()]}) == {}


def test_component_metadata_missing_fields_fail() -> None:
    import validate_component_metadata as vcm
    errs = vcm.validate_item(_meta_component(keywords=[], use_cases=[],
                                             component_type=None, visual_summary="  "))
    joined = " ".join(errs)
    assert "'keywords' is empty" in joined
    assert "'use_cases' is empty" in joined
    assert "'component_type' is blank" in joined
    assert "'visual_summary' is blank" in joined


def test_component_metadata_boilerplate_fails() -> None:
    import validate_component_metadata as vcm
    # auto-stage tag + Docling placeholder use/anti text must be rejected.
    errs = vcm.validate_item(_meta_component(
        tags=["strip", "auto-staged"],
        use_cases=["Review and publish this strip as a reusable component."],
        anti_use_cases=["Do not use before the Draft preview and metadata are reviewed."],
        retrieval_notes="Generated from region text, Docling label, source name, and page.",
    ))
    assert any("auto-stage/placeholder text" in e for e in errs), errs
    # An honest note that merely mentions Docling must NOT trip the gate.
    ok = vcm.validate_item(_meta_component(
        retrieval_notes="Region isolated manually; not a Docling auto-detected candidate."))
    assert ok == [], ok


def test_component_metadata_ocr_intent_fails() -> None:
    import validate_component_metadata as vcm
    errs = vcm.validate_item(_meta_component(
        intent=["2. THƯỜNG DÙNG ĐỂ BIỂU THỊ CÁC DẠNG CONTENT XOAY QUANH team"]))
    assert any("raw slide text/OCR" in e for e in errs), errs


def test_component_metadata_ignores_non_component_types() -> None:
    import validate_component_metadata as vcm
    # A template with deliberately thin metadata must NOT be gated.
    thin_template = {"id": "sun.deck.01-cover", "type": "template",
                     "intent": ["cover"], "tags": []}
    assert vcm.validate_item(thin_template) == []
    assert vcm.validate_registry({"items": [thin_template]}) == {}


def test_component_metadata_mapping_projection() -> None:
    import validate_component_metadata as vcm
    mapping = {
        "candidate_stable_id": "sun.component.demo-strip", "type": "component",
        "category": "component", "name": "Demo Metric Strip",
        "semantic_intent": ["statistics"], "tags": ["strip"],
        "content_structure": ["metric"], "component_type": "strip",
        "layout_role": "strip", "visual_summary": "KPI strip.",
        "keywords": ["kpi"], "use_cases": ["Show KPIs"],
        "anti_use_cases": ["No body text"], "retrieval_notes": "Pick for KPIs.",
        "quality_notes": "Reviewed.",
    }
    item = vcm.metadata_from_mapping(mapping)
    assert item["intent"] == ["statistics"], "semantic_intent must map to intent"
    assert vcm.validate_item(item) == []


def test_component_metadata_strict_requires_set_shape() -> None:
    import validate_component_metadata as vcm
    set_item = _meta_component(id="sun.component.role-card-set", category="component-set",
                               name="Role Card Set", tags=["roles", "personas"],
                               content_structure=["heading"], use_cases=["Show roles"])
    assert vcm.validate_item(set_item, strict=False) == []
    strict_errs = vcm.validate_item(set_item, strict=True)
    assert any("set-of-N" in e for e in strict_errs), strict_errs
    # Exposing the multiplicity clears the strict check.
    fixed = _meta_component(id="sun.component.role-card-set", category="component-set",
                            name="Role Card Set", tags=["cards", "set-of-4"])
    assert vcm.validate_item(fixed, strict=True) == []


def test_component_metadata_real_registry_good_components_pass() -> None:
    # The three hand-authored components in the live registry must pass; this
    # guards against the gate regressing into false positives on real data.
    import validate_component_metadata as vcm
    registry = read_text_slots.load_json(REGISTRY)
    by_id = {i["id"]: i for i in registry["items"]}
    for good in ("sun.component.lorem-ipsum-circle-badge-set",
                 "sun.component.foundation-top1-microsoft-overlap-circle-set",
                 "sun.component.goal-keyresult-task-hexagon-diagram"):
        assert vcm.validate_item(by_id[good]) == [], f"{good} should pass"


def test_component_metadata_live_registry_all_components_pass() -> None:
    import validate_component_metadata as vcm
    registry = read_text_slots.load_json(REGISTRY)
    failures = vcm.validate_registry(registry, strict=True)
    assert failures == {}, json.dumps(failures, indent=2, ensure_ascii=False)


def test_publish_blocks_weak_component_metadata_before_mutation() -> None:
    import importlib
    publish = importlib.import_module("publish_extraction")
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        extraction_dir = root / "extract"
        item_dir = extraction_dir / "items" / "weak-strip"
        (item_dir / "artifact").mkdir(parents=True)
        (item_dir / "preview").mkdir()
        (item_dir / "evidence").mkdir()
        (item_dir / "artifact" / "visual.svg").write_text("<svg/>", encoding="utf-8")
        (item_dir / "preview" / "thumbnail.png").write_bytes(b"not-a-real-png")
        (item_dir / "evidence" / "source-with-text.svg").write_text("<svg/>", encoding="utf-8")
        (item_dir / "mapping.json").write_text(json.dumps({
            "extraction_id": "publish-weak-demo",
            "item_id": "weak-strip",
            "candidate_stable_id": "sun.component.weak-strip",
            "name": "Weak Strip",
            "status": "staging",
            "type": "component",
            "category": "component",
            "brand": "sun-studio",
            # Auto-stage boilerplate + empty retrieval fields — must be blocked.
            "semantic_intent": ["weak strip", "picture candidate detected by Docling"],
            "tags": ["strip", "auto-staged"],
            "content_structure": [],
            "content_fields": {},
            "artifact_status": "ready",
            "approval": {"status": "approved"},
            "source": {
                "path": str(root / "source.pdf"),
                "slide_or_page": 1,
                "region": {"x": 0, "y": 0, "width": 1, "height": 1, "unit": "normalized"},
                "sha256": "source-hash",
            },
            "fingerprints": {
                "region_identity_sha256": "region-hash",
                "semantic_signature_sha256": "semantic-hash",
            },
        }), encoding="utf-8")
        registry = root / "visual-library.json"
        registry.write_text('{"items":[]}', encoding="utf-8")
        history = root / "history.json"
        history.write_text('{"attempts":[]}', encoding="utf-8")
        library = root / "library"
        old_argv = sys.argv[:]
        sys.argv = [
            "publish_extraction.py",
            "--extraction-dir", str(extraction_dir),
            "--item-id", "weak-strip",
            "--registry", str(registry),
            "--history", str(history),
            "--library-root", str(library),
        ]
        try:
            try:
                publish.main()
            except SystemExit as exc:
                assert "metadata gate failed" in str(exc).lower(), str(exc)
            else:
                raise AssertionError("weak component metadata must block publish")
        finally:
            sys.argv = old_argv
        # No registry, index, or library mutation may have occurred.
        assert read_text_slots.load_json(registry)["items"] == []
        assert not (root / "component-retrieval-index.jsonl").exists()
        assert not library.exists()


# --------------------------------------------------------------------------- #
# materialize_groups (classify_page_components + _common hash helpers)
# --------------------------------------------------------------------------- #
import json as _json_stdlib
from _common import normalized_bounds as _normalized_bounds
from _common import region_identity_hash as _region_identity_hash
from _common import semantic_signature_hash as _semantic_signature_hash


def test_group_bounds_to_normalized_region() -> None:
    canvas = {"w": 2938.83, "h": 2623.16}
    gb = {"x": 563.1, "y": 371.4, "w": 1867.5, "h": 586.3}
    region = _normalized_bounds({
        "x": gb["x"] / canvas["w"],
        "y": gb["y"] / canvas["h"],
        "width": gb["w"] / canvas["w"],
        "height": gb["h"] / canvas["h"],
        "unit": "normalized",
    })
    assert region["unit"] == "normalized"
    assert 0.0 < region["x"] < 1.0, region["x"]
    assert 0.0 < region["y"] < 1.0, region["y"]
    assert 0.0 < region["width"] < 1.0, region["width"]
    assert 0.0 < region["height"] < 1.0, region["height"]
    assert abs(region["x"] - 563.1 / 2938.83) < 1e-5
    assert abs(region["width"] - 1867.5 / 2938.83) < 1e-5


def test_materialized_mapping_fields() -> None:
    region = _normalized_bounds({
        "x": 0.19, "y": 0.14, "width": 0.64, "height": 0.22, "unit": "normalized",
    })
    rh = _region_identity_hash("sha_abc", "2", region, ["obj-1"])
    sh = _semantic_signature_hash(["Cover", "intro"])
    assert isinstance(rh, str) and len(rh) == 64, rh
    assert isinstance(sh, str) and len(sh) == 64, sh
    candidate = "sun.component.feature-step-shape-diagrams.g01"
    assert candidate.startswith("sun.component.")
    assert candidate.endswith(".g01")
    rh2 = _region_identity_hash("sha_abc", 2, region, ["obj-1"])
    assert rh == rh2, "int vs str slide_or_page must produce same hash"


def test_carved_slots_within_unit_and_subset() -> None:
    base_slots = [
        {"id": f"s{i}", "bounds": {"x": i * 0.1, "y": 0.1, "width": 0.08, "height": 0.05}}
        for i in range(10)
    ]
    region = {"x": 0.2, "y": 0.05, "width": 0.5, "height": 0.3}
    carved = [
        s for s in base_slots
        if (region["x"] <= s["bounds"]["x"] + s["bounds"]["width"] / 2 <= region["x"] + region["width"]
            and region["y"] <= s["bounds"]["y"] + s["bounds"]["height"] / 2 <= region["y"] + region["height"])
    ]
    assert len(carved) < len(base_slots), "carve should drop some slots"
    assert len(carved) > 0, "carve should keep some slots"
    for s in carved:
        cx = s["bounds"]["x"] + s["bounds"]["width"] / 2
        cy = s["bounds"]["y"] + s["bounds"]["height"] / 2
        assert region["x"] <= cx <= region["x"] + region["width"]
        assert region["y"] <= cy <= region["y"] + region["height"]


# --------------------------------------------------------------------------- #
# score_visual_items CLI — request/batch shape guards
# --------------------------------------------------------------------------- #
_SINGLE_REQUEST = {
    "request_id": "slide-01", "intent": ["timeline"], "tags": [],
    "content_structure": ["a"], "density": "medium", "brand": "sun",
    "required_exports": [],
}
_BATCH_REQUEST = {"job_id": "j", "brief": "b", "slides": [_SINGLE_REQUEST]}


def _run_scorer(tmp: Path, flag: str, payload) -> tuple[int, str, Path]:
    request_path = tmp / "in.json"
    output_path = tmp / "out.json"
    request_path.write_text(json.dumps(payload), encoding="utf-8")
    proc = subprocess.run(
        [sys.executable, str(SCRIPTS / "score_visual_items.py"),
         flag, str(request_path), "--output", str(output_path)],
        capture_output=True, text=True,
    )
    return proc.returncode, proc.stdout + proc.stderr, output_path


def _assert_scorer_rejects(flag: str, payload, note: str = "") -> None:
    """Every rejection must be non-zero AND leave no report behind."""
    with tempfile.TemporaryDirectory() as tmp:
        code, out, output_path = _run_scorer(Path(tmp), flag, payload)
        assert code != 0, (note, payload, out)
        assert not output_path.exists(), (note, payload, "wrote a report anyway")
        assert "reuse:" not in out, (note, payload, "printed a selection", out)


def test_malformed_single_request_never_reaches_the_scorer() -> None:
    """An unconstrained request scores generic assets at 90 — it must be refused.

    Root cause: `overlap_score()` returns 1.0 when the request contributes no
    terms (nothing asked for is trivially covered), so `{}` earned FULL semantic
    credit, cleared the reuse floor and selected `sun.asset.logo` at 90.0 with
    exit 0. Absence of `intent` is unscorable, not merely weak.
    """
    _assert_scorer_rejects("--request", {}, "empty object")
    _assert_scorer_rejects("--request", [], "JSON array")
    _assert_scorer_rejects("--request", "nope", "JSON string")
    _assert_scorer_rejects("--request", _BATCH_REQUEST, "batch envelope")
    # intent missing / wrong type / blank
    _assert_scorer_rejects("--request", {"query": "a cover slide"}, "no intent")
    _assert_scorer_rejects("--request", dict(_SINGLE_REQUEST, intent="cover"), "intent str")
    _assert_scorer_rejects("--request", dict(_SINGLE_REQUEST, intent=[]), "intent empty")
    _assert_scorer_rejects("--request", dict(_SINGLE_REQUEST, intent=["  "]), "intent blank")
    _assert_scorer_rejects("--request", dict(_SINGLE_REQUEST, intent=[1, 2]), "intent non-str")
    # other fields are type-checked so a typo fails loudly instead of being ignored
    _assert_scorer_rejects("--request", dict(_SINGLE_REQUEST, tags="cover"), "tags str")
    _assert_scorer_rejects("--request", dict(_SINGLE_REQUEST, content_structure={}), "cs dict")
    _assert_scorer_rejects("--request", dict(_SINGLE_REQUEST, item_count=0), "count 0")
    _assert_scorer_rejects("--request", dict(_SINGLE_REQUEST, item_count="3"), "count str")
    _assert_scorer_rejects("--request", dict(_SINGLE_REQUEST, item_count=True), "count bool")
    _assert_scorer_rejects("--request", dict(_SINGLE_REQUEST, content_shape=3), "shape int")


def test_malformed_batch_request_never_reaches_the_scorer() -> None:
    _assert_scorer_rejects("--batch-request", _SINGLE_REQUEST, "single request")
    _assert_scorer_rejects("--batch-request", [], "JSON array")
    _assert_scorer_rejects("--batch-request", {"job_id": "j", "slides": []}, "empty slides")
    _assert_scorer_rejects("--batch-request", {"job_id": "j", "slides": {}}, "slides dict")
    _assert_scorer_rejects("--batch-request", {"job_id": "j"}, "no slides")
    _assert_scorer_rejects("--batch-request",
                           {"job_id": "j", "slides": [_SINGLE_REQUEST, "nope"]},
                           "non-object slide entry")


def test_batch_with_a_malformed_later_slide_writes_no_partial_report() -> None:
    """Validation covers the WHOLE batch before slide 1 is scored.

    Otherwise a bad slide 3 leaves a report describing slides 1-2 that
    downstream gates would read as a complete selection.
    """
    good = dict(_SINGLE_REQUEST, request_id="slide-01")
    payload = {
        "job_id": "j",
        "slides": [good, dict(good, request_id="slide-02"),
                   {"request_id": "slide-03", "intent": "cover"}],  # malformed
    }
    with tempfile.TemporaryDirectory() as tmp:
        code, out, output_path = _run_scorer(Path(tmp), "--batch-request", payload)
        assert code != 0, out
        assert not output_path.exists(), "partial report written for slides 1-2"
        assert "slide-03" in out, out
        # Nothing was scored, so no per-slide decision line was printed.
        assert "slide-01:" not in out, out


def test_valid_request_and_batch_modes_still_write_their_reports() -> None:
    """Negative control: the guards must not break the two legitimate shapes."""
    with tempfile.TemporaryDirectory() as tmp:
        code, out, output_path = _run_scorer(Path(tmp), "--request", _SINGLE_REQUEST)
        assert code == 0, out
        report = json.loads(output_path.read_text(encoding="utf-8"))
        assert "decision" in report and "slides" not in report, report
        assert report["decision"]["action"] in ("reuse", "text-only"), report

    with tempfile.TemporaryDirectory() as tmp:
        code, out, output_path = _run_scorer(Path(tmp), "--batch-request", _BATCH_REQUEST)
        assert code == 0, out
        report = json.loads(output_path.read_text(encoding="utf-8"))
        assert len(report["slides"]) == 1, report
        assert report["slides"][0]["decision"]["action"] in ("reuse", "text-only"), report


# --------------------------------------------------------------------------- #
# export_pptx — selection inputs bound the export cache
# --------------------------------------------------------------------------- #
def _selection_report(action: str, item_id, rejected: list) -> dict:
    return {
        "job_id": "j",
        "generated_at": "2026-07-21T00:00:00Z",  # varies per run; must NOT count
        "rejected_items": rejected,
        "slides": [{"request_id": "slide-01",
                    "decision": {"action": action, "item_id": item_id}}],
    }


def test_selection_inputs_participate_in_the_export_fingerprint() -> None:
    """A changed chosen/rejected item set must invalidate stale export reuse.

    `html_sha` cannot bound this: a diagnostic re-score can change the decision
    or the rejected set while leaving byte-identical deck HTML.
    """
    with tempfile.TemporaryDirectory() as tmp:
        run_dir = Path(tmp)
        html = run_dir / "deck.html"
        html.write_text("<deck-stage></deck-stage>", encoding="utf-8")
        analysis = run_dir / "analysis"
        analysis.mkdir()
        report_path = analysis / "selection-report.json"

        def identity(payload: dict):
            report_path.write_text(json.dumps(payload), encoding="utf-8")
            return export_pptx.selection_identity(html)

        base = identity(_selection_report("text-only", None, []))

        # Only `generated_at` differs -> identity must be stable, or every
        # re-score would kill the cache and the reuse path would be dead code.
        same = _selection_report("text-only", None, [])
        same["generated_at"] = "2026-07-22T09:30:00Z"
        assert identity(same) == base, "timestamp must not change identity"

        # A different chosen item, and a different rejected set, must both move it.
        assert identity(_selection_report("reuse", "sun.component.x", [])) != base
        assert identity(_selection_report("text-only", None, ["sun.component.x"])) != base

        # Rejected-item ORDER is not a semantic change.
        a = identity(_selection_report("text-only", None, ["b", "a"]))
        b = identity(_selection_report("text-only", None, ["a", "b"]))
        assert a == b, "rejected_items order must not change identity"

        # An unreadable report must not silently authorise reuse.
        report_path.write_text("{not json", encoding="utf-8")
        assert export_pptx.selection_identity(html) == ["unreadable-selection-report"]

        # No selection report at all (non-slide-job export) stays None, so those
        # runs keep their existing cache behaviour unchanged.
        report_path.unlink()
        assert export_pptx.selection_identity(html) is None


# --------------------------------------------------------------------------- #
# scaffold_extraction — ID gating + analysis-dir coexistence
# --------------------------------------------------------------------------- #
import scaffold_extraction as scaffold_ex


def test_scaffold_rejects_docling_draft_ids() -> None:
    # Every placeholder analyze_with_docling.py can mint must be rejected so it
    # can never become a stable identity without a human rename.
    for bad in ("picture-p1-1", "figure-p2-3", "table-p10-1", "chart-px-1",
                "form-p3-2"):
        assert scaffold_ex._DOCLING_DRAFT_ID.match(bad), bad
    # Real semantic names (and the suggested renames) must pass.
    for ok in ("metric-card", "salary-table", "org-chart", "picture-frame",
               "table-of-contents"):
        assert not scaffold_ex._DOCLING_DRAFT_ID.match(ok), ok


def test_scaffold_still_rejects_positional_ids() -> None:
    # The pre-existing positional gate is unchanged by the refactor.
    for bad in ("page-01", "slide-3-full", "42", "top-left", "center"):
        assert scaffold_ex._BANNED_ID.match(bad), bad
    for ok in ("left-rail", "top-banner", "metric-card"):
        assert not scaffold_ex._BANNED_ID.match(ok), ok


def test_analyze_with_docling_emits_only_draft_ids() -> None:
    # Guard the contract between the two scripts: every candidate id the analyzer
    # produces must be caught by the scaffold draft gate.
    import analyze_with_docling as awd
    els = [{"page": p, "label": lbl, "text": "",
            "region": {"x": 0.1, "y": 0.1, "width": 0.2, "height": 0.2,
                       "unit": "normalized"}}
           for p, lbl in [(1, "picture"), (2, "table"), (10, "figure"),
                          (3, "form")]]
    items = awd.build_candidates(els, "demo", "component", None)
    assert items, "expected candidates from figure-like labels"
    for it in items:
        assert scaffold_ex._DOCLING_DRAFT_ID.match(it["item_id"]), it["item_id"]


def test_analyze_with_docling_skips_chart_candidates() -> None:
    import analyze_with_docling as awd

    els = [
        {"page": 1, "label": "chart", "text": "Pie chart",
         "region": {"x": 0.1, "y": 0.1, "width": 0.3,
                    "height": 0.3, "unit": "normalized"}},
        {"page": 1, "label": "picture", "text": "Reusable card",
         "region": {"x": 0.5, "y": 0.1, "width": 0.3,
                    "height": 0.3, "unit": "normalized"}},
    ]

    items = awd.build_candidates(els, "demo", "component", None)

    assert [item["item_id"] for item in items] == ["picture-p1-1"]


def test_analyze_with_docling_filters_tiny_candidates() -> None:
    import analyze_with_docling as awd
    els = [
        {"page": 1, "label": "picture", "text": "",
         "region": {"x": 0.1, "y": 0.1, "width": 0.02,
                    "height": 0.02, "unit": "normalized"}},
        {"page": 1, "label": "picture", "text": "",
         "region": {"x": 0.2, "y": 0.2, "width": 0.2,
                    "height": 0.2, "unit": "normalized"}},
    ]
    items = awd.build_candidates(els, "demo", "component", None)
    assert [item["item_id"] for item in items] == ["picture-p1-1"]


def test_analyze_with_docling_pdf_fallback_groups_rows() -> None:
    import analyze_with_docling as awd

    atoms = [
        {"kind": "text", "text": "2. header",
         "region": {"x": 0.1, "y": 0.08, "width": 0.4, "height": 0.02, "unit": "normalized"}},
        {"kind": "drawing", "text": "",
         "region": {"x": 0.19, "y": 0.14, "width": 0.20, "height": 0.22, "unit": "normalized"}},
        {"kind": "drawing", "text": "",
         "region": {"x": 0.41, "y": 0.14, "width": 0.20, "height": 0.22, "unit": "normalized"}},
        {"kind": "text", "text": "01 LOREM IPSUM",
         "region": {"x": 0.24, "y": 0.23, "width": 0.12, "height": 0.04, "unit": "normalized"}},
        {"kind": "text", "text": "GOAL KEY RESULT TASK",
         "region": {"x": 0.26, "y": 0.48, "width": 0.50, "height": 0.05, "unit": "normalized"}},
        {"kind": "text", "text": "Kết quả muốn đạt được",
         "region": {"x": 0.19, "y": 0.59, "width": 0.18, "height": 0.06, "unit": "normalized"}},
        {"kind": "text", "text": "FOUNDATION TOP1 MICROSOFT XIAOMI",
         "region": {"x": 0.27, "y": 0.78, "width": 0.45, "height": 0.09, "unit": "normalized"}},
    ]

    elements = awd.fallback_elements_from_atoms(3, atoms)
    assert [el["source"] for el in elements] == ["pymupdf-fallback"] * 3
    assert [round(el["region"]["y"], 2) for el in elements] == [0.13, 0.47, 0.77]
    items = awd.build_candidates(elements, "demo", "component", None)
    assert [item["item_id"] for item in items] == [
        "figure-p3-1",
        "figure-p3-2",
        "figure-p3-3",
    ]
    assert "PyMuPDF fallback" in items[0]["notes"]


def test_analyze_with_docling_fallback_keeps_uncovered_metric_row() -> None:
    import analyze_with_docling as awd

    fallback_rows = [
        {
            "page": 2,
            "label": "figure",
            "region": {
                "x": 0.14, "y": 0.39, "width": 0.68, "height": 0.24,
                "unit": "normalized",
            },
        },
        {
            "page": 2,
            "label": "figure",
            "region": {
                "x": 0.17, "y": 0.70, "width": 0.64, "height": 0.20,
                "unit": "normalized",
            },
        },
    ]
    existing_docling = [
        {"x": 0.151, "y": 0.413, "width": 0.151, "height": 0.208, "unit": "normalized"},
        {"x": 0.316, "y": 0.413, "width": 0.153, "height": 0.207, "unit": "normalized"},
        {"x": 0.485, "y": 0.413, "width": 0.150, "height": 0.208, "unit": "normalized"},
        {"x": 0.650, "y": 0.413, "width": 0.151, "height": 0.208, "unit": "normalized"},
        # Small arrow icons in the metric row must not suppress the broad row.
        {"x": 0.468, "y": 0.713, "width": 0.055, "height": 0.058, "unit": "normalized"},
        {"x": 0.470, "y": 0.833, "width": 0.053, "height": 0.059, "unit": "normalized"},
    ]

    kept = [
        row for row in fallback_rows
        if not awd._covered_by_existing_candidates(row["region"], existing_docling)
    ]

    assert kept == [fallback_rows[1]]


def test_analyze_with_docling_fallback_container_becomes_context() -> None:
    import analyze_with_docling as awd

    ai_visual = {
        "page": 4,
        "label": "picture",
        "text": "",
        "region": {
            "x": 0.211145, "y": 0.659546, "width": 0.551944,
            "height": 0.231766, "unit": "normalized",
        },
        "source": "docling",
    }
    broad_fallback = {
        "x": 0.100247, "y": 0.494247, "width": 0.674524,
        "height": 0.505753, "unit": "normalized",
    }
    metric_strip = {
        "x": 0.170, "y": 0.700, "width": 0.640,
        "height": 0.200, "unit": "normalized",
    }
    small_arrow = {
        "region": {
            "x": 0.470, "y": 0.833, "width": 0.053,
            "height": 0.059, "unit": "normalized",
        },
    }

    assert awd._contained_existing_candidate(broad_fallback, [ai_visual]) is ai_visual
    assert awd._covered_by_existing_candidates(broad_fallback, [ai_visual])
    assert awd._contained_existing_candidate(metric_strip, [small_arrow]) is None
    assert not awd._covered_by_existing_candidates(metric_strip, [small_arrow])

    awd._append_context_text(
        ai_visual,
        "2. CONTENT XOAY QUANH - build AI team and automation system",
    )
    items = awd.build_candidates([ai_visual], "demo", "component", None)

    assert [item["item_id"] for item in items] == ["picture-p4-1"]
    assert "CONTENT XOAY QUANH" in items[0]["semantic_intent"][0]


def test_analyze_with_docling_fallback_text_uses_reading_order() -> None:
    import analyze_with_docling as awd

    row = [
        {"kind": "text", "text": "+30%",
         "region": {"x": 0.55, "y": 0.72, "width": 0.14, "height": 0.06,
                    "unit": "normalized"}},
        {"kind": "text", "text": "Revenue",
         "region": {"x": 0.20, "y": 0.735, "width": 0.13, "height": 0.03,
                    "unit": "normalized"}},
        {"kind": "text", "text": "+30%",
         "region": {"x": 0.55, "y": 0.84, "width": 0.14, "height": 0.06,
                    "unit": "normalized"}},
        {"kind": "text", "text": "Team Size",
         "region": {"x": 0.20, "y": 0.835, "width": 0.13, "height": 0.03,
                    "unit": "normalized"}},
    ]

    assert awd._text_lines_for_row(row) == ["Revenue +30%", "Team Size +30%"]


def test_analyze_with_docling_merges_header_and_visual_rows() -> None:
    import analyze_with_docling as awd

    rows = [
        {
            "page": 5,
            "label": "figure",
            "text": "This is a contributors slide. Insert your team here.",
            "region": {"x": 0.28, "y": 0.12, "width": 0.42, "height": 0.13,
                       "unit": "normalized"},
            "source": "pymupdf-fallback",
        },
        {
            "page": 5,
            "label": "figure",
            "text": "Patrick E. Shorey Mary T. Middleton William R. Hudson",
            "region": {"x": 0.24, "y": 0.265, "width": 0.50, "height": 0.23,
                       "unit": "normalized"},
            "source": "pymupdf-fallback",
        },
    ]

    merged = awd._merge_header_visual_rows(rows)

    assert len(merged) == 1
    assert merged[0]["region"]["y"] < rows[0]["region"]["y"]
    assert merged[0]["region"]["height"] > 0.35
    assert "contributors slide" in merged[0]["text"]
    assert "Patrick" in merged[0]["text"]


def test_analyze_with_docling_icon_sheet_candidate_covers_full_glyph_grid() -> None:
    import analyze_with_docling as awd

    atoms = []
    for row in range(8):
        for col in range(8):
            atoms.append({
                "kind": "drawing",
                "text": "",
                "region": {
                    "x": 0.10 + col * 0.05,
                    "y": 0.13 + row * 0.07,
                    "width": 0.012,
                    "height": 0.018,
                    "unit": "normalized",
                },
            })

    element = awd._icon_sheet_element_from_atoms(
        1, atoms, "ICON\n1. NHUNG ICON HAY XUAT HIEN")

    assert element is not None
    assert element["region"]["x"] < 0.09
    assert element["region"]["y"] < 0.08
    assert element["region"]["width"] > 0.36
    assert element["region"]["height"] > 0.55


def test_analyze_with_docling_pdf_page_mode_survives_one_page_failure() -> None:
    import analyze_with_docling as awd

    class _Size:
        width = 100
        height = 100

    class _Page:
        size = _Size()

    class _Label:
        value = "picture"

    class _BBox:
        l = 10
        t = 10
        r = 50
        b = 50

    class _Prov:
        page_no = 1
        bbox = _BBox()

    class _Item:
        label = _Label()
        text = "Reusable visual"
        prov = [_Prov()]

    class _Doc:
        pages = {1: _Page()}

        def iterate_items(self):
            return [(_Item(), 0)]

    class _Result:
        document = _Doc()

    class _Converter:
        def __init__(self) -> None:
            self.page_ranges: list[tuple[int, int]] = []

        def convert(self, source, **kwargs):
            page_range = kwargs["page_range"]
            self.page_ranges.append(page_range)
            if page_range == (2, 2):
                raise RuntimeError("page failed")
            return _Result()

    converter = _Converter()
    old = awd._page_numbers_for_source
    awd._page_numbers_for_source = lambda source, pages: ([1, 2, 3], [])
    try:
        elements, warnings, stats = awd.analyze_source(
            converter, Path("demo.pdf"), (1, 3))
    finally:
        awd._page_numbers_for_source = old

    assert converter.page_ranges == [(1, 1), (2, 2), (3, 3)]
    assert [el["page"] for el in elements] == [1, 3]
    assert any("page 2" in warning for warning in warnings)
    assert stats["docling_mode"] == "page-by-page"
    assert stats["docling_pages_attempted"] == 3
    assert stats["docling_pages_failed"] == 1


def test_scaffold_rejects_docling_draft_without_polluting_analysis_dir() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        source = root / "source.pdf"
        source.write_text("fake", encoding="utf-8")
        output_root = root / "outputs"
        analysis = output_root / "docling-demo" / "analysis"
        analysis.mkdir(parents=True)
        request_path = root / "request.json"
        request_path.write_text(json.dumps({
            "extraction_id": "docling-demo",
            "source_path": str(source),
            "items": [{
                "item_id": "picture-p1-1",
                "slide_or_page": 1,
                "region": {"x": 0.1, "y": 0.1, "width": 0.2,
                           "height": 0.2, "unit": "normalized"},
                "requested_type": "component",
                "semantic_intent": ["picture candidate detected by Docling"],
            }],
        }), encoding="utf-8")
        history = root / "history.json"
        registry = root / "registry.json"
        history.write_text('{"attempts":[]}', encoding="utf-8")
        registry.write_text('{"items":[]}', encoding="utf-8")
        old_argv = sys.argv[:]
        try:
            sys.argv = [
                "scaffold_extraction.py", "--request", str(request_path),
                "--output-root", str(output_root), "--history", str(history),
                "--registry", str(registry),
            ]
            try:
                scaffold_ex.main()
            except SystemExit as exc:
                assert "Docling draft placeholder" in str(exc)
            else:
                raise AssertionError("expected Docling draft placeholder rejection")
        finally:
            sys.argv = old_argv
        assert analysis.exists(), "analysis/ should be preserved"
        assert not (output_root / "docling-demo" / "request.json").exists()
        assert not (output_root / "docling-demo" / "items").exists()


# --------------------------------------------------------------------------- #
# candidate_review — rename / metadata / approval (analysis-only)
# --------------------------------------------------------------------------- #
import candidate_review as crv

_EXTRACTION_SCHEMA = SCRIPTS.parent / "schemas" / "extraction-request.schema.json"


def _review_fixture(tmp: Path, candidate_id: str = "picture-p1-1") -> tuple[Path, str]:
    """Create an extractions root with one analysis run carrying a placeholder
    candidate. Returns (root, extraction_id)."""
    extraction_id = "docling-demo"
    adir = tmp / extraction_id / "analysis"
    adir.mkdir(parents=True)
    (adir / "candidate-extraction-request.json").write_text(json.dumps({
        "extraction_id": extraction_id,
        "source_path": "input/Demo.pdf",
        "items": [{
            "item_id": candidate_id,
            "slide_or_page": 1,
            "region": {"x": 0.5, "y": 0.0, "width": 0.4, "height": 0.9,
                       "unit": "normalized"},
            "object_ids": [],
            "requested_type": "component",
            "semantic_intent": ["picture candidate detected by Docling"],
            "notes": "DRAFT candidate from Docling auto-detect.",
            "replacement_for": None,
        }],
    }), encoding="utf-8")
    (adir / "page-analysis.json").write_text('{"elements": []}', encoding="utf-8")
    (adir / "docling-report.json").write_text('{"candidate_count": 1}', encoding="utf-8")
    return tmp, extraction_id


def _valid_metadata(item_id: str = "kickoff-2026-hero-visual") -> dict:
    return {
        "item_id": item_id,
        "display_name": "Kick-off 2026 hero visual",
        "requested_type": "component",
        "component_type": "hero",
        "layout_role": "full-bleed",
        "visual_summary": "A tall orange hero illustration on the right column.",
        "semantic_intent": ["kickoff hero", "goal setting cover"],
        "content_structure": ["illustration"],
        "tags": ["hero", "orange"],
        "keywords": ["kickoff", "2026"],
        "use_cases": ["cover slide"],
        "anti_use_cases": ["dense data slide"],
        "quality_notes": "",
        "retrieval_notes": "",
    }


def test_candidate_placeholder_id_cannot_be_approved() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        root, eid = _review_fixture(Path(tmp))
        # Save valid metadata but keep the Docling placeholder as item_id.
        meta = _valid_metadata(item_id="picture-p1-1")
        crv.save_review(eid, "picture-p1-1", meta, reviewer="t", root=root)
        try:
            crv.approve(eid, "picture-p1-1", reviewer="t", root=root)
        except crv.CandidateValidationError as exc:
            assert any("placeholder" in e.lower() for e in exc.errors), exc.errors
        else:
            raise AssertionError("placeholder item_id must not be approvable")
        # No approved artifact written.
        assert not (root / eid / "analysis" / "approved").exists()


def test_candidate_positional_id_rejected() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        root, eid = _review_fixture(Path(tmp))
        crv.save_review(eid, "picture-p1-1", _valid_metadata("top-left"),
                        reviewer="t", root=root)
        errors = crv.validate_review(crv.get_candidates(eid, root=root)
                                     ["candidates"][0]["review"])
        assert any("positional" in e.lower() or "generic" in e.lower() for e in errors), errors


def test_candidate_required_metadata_enforced() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        root, eid = _review_fixture(Path(tmp))
        # rename only, leave all metadata empty
        crv.save_review(eid, "picture-p1-1", {"item_id": "kickoff-hero"},
                        reviewer="t", root=root)
        try:
            crv.approve(eid, "picture-p1-1", reviewer="t", root=root)
        except crv.CandidateValidationError as exc:
            joined = " ".join(exc.errors).lower()
            assert "display name" in joined and "visual summary" in joined, exc.errors
        else:
            raise AssertionError("missing required metadata must block approval")


def test_candidate_approve_writes_schema_compatible_request() -> None:
    schema = read_text_slots.load_json(_EXTRACTION_SCHEMA)
    item_schema = schema["properties"]["items"]["items"]
    allowed = set(item_schema["properties"])
    required = set(item_schema["required"])
    top_required = set(schema["required"])

    with tempfile.TemporaryDirectory() as tmp:
        root, eid = _review_fixture(Path(tmp))
        crv.save_review(eid, "picture-p1-1", _valid_metadata(), reviewer="t", root=root)
        result = crv.approve(eid, "picture-p1-1", reviewer="t", root=root)

        approved_path = root / eid / "analysis" / "approved" / "kickoff-2026-hero-visual.extraction-request.json"
        assert approved_path.is_file(), result
        req = read_text_slots.load_json(approved_path)
        assert top_required <= set(req), req
        item = req["items"][0]
        assert required <= set(item), item
        assert set(item) <= allowed, f"extra keys not in schema: {set(item) - allowed}"
        assert item["item_id"] == "kickoff-2026-hero-visual"
        # The approved request must also pass the live scaffold gate.
        scaffold_ex.validate_request_item(item)
        # review status updated, reviewer recorded.
        cand = crv.get_candidates(eid, root=root)["candidates"][0]
        assert cand["review"]["review_status"] == "approved_for_extraction"
        assert cand["review"]["reviewer"] == "t"


def test_candidate_reject_produces_no_approved_request() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        root, eid = _review_fixture(Path(tmp))
        crv.save_review(eid, "picture-p1-1", _valid_metadata(), reviewer="t", root=root)
        crv.approve(eid, "picture-p1-1", reviewer="t", root=root)
        approved_path = root / eid / "analysis" / "approved" / "kickoff-2026-hero-visual.extraction-request.json"
        assert approved_path.is_file()
        # Rejecting must drop the stale approved artifact and flip status.
        review = crv.reject(eid, "picture-p1-1", "wrong crop", reviewer="t", root=root)
        assert review["review_status"] == "rejected"
        assert not approved_path.exists(), "reject must remove the approved request"
        # A reject with no reason is refused.
        try:
            crv.reject(eid, "picture-p1-1", "", root=root)
        except crv.CandidateError:
            pass
        else:
            raise AssertionError("reject without a reason must fail")


def test_candidate_review_preserves_analysis_files() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        root, eid = _review_fixture(Path(tmp))
        adir = root / eid / "analysis"
        before = (adir / "candidate-extraction-request.json").read_text(encoding="utf-8")
        page_before = (adir / "page-analysis.json").read_text(encoding="utf-8")
        crv.save_review(eid, "picture-p1-1", _valid_metadata(), reviewer="t", root=root)
        crv.approve(eid, "picture-p1-1", reviewer="t", root=root)
        assert (adir / "candidate-extraction-request.json").read_text(encoding="utf-8") == before
        assert (adir / "page-analysis.json").read_text(encoding="utf-8") == page_before
        assert (adir / "docling-report.json").exists()


def test_candidate_pdf_preview_is_generated_and_reused() -> None:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return
    with tempfile.TemporaryDirectory() as tmp:
        tmpp = Path(tmp)
        source = tmpp / "source.pdf"
        doc = fitz.open()
        page = doc.new_page(width=200, height=100)
        page.draw_rect(fitz.Rect(50, 20, 150, 80), color=(1, 0, 0), fill=(1, 0.8, 0.7))
        page.insert_text((58, 55), "Preview", fontsize=16, color=(0, 0, 0))
        doc.save(source)
        doc.close()

        root = tmpp / "ext"
        eid = "docling-preview-demo"
        adir = root / eid / "analysis"
        adir.mkdir(parents=True)
        (adir / "candidate-extraction-request.json").write_text(json.dumps({
            "extraction_id": eid,
            "source_path": str(source),
            "items": [{
                "item_id": "picture-p1-1",
                "slide_or_page": 1,
                "region": {"x": 0.25, "y": 0.2, "width": 0.5, "height": 0.6,
                           "unit": "normalized"},
                "object_ids": [],
                "requested_type": "component",
                "semantic_intent": ["preview candidate"],
                "notes": "preview smoke",
                "replacement_for": None,
            }],
        }), encoding="utf-8")

        result = crv.get_candidates(eid, root=root)
        preview = result["candidates"][0]["preview"]
        assert preview["status"] == "ready", preview
        png = adir / "previews" / "picture-p1-1.png"
        assert png.is_file(), preview
        assert png.read_bytes().startswith(b"\x89PNG"), "preview must be a PNG"
        first_mtime = png.stat().st_mtime_ns

        second = crv.get_candidates(eid, root=root)["candidates"][0]["preview"]
        assert second["path"] == preview["path"]
        assert png.stat().st_mtime_ns == first_mtime, "existing preview should be reused"


def test_candidate_preview_unavailable_for_non_pdf_source() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        tmpp = Path(tmp)
        root = tmpp / "ext"
        eid = "docling-preview-fallback"
        adir = root / eid / "analysis"
        adir.mkdir(parents=True)
        (adir / "candidate-extraction-request.json").write_text(json.dumps({
            "extraction_id": eid,
            "source_path": "source.pptx",
            "items": [{
                "item_id": "picture-p1-1",
                "slide_or_page": 1,
                "region": {"x": 0.25, "y": 0.2, "width": 0.5, "height": 0.6,
                           "unit": "normalized"},
                "object_ids": [],
                "requested_type": "component",
                "semantic_intent": ["preview candidate"],
                "notes": "preview fallback",
                "replacement_for": None,
            }],
        }), encoding="utf-8")

        preview = crv.get_candidates(eid, root=root)["candidates"][0]["preview"]
        assert preview["status"] == "unavailable", preview
        assert "PDF sources only" in preview["reason"]
        assert not (adir / "previews").exists()


def test_candidate_preview_unavailable_for_malformed_region() -> None:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return
    with tempfile.TemporaryDirectory() as tmp:
        tmpp = Path(tmp)
        source = tmpp / "source.pdf"
        doc = fitz.open()
        doc.new_page(width=200, height=100)
        doc.save(source)
        doc.close()

        root = tmpp / "ext"
        eid = "docling-preview-bad-region"
        adir = root / eid / "analysis"
        adir.mkdir(parents=True)
        (adir / "candidate-extraction-request.json").write_text(json.dumps({
            "extraction_id": eid,
            "source_path": str(source),
            "items": [{
                "item_id": "picture-p1-1",
                "slide_or_page": 1,
                "region": {"x": 0.25, "unit": "normalized"},
                "object_ids": [],
                "requested_type": "component",
                "semantic_intent": ["preview candidate"],
                "notes": "bad region",
                "replacement_for": None,
            }],
        }), encoding="utf-8")

        preview = crv.get_candidates(eid, root=root)["candidates"][0]["preview"]
        assert preview["status"] == "unavailable", preview
        assert "Preview render failed" in preview["reason"]
        assert not (adir / "previews").exists()


def test_candidate_review_does_not_touch_registry_or_library() -> None:
    # candidate_review must only write under the analysis dir; the real registry,
    # compact registry, history, and library must be byte-identical afterwards.
    repo = SCRIPTS.parents[1]
    watched = [
        repo / "slide-system" / "registries" / "visual-library.json",
        repo / "slide-system" / "registries" / "visual-library-compact.json",
        repo / "slide-system" / "registries" / "extraction-history.json",
    ]
    before = {p: p.read_bytes() for p in watched if p.exists()}
    with tempfile.TemporaryDirectory() as tmp:
        root, eid = _review_fixture(Path(tmp))
        crv.save_review(eid, "picture-p1-1", _valid_metadata(), reviewer="t", root=root)
        crv.approve(eid, "picture-p1-1", reviewer="t", root=root)
        crv.reject(eid, "picture-p1-1", "redo", reviewer="t", root=root)
    after = {p: p.read_bytes() for p in watched if p.exists()}
    assert before == after, "candidate review must not mutate registry/history/library"


def test_candidate_invalid_extraction_id_rejected() -> None:
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        for bad in ("../escape", "..", "a/b", "/etc", "bad id"):
            try:
                crv.get_candidates(bad, root=root)
            except crv.CandidateError:
                pass
            else:
                raise AssertionError(f"invalid extraction id must be rejected: {bad!r}")


def test_candidate_editing_resets_approval() -> None:
    # Editing an approved candidate must revert it to pending and drop the stale
    # approved request, so an approval never outlives the metadata it was built on.
    with tempfile.TemporaryDirectory() as tmp:
        root, eid = _review_fixture(Path(tmp))
        crv.save_review(eid, "picture-p1-1", _valid_metadata(), reviewer="t", root=root)
        crv.approve(eid, "picture-p1-1", reviewer="t", root=root)
        approved_path = root / eid / "analysis" / "approved" / "kickoff-2026-hero-visual.extraction-request.json"
        assert approved_path.is_file()
        crv.save_review(eid, "picture-p1-1", {"visual_summary": "edited"},
                        reviewer="t", root=root)
        assert not approved_path.exists(), "editing must drop the stale approved request"
        cand = crv.get_candidates(eid, root=root)["candidates"][0]
        assert cand["review"]["review_status"] == "pending"


def test_candidate_multiple_approvals_scaffold_without_collision() -> None:
    # Regression: every approved request from one run must get its own scaffold
    # extraction id, else the second candidate fails with "already exists".
    import scaffold_extraction as sx
    with tempfile.TemporaryDirectory() as tmp:
        tmpp = Path(tmp)
        source = tmpp / "Demo.pdf"
        source.write_bytes(b"%PDF-1.4 fake source")
        eid = "docling-demo"
        root = tmpp / "ext"
        adir = root / eid / "analysis"
        adir.mkdir(parents=True)
        (adir / "candidate-extraction-request.json").write_text(json.dumps({
            "extraction_id": eid,
            "source_path": str(source),
            "items": [
                {"item_id": f"picture-p1-{i}", "slide_or_page": 1,
                 "region": {"x": 0.1 * i, "y": 0.1, "width": 0.3, "height": 0.3,
                            "unit": "normalized"},
                 "object_ids": [], "requested_type": "component",
                 "semantic_intent": ["pic"], "replacement_for": None}
                for i in (1, 2)
            ],
        }), encoding="utf-8")

        for i in (1, 2):
            crv.save_review(eid, f"picture-p1-{i}", _valid_metadata(f"hero-{i}"),
                            reviewer="t", root=root)
            crv.approve(eid, f"picture-p1-{i}", reviewer="t", root=root)

        req1 = read_text_slots.load_json(adir / "approved" / "hero-1.extraction-request.json")
        req2 = read_text_slots.load_json(adir / "approved" / "hero-2.extraction-request.json")
        assert req1["extraction_id"] == "docling-demo-hero-1", req1["extraction_id"]
        assert req1["extraction_id"] != req2["extraction_id"]

        out_root = tmpp / "out"
        hist = tmpp / "history.json"; hist.write_text('{"attempts":[]}', encoding="utf-8")
        reg = tmpp / "registry.json"; reg.write_text('{"items":[]}', encoding="utf-8")
        for name in ("hero-1", "hero-2"):
            req = adir / "approved" / f"{name}.extraction-request.json"
            old_argv = sys.argv[:]
            sys.argv = ["scaffold_extraction.py", "--request", str(req),
                        "--output-root", str(out_root), "--history", str(hist),
                        "--registry", str(reg)]
            try:
                assert sx.main() == 0, name
            finally:
                sys.argv = old_argv
        # Both scaffolded into separate, non-colliding output dirs.
        assert (out_root / "docling-demo-hero-1" / "items" / "hero-1").is_dir()
        assert (out_root / "docling-demo-hero-2" / "items" / "hero-2").is_dir()


def test_candidate_rename_removes_old_approved_artifact() -> None:
    # Renaming an approved candidate must not orphan the old item_id's request.
    with tempfile.TemporaryDirectory() as tmp:
        root, eid = _review_fixture(Path(tmp))
        crv.save_review(eid, "picture-p1-1", _valid_metadata("hero-a"), reviewer="t", root=root)
        crv.approve(eid, "picture-p1-1", reviewer="t", root=root)
        old = root / eid / "analysis" / "approved" / "hero-a.extraction-request.json"
        assert old.is_file()
        meta = _valid_metadata("hero-b")
        crv.save_review(eid, "picture-p1-1", meta, reviewer="t", root=root)
        assert not old.exists(), "old-name approved request must be removed on rename"
        crv.approve(eid, "picture-p1-1", reviewer="t", root=root)
        new = root / eid / "analysis" / "approved" / "hero-b.extraction-request.json"
        assert new.is_file() and not old.exists()


def test_auto_stage_candidates_creates_reviewable_draft() -> None:
    import importlib
    import build_component_catalog as bcc

    asc = importlib.import_module("auto_stage_candidates")
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return

    with tempfile.TemporaryDirectory() as tmp:
        tmpp = Path(tmp)
        source = tmpp / "Kickoff 2026.pdf"
        doc = fitz.open()
        page = doc.new_page(width=300, height=180)
        page.draw_rect(fitz.Rect(45, 35, 245, 135), color=(1, 0.2, 0.05),
                       fill=(1, 0.85, 0.76))
        page.insert_text((70, 90), "Kickoff Hero", fontsize=20, color=(0, 0, 0))
        doc.save(source)
        doc.close()

        output_root = tmpp / "component-extractions"
        eid = "docling-auto-demo"
        adir = output_root / eid / "analysis"
        adir.mkdir(parents=True)
        (adir / "candidate-extraction-request.json").write_text(json.dumps({
            "extraction_id": eid,
            "source_path": str(source),
            "items": [{
                "item_id": "picture-p1-1",
                "slide_or_page": 1,
                "region": {"x": 0.1, "y": 0.12, "width": 0.78, "height": 0.72,
                           "unit": "normalized"},
                "object_ids": [],
                "requested_type": "component",
                "semantic_intent": ["kickoff hero detected by Docling"],
                "notes": "Auto-stage this detected hero visual into Draft.",
                "replacement_for": None,
            }],
        }), encoding="utf-8")

        hist = tmpp / "history.json"
        hist.write_text('{"attempts":[]}', encoding="utf-8")
        reg = tmpp / "registry.json"
        reg.write_text('{"items":[]}', encoding="utf-8")
        summary = asc.stage_run(
            eid,
            root=output_root,
            output_root=output_root,
            history=hist,
            registry=reg,
            rebuild_catalog=False,
        )

        assert summary["staged"] == 1, summary
        staged = summary["items"][0]
        item_id = staged["item_id"]
        assert item_id != "picture-p1-1"
        assert not scaffold_ex._DOCLING_DRAFT_ID.match(item_id)
        item_dir = output_root / staged["extraction_id"] / "items" / item_id
        mapping = read_text_slots.load_json(item_dir / "mapping.json")
        assert mapping["status"] == "staging"
        assert mapping["candidate_stable_id"].startswith("sun.component.")
        assert mapping["source"]["candidate_id"] == "picture-p1-1"
        assert mapping["review"]["mode"] == "auto-staged"
        assert (item_dir / "artifact" / "visual.svg").is_file()
        assert (item_dir / "artifact" / "text-slots.json").is_file()
        assert (item_dir / "evidence" / "source-with-text.svg").is_file()
        assert (item_dir / "preview" / "thumbnail.png").is_file()

        catalog_path = tmpp / "catalog-data.json"
        old_argv = sys.argv[:]
        sys.argv = [
            "build_component_catalog.py",
            "--registry", str(reg),
            "--extractions", str(output_root),
            "--output", str(catalog_path),
        ]
        try:
            assert bcc.main() == 0
        finally:
            sys.argv = old_argv
        catalog = read_text_slots.load_json(catalog_path)
        draft = next(item for item in catalog["items"]
                     if item["id"] == mapping["candidate_stable_id"])
        assert draft["status"] == "staging"
        assert draft["publish_readiness"]["ready"], draft["publish_readiness"]
        assert draft["images"], "Draft must have a visual preview for final review"
        assert draft["component_type"] == "card"
        assert draft["layout_role"]
        assert draft["keywords"]
        assert draft["use_cases"]

        dup_eid = "docling-auto-demo-duplicate"
        dup_adir = output_root / dup_eid / "analysis"
        dup_adir.mkdir(parents=True)
        (dup_adir / "candidate-extraction-request.json").write_text(json.dumps({
            "extraction_id": dup_eid,
            "source_path": str(source),
            "items": [{
                "item_id": "picture-p1-1",
                "slide_or_page": 1,
                "region": {"x": 0.1, "y": 0.12, "width": 0.78, "height": 0.72,
                           "unit": "normalized"},
                "object_ids": [],
                "requested_type": "component",
                "semantic_intent": ["same hero detected by a later Docling run"],
                "notes": "This region is already staged and must not duplicate.",
                "replacement_for": None,
            }],
        }), encoding="utf-8")
        duplicate = asc.stage_run(
            dup_eid,
            root=output_root,
            output_root=output_root,
            history=hist,
            registry=reg,
            rebuild_catalog=False,
            build_artifacts=False,
        )
        assert duplicate["staged"] == 0, duplicate
        assert duplicate["skipped"] == 1, duplicate
        assert duplicate["items"][0]["status"] == "already_staged_region"
        assert duplicate["items"][0]["stable_id"] == mapping["candidate_stable_id"]

        second = asc.stage_run(
            eid,
            root=output_root,
            output_root=output_root,
            history=hist,
            registry=reg,
            rebuild_catalog=False,
            build_artifacts=False,
        )
        assert second["staged"] == 0, second
        assert second["skipped"] == 1, second
        assert second["items"][0]["status"] in {"already_staged", "already_staged_region"}


def test_auto_stage_skips_chart_candidates_from_existing_analysis() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    assert asc._auto_stage_skip_reason({"item_id": "chart-p1-1"})
    assert asc._auto_stage_skip_reason({
        "item_id": "picture-p1-1",
        "semantic_intent": ["pie chart candidate detected by Docling"],
    })
    assert asc._auto_stage_skip_reason({
        "item_id": "picture-p1-1",
        "semantic_intent": ["org chart radial team structure"],
    }) is None

    with tempfile.TemporaryDirectory() as tmp:
        tmpp = Path(tmp)
        output_root = tmpp / "component-extractions"
        eid = "docling-chart-demo"
        adir = output_root / eid / "analysis"
        adir.mkdir(parents=True)
        (adir / "candidate-extraction-request.json").write_text(json.dumps({
            "extraction_id": eid,
            "source_path": str(tmpp / "source.pdf"),
            "items": [{
                "item_id": "chart-p1-1",
                "slide_or_page": 1,
                "region": {"x": 0.1, "y": 0.1, "width": 0.3,
                           "height": 0.3, "unit": "normalized"},
                "object_ids": [],
                "requested_type": "component",
                "semantic_intent": ["chart candidate detected by Docling"],
                "notes": "DRAFT candidate from Docling auto-detect.",
                "replacement_for": None,
            }],
        }), encoding="utf-8")

        hist = tmpp / "history.json"
        reg = tmpp / "registry.json"
        reg.write_text('{"items":[]}', encoding="utf-8")
        summary = asc.stage_run(
            eid,
            root=output_root,
            output_root=output_root,
            history=hist,
            registry=reg,
            rebuild_catalog=False,
            build_artifacts=False,
        )

        assert summary["staged"] == 0, summary
        assert summary["skipped"] == 1, summary
        assert summary["items"][0]["reason"] == "chart candidates are skipped by auto-detect"
        assert not (adir / "approved").exists()


def test_auto_stage_skips_duplicate_component_patterns_across_pages() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    first = {
        "item_id": "picture-p3-1",
        "slide_or_page": 3,
        "region": {"x": 0.05, "y": 0.08, "width": 0.42,
                   "height": 0.32, "unit": "normalized"},
        "semantic_intent": ["Goal setting title card with highlighted subtitle"],
    }
    second = {
        "item_id": "picture-p4-1",
        "slide_or_page": 4,
        "region": {"x": 0.052, "y": 0.081, "width": 0.58,
                   "height": 0.29, "unit": "normalized"},
        "semantic_intent": ["Check-in title card with highlighted subtitle"],
    }
    assert asc._duplicate_pattern_signature("input/goal-setting.pdf", first) == (
        asc._duplicate_pattern_signature("input/goal-setting.pdf", second)
    )
    same_page_neighbour = {
        **first,
        "item_id": "picture-p3-2",
        "region": {"x": 0.55, "y": 0.08, "width": 0.42,
                   "height": 0.32, "unit": "normalized"},
    }
    assert asc._duplicate_pattern_signature("input/goal-setting.pdf", first) != (
        asc._duplicate_pattern_signature("input/goal-setting.pdf", same_page_neighbour)
    )

    with tempfile.TemporaryDirectory() as tmp:
        tmpp = Path(tmp)
        source = tmpp / "Goal Setting 2026.pptx"
        source.write_bytes(b"fake pptx source for hashing only")
        output_root = tmpp / "component-extractions"
        eid = "docling-duplicate-pattern-demo"
        adir = output_root / eid / "analysis"
        adir.mkdir(parents=True)
        (adir / "candidate-extraction-request.json").write_text(json.dumps({
            "extraction_id": eid,
            "source_path": str(source),
            "items": [
                {
                    **first,
                    "object_ids": [],
                    "requested_type": "component",
                    "notes": "Year-end evaluation title card",
                    "replacement_for": None,
                },
                {
                    **second,
                    "object_ids": [],
                    "requested_type": "component",
                    "notes": "Quarterly check-in title card",
                    "replacement_for": None,
                },
            ],
        }), encoding="utf-8")

        hist = tmpp / "history.json"
        reg = tmpp / "registry.json"
        reg.write_text('{"items":[]}', encoding="utf-8")
        summary = asc.stage_run(
            eid,
            root=output_root,
            output_root=output_root,
            history=hist,
            registry=reg,
            rebuild_catalog=False,
            build_artifacts=False,
        )

        assert summary["staged"] == 1, summary
        assert summary["skipped"] == 1, summary
        duplicate = summary["items"][1]
        assert duplicate["status"] == "skipped_duplicate_pattern"
        assert duplicate["duplicate_of_candidate_id"] == "picture-p3-1"


def test_auto_stage_cli_reads_analysis_from_output_root() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    with tempfile.TemporaryDirectory() as tmp:
        tmpp = Path(tmp)
        output_root = tmpp / "component-extractions"
        eid = "docling-cli-demo"
        adir = output_root / eid / "analysis"
        adir.mkdir(parents=True)
        source = tmpp / "source.pdf"
        source.write_bytes(b"%PDF-1.4\n% test source\n")
        (adir / "candidate-extraction-request.json").write_text(json.dumps({
            "extraction_id": eid,
            "source_path": str(source),
            "items": [{
                "item_id": "picture-p1-1",
                "slide_or_page": 1,
                "region": {"x": 0.1, "y": 0.12, "width": 0.78, "height": 0.72,
                           "unit": "normalized"},
                "object_ids": [],
                "requested_type": "component",
                "semantic_intent": ["cli output root hero detected by Docling"],
                "notes": "Auto-stage this detected hero visual into Draft.",
                "replacement_for": None,
            }],
        }), encoding="utf-8")
        hist = tmpp / "history.json"
        reg = tmpp / "registry.json"
        reg.write_text('{"items":[]}', encoding="utf-8")

        rc = asc.main([
            eid,
            "--output-root", str(output_root),
            "--history", str(hist),
            "--registry", str(reg),
            "--no-catalog",
            "--no-artifacts",
        ])

        assert rc == 0, "CLI must read analysis from --output-root"
        assert hist.is_file(), "CLI should initialize a missing custom history file"
        run_dirs = [p for p in output_root.iterdir() if p.is_dir() and p.name != eid]
        assert run_dirs, "CLI should create a staged extraction dir"


def test_auto_stage_decomposes_large_cards_as_layout_rows() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "goal-card"
        item.mkdir(parents=True)
        (item / "mapping.json").write_text(_json.dumps({
            "component_type": "card",
            "source": {"region": {"width": 0.64, "height": 0.80}},
        }), encoding="utf-8")
        assert asc._decompose_mode(item) == "layout-row-groups"

        (item / "mapping.json").write_text(_json.dumps({
            "component_type": "card",
            "source": {"region": {"width": 0.20, "height": 0.20}},
        }), encoding="utf-8")
        assert asc._decompose_mode(item) is None

        (item / "mapping.json").write_text(_json.dumps({
            "component_type": "strip",
            "source": {"region": {"width": 0.20, "height": 0.20}},
        }), encoding="utf-8")
        assert asc._decompose_mode(item) == "cards"


def test_auto_stage_decomposes_tables_and_broad_visuals_as_layout_rows() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    with tempfile.TemporaryDirectory() as tmp:
        item = Path(tmp) / "items" / "compound-region"
        item.mkdir(parents=True)

        def write_mapping(component_type: str, width: float, height: float) -> None:
            (item / "mapping.json").write_text(_json.dumps({
                "component_type": component_type,
                "source": {"region": {"width": width, "height": height}},
            }), encoding="utf-8")

        write_mapping("table", 0.72, 0.38)
        assert asc._decompose_mode(item) == "layout-row-groups"
        write_mapping("visual", 0.55, 0.23)
        assert asc._decompose_mode(item) == "layout-row-groups"
        write_mapping("component", 0.55, 0.23)
        assert asc._decompose_mode(item) == "layout-row-groups"
        write_mapping("visual", 0.30, 0.20)
        assert asc._decompose_mode(item) is None


def test_auto_stage_semantic_ids_fallback_avoids_full_source_slug() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    source = "input/SUN.SLIDE.pdf"
    item_a = {
        "item_id": "picture-p20-1",
        "slide_or_page": 20,
        "semantic_intent": ["picture candidate detected by Docling"],
        "notes": "DRAFT candidate from Docling auto-detect. Rename item_id.",
    }
    item_b = {
        "item_id": "picture-p21-1",
        "slide_or_page": 21,
        "semantic_intent": ["picture candidate detected by Docling"],
        "notes": "DRAFT candidate from Docling auto-detect. Rename item_id.",
    }
    goal_card = {
        "item_id": "picture-p9-1",
        "slide_or_page": 9,
        "region": {"x": 0.1, "y": 0.1, "width": 0.25, "height": 0.35,
                   "unit": "normalized"},
        "semantic_intent": ["picture candidate detected by Docling"],
        "notes": "DRAFT candidate from Docling auto-detect. Rename item_id.",
    }

    used: set[str] = set()
    id_a = asc.semantic_item_id(source, item_a, used)
    id_b = asc.semantic_item_id(source, item_b, used)
    goal_id = asc.semantic_item_id("input/Kick_off_GOAL_SETTING_2026-2.pdf", goal_card, set())

    assert id_a == "detected-icon-1"
    assert id_b == "detected-icon-1-2"
    assert goal_id == "goal-setting-card-1"
    assert id_a != id_b
    assert not scaffold_ex._DOCLING_DRAFT_ID.match(id_a)
    assert not scaffold_ex._DOCLING_DRAFT_ID.match(id_b)
    assert "kick-off" not in goal_id and "2026" not in goal_id


def test_auto_stage_semantic_ids_use_page_context_before_source_fallback() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    item = {
        "item_id": "picture-p7-1",
        "slide_or_page": 7,
        "region": {"x": 0.14, "y": 0.41, "width": 0.70, "height": 0.50,
                   "unit": "normalized"},
        "region_text": "Patrick E. Shorey\nMary T. Middleton",
        "page_text": "This is a contributors slide. Insert your team here.",
        "semantic_intent": ["picture candidate detected by Docling"],
        "notes": "DRAFT candidate from Docling auto-detect. Rename item_id.",
    }

    item_id = asc.semantic_item_id("input/Sun.Presentation.pdf", item, set())

    assert item_id == "contributors-team-visual", item_id
    assert not item_id.startswith("source-")


def test_auto_stage_semantic_ids_do_not_emit_source_visual_for_generic_pdf() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    item = {
        "item_id": "figure-p16-1",
        "slide_or_page": 16,
        "region": {"x": 0.04, "y": 0.2, "width": 0.91, "height": 0.52,
                   "unit": "normalized"},
        "semantic_intent": ["figure candidate detected by PyMuPDF fallback"],
        "notes": "DRAFT candidate from PyMuPDF fallback auto-detect.",
    }

    item_id = asc.semantic_item_id("input/Sun.Presentation.pdf", item, set())

    assert item_id == "detected-visual-1", item_id
    assert not item_id.startswith("source-")


def test_auto_stage_semantic_ids_use_region_text_before_source_name() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    source = "input/GUIDLINE_PRESENTATION_SUN.pdf"
    role_card = {
        "item_id": "picture-p2-2",
        "slide_or_page": 2,
        "region": {"x": 0.15, "y": 0.41, "width": 0.15, "height": 0.2,
                   "unit": "normalized"},
        "region_text": (
            "Chuyen muc tieu cong ty thanh huong di ro rang\n"
            "TRANSLATOR\n"
            "cho team"
        ),
    }
    level_strip = {
        "item_id": "picture-p2-1",
        "slide_or_page": 2,
        "region": {"x": 0.15, "y": 0.16, "width": 0.66, "height": 0.18,
                   "unit": "normalized"},
        "region_text": "AI Coding Assistants\nLevel 1\nAgent Networks\nLevel 4",
    }

    role_id = asc.semantic_item_id(source, role_card, set())
    strip_id = asc.semantic_item_id(source, level_strip, set())
    metadata = asc.metadata_for(source, role_card, role_id)

    assert role_id == "translator-card"
    assert strip_id == "ai-coding-assistants-levels-strip"
    assert "guidline" not in role_id
    assert metadata["component_type"] == "card"
    assert metadata["keywords"][:2] == ["translator", "card"]


def test_auto_stage_semantic_ids_translate_vietnamese_hints_to_english() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    recruitment_item = {
        "item_id": "picture-p5-1",
        "slide_or_page": 5,
        "region": {"x": 0.1, "y": 0.1, "width": 0.2, "height": 0.3,
                   "unit": "normalized"},
        "region_text": "HIEU MUC TIEU TUYEN DUNG",
    }
    team_item = {
        "item_id": "picture-p4-2",
        "slide_or_page": 4,
        "region": {"x": 0.1, "y": 0.1, "width": 0.5, "height": 0.3,
                   "unit": "normalized"},
        "region_text": "XAY DUNG DOI NGU",
    }

    item_id = asc.semantic_item_id("input/interview-workshop.pdf", recruitment_item, set())
    team_id = asc.semantic_item_id("input/GUIDLINE_PRESENTATION_SUN.pdf", team_item, set())
    team_meta = asc.metadata_for("input/GUIDLINE_PRESENTATION_SUN.pdf", team_item, team_id)

    assert item_id == "recruitment-goal-card"
    assert "tuyen" not in item_id and "dung" not in item_id
    assert item_id.isascii()
    assert team_id == "team-visual"
    assert "xay" not in team_id and "ngu" not in team_id
    assert "xay" not in team_meta["keywords"]
    assert "dung" not in team_meta["keywords"]


def test_auto_stage_semantic_ids_translate_salary_benefit_vietnamese() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    salary_item = {
        "item_id": "picture-p2-1",
        "slide_or_page": 2,
        "region": {"x": 0.1, "y": 0.1, "width": 0.35, "height": 0.35,
                   "unit": "normalized"},
        "region_text": "Lương phúc lợi\nQuyền lợi dài hạn",
        "semantic_intent": ["picture candidate detected by Docling"],
    }
    investment_item = {
        "item_id": "picture-p3-1",
        "slide_or_page": 3,
        "region": {"x": 0.1, "y": 0.1, "width": 0.35, "height": 0.35,
                   "unit": "normalized"},
        "region_text": "Một bước đầu tư",
        "semantic_intent": ["picture candidate detected by Docling"],
    }
    subtitle_item = {
        "item_id": "figure-p4-1",
        "slide_or_page": 4,
        "region": {"x": 0.1, "y": 0.1, "width": 0.55, "height": 0.25,
                   "unit": "normalized"},
        "region_text": "goes sub tittle",
        "semantic_intent": ["figure candidate detected by PyMuPDF fallback"],
    }

    assert asc.semantic_item_id("input/Salary.pdf", salary_item, set()) == "salary-benefits-long-term-card"
    assert asc.semantic_item_id("input/Salary.pdf", investment_item, set()) == "investment-card"
    assert asc.semantic_item_id("input/Sun.Presentation.pdf", subtitle_item, set()) == "subtitle-visual"


def test_auto_stage_metadata_keeps_context_intent_with_region_text() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    item = {
        "item_id": "picture-p4-2",
        "slide_or_page": 4,
        "region": {"x": 0.21, "y": 0.66, "width": 0.55, "height": 0.23,
                   "unit": "normalized"},
        "region_text": "XAY DUNG DOI NGU AI\nXAY DUNG HE THONG TU DONG HOA",
        "semantic_intent": [
            "2. CONTENT XOAY QUANH - build AI team and automation system",
        ],
    }

    metadata = asc.metadata_for(
        "input/GUIDLINE_PRESENTATION_SUN.pdf", item, "team-visual")

    assert metadata["semantic_intent"][0] == "team visual"
    assert any("CONTENT XOAY QUANH" in value
               for value in metadata["semantic_intent"])
    assert "content" in metadata["keywords"]


def test_auto_stage_semantic_ids_use_intent_when_region_text_missing() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    item = {
        "item_id": "picture-p5-1",
        "slide_or_page": 5,
        "region": {"x": 0.25, "y": 0.28, "width": 0.47, "height": 0.21,
                   "unit": "normalized"},
        "region_text": (
            "Patrick E. Shorey\nRecreational therapist\n"
            "Mary T. Middleton\nPhysical meteorologist"
        ),
        "semantic_intent": [
            "This is a contributors slide. Insert your team here. Lorem ipsum dolor sit amet.",
        ],
    }

    item_id = asc.semantic_item_id("input/GUIDLINE_PRESENTATION_SUN.pdf", item, set())
    metadata = asc.metadata_for("input/GUIDLINE_PRESENTATION_SUN.pdf", item, item_id)

    assert item_id == "contributors-team-visual"
    assert not item_id.startswith("source-")
    assert metadata["keywords"][:3] == ["contributors", "team", "visual"]


def test_auto_stage_semantic_ids_filter_mixed_vietnamese_prose() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    item = {
        "item_id": "figure-p4-2",
        "slide_or_page": 4,
        "region": {"x": 0.1, "y": 0.4, "width": 0.5, "height": 0.2,
                   "unit": "normalized"},
        "region_text": "Đảm bảo goal thực tế",
    }

    item_id = asc.semantic_item_id("input/GUIDLINE_PRESENTATION_SUN.pdf", item, set())

    assert item_id == "goal-strip"
    assert "dam" not in item_id and "bao" not in item_id
    assert "thuc" not in item_id and "te" not in item_id


def test_auto_stage_semantic_ids_level_series_without_content_rule() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    item = {
        "item_id": "picture-p2-1",
        "slide_or_page": 2,
        "region": {"x": 0.1, "y": 0.2, "width": 0.7, "height": 0.18,
                   "unit": "normalized"},
        "region_text": "Design Operations\nLevel 1\nReview System\nLevel 2",
    }

    item_id = asc.semantic_item_id("input/source.pdf", item, set())

    assert item_id == "design-operations-review-system-strip"
    assert "source" not in item_id


def test_auto_stage_semantic_ids_metric_series_uses_labels_and_strip() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    item = {
        "item_id": "figure-p2-6",
        "slide_or_page": 2,
        "region": {"x": 0.19, "y": 0.68, "width": 0.52, "height": 0.24,
                   "unit": "normalized"},
        "region_text": "+30%\nRevenue\nTeam Size\n(110 Members)\n+30%",
    }

    item_id = asc.semantic_item_id("input/GUIDLINE_PRESENTATION_SUN.pdf", item, set())
    metadata = asc.metadata_for("input/GUIDLINE_PRESENTATION_SUN.pdf", item, item_id)

    assert item_id == "revenue-team-size-metric-strip"
    assert metadata["component_type"] == "strip"
    assert metadata["keywords"][:4] == ["revenue", "team", "size", "metric"]


def test_auto_stage_icon_reference_uses_page_context() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    item = {
        "item_id": "figure-p1-1",
        "slide_or_page": 1,
        "region": {"x": 0.54, "y": 0.17, "width": 0.38, "height": 0.15,
                   "unit": "normalized"},
        "region_text": "BOD\nPersonal\nCompany\nLearning & Sharing",
        "page_text": "ICON\n1. NHUNG ICON HAY XUAT HIEN",
        "semantic_intent": ["BOD Personal Company Learning Sharing"],
    }

    item_id = asc.semantic_item_id("input/GUIDLINE_PRESENTATION_SUN.pdf", item, set())
    metadata = asc.metadata_for("input/GUIDLINE_PRESENTATION_SUN.pdf", item, item_id)

    assert item_id == "icon-reference-sheet"
    assert metadata["component_type"] == "icon"
    assert metadata["layout_role"] == "icon reference sheet"
    assert "icon-set" in metadata["tags"]

    with tempfile.TemporaryDirectory() as tmp:
        item_dir = Path(tmp)
        (item_dir / "mapping.json").write_text(json.dumps({
            "component_type": "icon",
        }), encoding="utf-8")
        assert asc._is_icon_sheet_item(item_dir)


def test_auto_stage_overrides_stale_history_stable_id_for_auto_drafts() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        item_dir = tmp_path / "items" / "team-visual"
        item_dir.mkdir(parents=True)
        (item_dir / "mapping.json").write_text(json.dumps({
            "candidate_stable_id": "sun.component.xay-dung-oi-ngu-visual",
            "status": "staging",
            "source": {},
        }), encoding="utf-8")
        review = asc.metadata_for(
            "input/GUIDLINE_PRESENTATION_SUN.pdf",
            {"item_id": "picture-p4-2", "slide_or_page": 4,
             "region": {"x": 0.1, "y": 0.1, "width": 0.5, "height": 0.3,
                        "unit": "normalized"},
             "region_text": "XAY DUNG DOI NGU"},
            "team-visual",
        )
        review["candidate_id"] = "picture-p4-2"
        asc._augment_mapping(item_dir, review, "docling-run", {"item_id": "picture-p4-2"})
        mapping = read_text_slots.load_json(item_dir / "mapping.json")
        assert mapping["candidate_stable_id"] == "sun.component.team-visual"

        history = tmp_path / "history.json"
        history.write_text(json.dumps({"attempts": [{
            "extraction_id": "docling-run-team-visual",
            "item_id": "team-visual",
            "stable_id": "sun.component.xay-dung-oi-ngu-visual",
        }]}), encoding="utf-8")
        asc._sync_history_stable_id(
            history, "docling-run-team-visual", "team-visual",
            mapping["candidate_stable_id"],
        )
        synced = read_text_slots.load_json(history)
        assert synced["attempts"][0]["stable_id"] == "sun.component.team-visual"


def test_auto_stage_semantic_ids_suffix_existing_component_names() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    item = {
        "item_id": "picture-p2-2",
        "slide_or_page": 2,
        "region": {"x": 0.15, "y": 0.41, "width": 0.15, "height": 0.2,
                   "unit": "normalized"},
        "region_text": "TRANSLATOR",
    }

    item_id = asc.semantic_item_id(
        "input/GUIDLINE_PRESENTATION_SUN.pdf",
        item,
        {"translator-card"},
    )

    assert item_id == "translator-card-2"


def test_auto_stage_clusters_same_page_by_role_and_layout_row() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")

    def _record(item_id: str, x: float, y: float, width: float, height: float) -> dict:
        return {
            "candidate_id": item_id,
            "item": {
                "item_id": item_id,
                "slide_or_page": 2,
                "region": {"x": x, "y": y, "width": width, "height": height,
                           "unit": "normalized"},
            },
            "review": {"item_id": item_id, "display_name": item_id.replace("-", " ").title()},
            "item_dir": "unused",
            "stable_id": f"sun.component.{item_id}",
        }

    records = [
        _record("ai-coding-maturity-levels-strip", 0.05, 0.12, 0.8, 0.2),
        _record("translator-card", 0.05, 0.5, 0.18, 0.25),
        _record("coach-card", 0.28, 0.5, 0.18, 0.25),
    ]

    clusters = asc._cluster_staged_records(records)

    assert len(clusters) == 1, clusters
    assert [record["review"]["item_id"] for record in clusters[0]] == [
        "translator-card",
        "coach-card",
    ]


def test_auto_stage_group_records_keep_docling_candidate_order() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    records = []
    for candidate_id, item_id in [
        ("picture-p2-5", "coach-card"),
        ("picture-p2-2", "translator-card"),
        ("picture-p2-4", "driver-card"),
        ("picture-p2-3", "strategist-card"),
    ]:
        records.append({
            "candidate_id": candidate_id,
            "item": {
                "item_id": candidate_id,
                "slide_or_page": 2,
                "region": {"x": 0.5, "y": 0.5, "width": 0.15, "height": 0.2,
                           "unit": "normalized"},
            },
            "review": {"item_id": item_id, "display_name": item_id.replace("-", " ").title()},
            "item_dir": "unused",
            "stable_id": f"sun.component.{item_id}",
        })

    assert [r["review"]["item_id"] for r in asc._sort_group_records(records)] == [
        "translator-card",
        "strategist-card",
        "driver-card",
        "coach-card",
    ]


def test_auto_stage_groups_related_candidates_as_carousel_draft() -> None:
    import importlib
    import build_component_catalog as bcc

    asc = importlib.import_module("auto_stage_candidates")
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return

    with tempfile.TemporaryDirectory() as tmp:
        tmpp = Path(tmp)
        source = tmpp / "Role cards.pdf"
        doc = fitz.open()
        page = doc.new_page(width=400, height=220)
        page.draw_rect(fitz.Rect(40, 40, 170, 180), color=(1, 0.4, 0.2),
                       fill=(1, 0.9, 0.8))
        page.insert_text((62, 95), "TRANSLATOR", fontsize=18, color=(0, 0, 0))
        page.draw_rect(fitz.Rect(230, 40, 360, 180), color=(0.1, 0.2, 1),
                       fill=(0.8, 0.9, 1))
        page.insert_text((270, 95), "COACH", fontsize=18, color=(0, 0, 0))
        doc.save(source)
        doc.close()

        output_root = tmpp / "component-extractions"
        eid = "docling-auto-group-demo"
        adir = output_root / eid / "analysis"
        adir.mkdir(parents=True)
        (adir / "candidate-extraction-request.json").write_text(json.dumps({
            "extraction_id": eid,
            "source_path": str(source),
            "items": [
                {
                    "item_id": "picture-p1-1",
                    "slide_or_page": 1,
                    "region": {"x": 0.1, "y": 0.18, "width": 0.35, "height": 0.66,
                               "unit": "normalized"},
                    "object_ids": [],
                    "requested_type": "component",
                    "semantic_intent": ["picture candidate detected by Docling"],
                    "notes": "Translator card",
                    "replacement_for": None,
                },
                {
                    "item_id": "picture-p1-2",
                    "slide_or_page": 1,
                    "region": {"x": 0.55, "y": 0.18, "width": 0.35, "height": 0.66,
                               "unit": "normalized"},
                    "object_ids": [],
                    "requested_type": "component",
                    "semantic_intent": ["picture candidate detected by Docling"],
                    "notes": "Coach card",
                    "replacement_for": None,
                },
            ],
        }), encoding="utf-8")

        hist = tmpp / "history.json"
        hist.write_text('{"attempts":[]}', encoding="utf-8")
        reg = tmpp / "registry.json"
        reg.write_text('{"items":[]}', encoding="utf-8")
        summary = asc.stage_run(
            eid,
            root=output_root,
            output_root=output_root,
            history=hist,
            registry=reg,
            rebuild_catalog=False,
        )
        assert summary["staged"] == 2, summary
        assert summary["grouped"] == 1, summary
        group = summary["group_item"]
        assert group["item_id"] == "translator-coach-card-set"
        group_dir = output_root / group["extraction_id"] / "items" / group["item_id"]
        group_mapping = read_text_slots.load_json(group_dir / "mapping.json")
        assert group_mapping["component_type"] == "component-set"
        assert len(group_mapping["collection_children"]) == 2
        assert (group_dir / "artifact" / "components" / "components-manifest.json").is_file()
        group_thumb = group_dir / "preview" / "thumbnail.png"
        assert group_thumb.is_file()
        png = group_thumb.read_bytes()
        assert png[:8] == b"\x89PNG\r\n\x1a\n"
        thumb_width = int.from_bytes(png[16:20], "big")
        thumb_height = int.from_bytes(png[20:24], "big")
        assert thumb_width > thumb_height, (thumb_width, thumb_height)
        compact = asc.compact_summary(summary)
        assert "artifact_log" not in compact["items"][0]
        assert compact["items"][0]["artifact_log_lines"] > 0

        catalog_path = tmpp / "catalog-data.json"
        old_argv = sys.argv[:]
        sys.argv = [
            "build_component_catalog.py",
            "--registry", str(reg),
            "--extractions", str(output_root),
            "--output", str(catalog_path),
        ]
        try:
            assert bcc.main() == 0
        finally:
            sys.argv = old_argv
        catalog = read_text_slots.load_json(catalog_path)
        ids = [item["id"] for item in catalog["items"]]
        assert "sun.component.translator-card" not in ids
        assert "sun.component.coach-card" not in ids
        draft = next(item for item in catalog["items"]
                     if item["id"] == "sun.component.translator-coach-card-set")
        assert draft["publish_readiness"]["ready"], draft["publish_readiness"]
        assert [image["label"] for image in draft["images"]][:6] == [
            "Full component",
            "Full component (Text-free)",
            "Translator Card",
            "Translator Card (Text-free)",
            "Coach Card",
            "Coach Card (Text-free)",
        ]


def test_auto_stage_group_text_free_svg_rewrites_component_asset_refs() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    with tempfile.TemporaryDirectory() as tmp:
        tmpp = Path(tmp)
        child_artifact = tmpp / "child" / "artifact"
        child_assets = child_artifact / "assets"
        child_assets.mkdir(parents=True)
        (child_assets / "icon.png").write_bytes(b"png")
        child_svg = child_artifact / "visual.svg"
        child_svg.write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" '
            'xmlns:xlink="http://www.w3.org/1999/xlink">'
            '<defs><g id="mask"/></defs>'
            '<image xlink:href="assets/icon.png"/>'
            '<use xlink:href="#mask"/></svg>',
            encoding="utf-8",
        )

        parent_artifact = tmpp / "parent" / "artifact"
        dest_svg = parent_artifact / "components" / "child-text-free.svg"
        asc._copy_svg_with_assets(
            child_svg, dest_svg, parent_artifact / "assets", "../assets/")

        copied = dest_svg.read_text(encoding="utf-8")
        assert 'xlink:href="../assets/icon.png"' in copied, copied
        assert 'xlink:href="assets/icon.png"' not in copied, copied
        assert 'xlink:href="#mask"' in copied, copied
        assert (parent_artifact / "assets" / "icon.png").read_bytes() == b"png"


def test_auto_stage_existing_stable_ids_ignore_skipped_outputs() -> None:
    import importlib

    asc = importlib.import_module("auto_stage_candidates")
    with tempfile.TemporaryDirectory() as tmp:
        root = Path(tmp)
        skipped = root / "skip" / "items" / "blank" / "mapping.json"
        active = root / "active" / "items" / "card" / "mapping.json"
        skipped.parent.mkdir(parents=True)
        active.parent.mkdir(parents=True)
        skipped.write_text(json.dumps({
            "status": "skipped",
            "candidate_stable_id": "sun.component.blank",
        }), encoding="utf-8")
        active.write_text(json.dumps({
            "status": "staging",
            "candidate_stable_id": "sun.component.card",
        }), encoding="utf-8")

        assert asc._existing_stable_ids(root) == {"sun.component.card"}


def test_catalog_has_no_candidate_review_top_tab() -> None:
    html = (SCRIPTS.parents[1] / "slide-system" / "catalog" / "index.html").read_text(encoding="utf-8")
    assert 'data-section="review"' not in html
    assert 'id="section-review"' not in html


def test_catalog_server_exposes_no_candidate_review_routes() -> None:
    source = (SCRIPTS.parents[1] / "slide-system" / "catalog" / "catalog_server.py").read_text(encoding="utf-8")
    assert "/api/candidates" not in source
    assert "_serve_candidate" not in source
    assert "_candidate_segments" not in source


def test_catalog_server_parses_stage_candidate_booleans() -> None:
    path = SCRIPTS.parents[1] / "slide-system" / "catalog" / "catalog_server.py"
    spec = importlib.util.spec_from_file_location("catalog_server_under_test", path)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    assert module.body_bool({}, "build_artifacts", True) is True
    assert module.body_bool({"build_artifacts": False}, "build_artifacts", True) is False
    assert module.body_bool({"build_artifacts": "false"}, "build_artifacts", True) is False
    assert module.body_bool({"build_artifacts": "0"}, "build_artifacts", True) is False
    assert module.body_bool({"build_artifacts": "true"}, "build_artifacts", False) is True
    try:
        module.body_bool({"build_artifacts": "sometimes"}, "build_artifacts", True)
    except ValueError as exc:
        assert "build_artifacts" in str(exc)
    else:
        raise AssertionError("invalid boolean strings must fail")


# --------------------------------------------------------------------------- #
def _run_all() -> int:
    tests = [v for k, v in sorted(globals().items()) if k.startswith("test_") and callable(v)]
    failed = 0
    for t in tests:
        try:
            t()
            print(f"  PASS  {t.__name__}")
        except Exception as exc:  # noqa: BLE001
            failed += 1
            print(f"  FAIL  {t.__name__}: {exc}")
    print(f"\n{len(tests) - failed}/{len(tests)} passed")
    return 1 if failed else 0


if __name__ == "__main__":
    raise SystemExit(_run_all())
