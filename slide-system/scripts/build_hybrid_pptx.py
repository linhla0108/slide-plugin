#!/usr/bin/env python3
"""Build an editable hybrid PPTX from slide renders + DOM text layout.

This is the generalised version of the per-job build scripts (e.g.
build_v3_hybrid_editable.py). It works the same way:

  1. Each slide render (slide-01-bg.png …) becomes a full-slide background image.
  2. Every text item in export-layout.json is re-created as a native PowerPoint
     text box — so slide copy stays editable in PowerPoint / Keynote.

The result is a "hybrid editable" PPTX: visually identical to the HTML deck,
with native editable text. Complex styled text (gradient fills, custom SVG
glyphs, blend-mode effects) may look slightly different when edited because
PowerPoint cannot replicate all CSS effects — that is expected and documented.

This script needs only standard pip packages (no Claude Code required):
    pip install python-pptx Pillow

Usage:
    python3 build_hybrid_pptx.py \\
        --layout   qa/export-layout.json \\
        --renders  qa/export-renders/ \\
        --slides   8 \\
        --output   exports/deck-editable.pptx

Options:
    --layout      Path to export-layout.json produced by capture-slides.js
    --renders     Directory containing slide-01-bg.png, slide-02-bg.png …
    --slides      Total slide count
    --output      Output .pptx path
    --layout-w    Pixel width used when the layout was captured (default: 1920)
    --layout-h    Pixel height used when the layout was captured (default: 1080)
    --font        Font name to embed in text boxes (default: Proxima Nova)
    --fallback-font  Font name used when --font is unavailable on the viewer's
                  machine (default: Arial). PowerPoint substitutes this
                  automatically if Proxima Nova is not installed.
"""

from __future__ import annotations

import argparse
import json
import re
import zipfile
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    raise SystemExit(
        "python-pptx not found. Run:  pip install python-pptx\n"
        "(Also install Pillow if not present: pip install Pillow)"
    )

SLIDE_W = Inches(13.333333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)


def slide_dimensions(canvas_w: float, canvas_h: float) -> tuple[float, float]:
    """Fit the source aspect ratio within the standard 13.333 x 7.5in box."""
    if canvas_w <= 0 or canvas_h <= 0:
        raise ValueError("Canvas dimensions must be positive")
    aspect = canvas_w / canvas_h
    widescreen = 13.333333 / 7.5
    if abs(aspect - 16 / 9) < 1e-6:
        return 13.333333, 7.5
    if aspect >= widescreen:
        return 13.333333, 13.333333 / aspect
    return 7.5 * aspect, 7.5


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description=__doc__,
                                formatter_class=argparse.RawDescriptionHelpFormatter)
    p.add_argument("--layout",       required=True, help="Path to export-layout.json")
    p.add_argument("--renders",      required=True, help="Directory with slide-XX-bg.png files")
    p.add_argument("--slides",       required=True, type=int, help="Number of slides")
    p.add_argument("--output",       required=True, help="Output .pptx path")
    p.add_argument("--layout-w",     type=float, default=1920.0,
                   help="Pixel width of the layout coordinate space (default: 1920)")
    p.add_argument("--layout-h",     type=float, default=1080.0,
                   help="Pixel height of the layout coordinate space (default: 1080)")
    p.add_argument("--font",         default="Proxima Nova",
                   help="Primary font name for text boxes (default: Proxima Nova)")
    p.add_argument("--fallback-font", default="Arial",
                   help="Fallback font if primary is unavailable (default: Arial)")
    return p.parse_args()


def parse_css_color(value: str) -> RGBColor:
    """Convert css rgb(...) / rgba(...) / #rrggbb to RGBColor."""
    if value.startswith("#"):
        hex_val = value.lstrip("#")
        if len(hex_val) == 3:
            hex_val = "".join(c * 2 for c in hex_val)
        r, g, b = int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16)
        return RGBColor(r, g, b)
    nums = [float(n) for n in re.findall(r"[\d.]+", value)]
    if len(nums) < 3:
        return RGBColor(248, 250, 252)
    r, g, b = [max(0, min(255, round(n))) for n in nums[:3]]
    return RGBColor(r, g, b)


def px_to_inches(px_val: float, layout_px: float, slide_inches: float) -> float:
    """Convert layout-space pixels to slide inches."""
    return (px_val / layout_px) * float(slide_inches / Inches(1))


def font_pt(item: dict, layout_w: float) -> float:
    """Convert CSS fontSize (px) to points for the slide.

    VERIFIED against the proven Phase 1 build (build_v3_hybrid_editable.py):
    a 1920px-wide canvas maps to a 13.333in (=960pt) slide, so the exact
    px->pt factor is 960/1920 = 0.5 — NOT the generic 96dpi 0.75.

    fontSize in export-layout.json is the UNSCALED computed CSS px (e.g. 102px),
    captured with `noscale` set so it is already in the authored design-px
    space. The factor below is design_pt_width / design_px_width.
    """
    css_px = float(str(item.get("fontSize", "18px")).replace("px", "").strip())
    SLIDE_W_PT = 13.333333 * 72.0  # 960pt
    factor = SLIDE_W_PT / layout_w  # = 0.5 when layout_w == 1920
    return max(4.0, round(css_px * factor, 1))


def box_inches(item: dict, layout_w: float, layout_h: float) -> tuple[float, float, float, float]:
    """Return (left, top, width, height) in inches for the slide."""
    slide_w_in = 13.333333
    slide_h_in = 7.5
    x = item["x"] / layout_w * slide_w_in
    y = item["y"] / layout_h * slide_h_in
    w = max(item["w"] / layout_w * slide_w_in, 0.08)
    h = max(item["h"] / layout_h * slide_h_in * 1.35, 0.14)  # 35% taller for wrapping
    return x, y, w, h


def add_text_box(slide, item: dict, font_name: str, layout_w: float, layout_h: float) -> None:
    x, y, w, h = box_inches(item, layout_w, layout_h)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = f"Editable: {item['text'][:40]}"

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    para = tf.paragraphs[0]
    align = item.get("align", "start")
    para.alignment = PP_ALIGN.CENTER if align == "center" else (
        PP_ALIGN.RIGHT if align == "right" else PP_ALIGN.LEFT
    )

    run = para.add_run()
    run.text = item["text"]
    run.font.name = font_name
    run.font.size = Pt(font_pt(item, layout_w))
    run.font.bold = int(str(item.get("fontWeight", "400")).replace("bold", "700").replace("normal", "400")) >= 700
    run.font.color.rgb = parse_css_color(item.get("color", "rgb(248,250,252)"))


# ---------------------------------------------------------------------------
# LAYERED MODE (v2 — EXPORT-PPTX-3LAYER-PLAN.md §3.2)
# Reads ONE export-manifest.json and composes base → (overlay|text interleaved
# by z) per slide. The flat path below is the frozen v1 code and must not
# change (isolation rule #1); that is why this is a separate function set.
# ---------------------------------------------------------------------------

def parse_layered_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Layered (3-layer) PPTX build from export-manifest.json")
    p.add_argument("--manifest", required=True, help="Path to export-manifest.json (layered)")
    p.add_argument("--renders",  required=True, help="Directory with the PNGs the manifest names")
    p.add_argument("--output",   required=True, help="Output .pptx path")
    p.add_argument("--font",     default="Proxima Nova")
    p.add_argument("--fallback-font", default="Arial")
    p.add_argument("--vector-root", default=None,
                   help="Directory vector_source paths resolve against (default: deck dir = renders dir)")
    return p.parse_args()


GENERIC_FAMILIES = {"serif", "sans-serif", "monospace", "cursive", "fantasy",
                    "system-ui", "ui-serif", "ui-sans-serif", "ui-monospace"}


def first_font_family(font_family: str, fallback: str) -> str:
    """First concrete family from a CSS font-family list (skip generics)."""
    for token in str(font_family or "").split(","):
        name = token.strip().strip('"').strip("'")
        if name and name.lower() not in GENERIC_FAMILIES:
            return name
    return fallback


def set_run_typefaces(run, name: str) -> None:
    """latin + ea + cs typefaces — PowerPoint substitutes the ea/cs face for
    Vietnamese diacritic glyphs unless all three are pinned to the same font."""
    from pptx.oxml.ns import qn
    run.font.name = name
    rpr = run.font._rPr
    for tag in ("a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", name)


def set_letter_spacing(run, letter_spacing: str, canvas_w: float,
                       slide_w_in: float = 13.333333) -> None:
    """CSS letter-spacing px → DrawingML spc (hundredths of a point)."""
    raw = str(letter_spacing or "").replace("px", "").strip()
    if not raw or raw == "normal":
        return
    try:
        px = float(raw)
    except ValueError:
        return
    pt = px * (slide_w_in * 72.0) / canvas_w  # same px→pt factor as font_pt()
    run.font._rPr.set("spc", str(int(round(pt * 100))))


SVG_BLIP_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"


def embed_svg_blip(slide, picture, svg_path: Path) -> None:
    """Attach the source SVG to an existing picture shape (PowerPoint 2016+
    renders the vector; every other viewer keeps the PNG fallback that is
    already the picture's blip). python-pptx has no API for this — the part
    and the <asvg:svgBlip> extension are injected manually."""
    from lxml import etree
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.opc.package import Part
    from pptx.oxml.ns import qn

    package = slide.part.package
    partname = package.next_partname("/ppt/media/image%d.svg")
    svg_part = Part(partname, "image/svg+xml", package, svg_path.read_bytes())
    rid = slide.part.relate_to(svg_part, RT.IMAGE)

    blip = picture._element.blipFill.find(qn("a:blip"))
    ext_lst = blip.find(qn("a:extLst"))
    if ext_lst is None:
        ext_lst = etree.SubElement(blip, qn("a:extLst"))
    ext = etree.SubElement(ext_lst, qn("a:ext"))
    ext.set("uri", SVG_BLIP_EXT_URI)
    # nsmap pins the conventional "asvg" prefix — lxml would otherwise emit
    # an auto prefix (ns0:), which is spec-valid but not what Office writes.
    svg_blip = etree.SubElement(ext, f"{{{ASVG_NS}}}svgBlip", nsmap={"asvg": ASVG_NS})
    svg_blip.set(qn("r:embed"), rid)


def add_native_shape(slide, payload: dict, canvas_w: float, canvas_h: float,
                     slide_w_in: float = 13.333333, slide_h_in: float = 7.5):
    """Real PPTX autoshape for a simple solid element — scales losslessly."""
    from pptx.enum.shapes import MSO_SHAPE

    b = payload["bounds"]
    x, y = b["x"] / canvas_w * slide_w_in, b["y"] / canvas_h * slide_h_in
    w, h = b["w"] / canvas_w * slide_w_in, b["h"] / canvas_h * slide_h_in
    radius = float(payload.get("radius") or 0)
    kind = payload.get("shape", "rect")
    if kind == "ellipse":
        shape_type = MSO_SHAPE.OVAL
    elif radius > 0:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        shape_type = MSO_SHAPE.RECTANGLE

    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = f"Native: {payload['id']}"
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE and min(b["w"], b["h"]) > 0:
        shape.adjustments[0] = max(0.0, min(0.5, radius / min(b["w"], b["h"])))
    shape.fill.solid()
    shape.fill.fore_color.rgb = parse_css_color(payload.get("fill", "rgb(200,200,200)"))
    border_w = float(payload.get("border_width") or 0)
    if border_w > 0:
        shape.line.color.rgb = parse_css_color(payload.get("border_color", "rgb(0,0,0)"))
        shape.line.width = Pt(border_w * (slide_w_in * 72.0) / canvas_w)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def apply_text_transform(text: str, transform: str) -> str:
    if transform == "uppercase":
        return text.upper()
    if transform == "lowercase":
        return text.lower()
    if transform == "capitalize":
        return " ".join(w[:1].upper() + w[1:] if w else w for w in text.split(" "))
    return text


def add_text_box_v2(slide, item: dict, font_name: str, canvas_w: float, canvas_h: float,
                    slide_w_in: float = 13.333333, slide_h_in: float = 7.5) -> None:
    """Layered text box: no 1.35 height hack — exact line spacing instead,
    and computed text-transform applied (P1 wrong-characters fix)."""
    x = item["x"] / canvas_w * slide_w_in
    y = item["y"] / canvas_h * slide_h_in
    w = max(item["w"] / canvas_w * slide_w_in, 0.08)
    h = max(item["h"] / canvas_h * slide_h_in * 1.05, 0.14)  # 5% slack only
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = f"Editable: {item['text'][:40]}"

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.auto_size = None
    for side in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, side, Inches(0))

    align = item.get("align", "start")
    alignment = PP_ALIGN.CENTER if align == "center" else (
        PP_ALIGN.RIGHT if align in ("right", "end") else PP_ALIGN.LEFT
    )
    font_px = float(str(item.get("fontSize", "18px")).replace("px", "").strip() or 18)
    lh_raw = str(item.get("lineHeight", "")).replace("px", "").strip()
    try:
        line_px = float(lh_raw)
    except ValueError:
        line_px = font_px * 1.2
    line_spacing = round(line_px / font_px, 3) if font_px > 0 else None

    content = apply_text_transform(item["text"], str(item.get("textTransform", "none")))
    # Explicit <br> breaks arrive as \n from capture — one paragraph per line.
    for index, line in enumerate(content.split("\n")):
        para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        para.alignment = alignment
        if line_spacing:
            para.line_spacing = line_spacing
        run = para.add_run()
        run.text = line
        set_run_typefaces(run, first_font_family(item.get("fontFamily", ""), font_name))
        set_letter_spacing(run, item.get("letterSpacing", ""), canvas_w, slide_w_in)
        css_px = float(str(item.get("fontSize", "18px")).replace("px", "").strip())
        run.font.size = Pt(max(4.0, round(css_px * slide_w_in * 72.0 / canvas_w, 1)))
        run.font.bold = int(str(item.get("fontWeight", "400")).replace("bold", "700").replace("normal", "400")) >= 700
        run.font.color.rgb = parse_css_color(item.get("color", "rgb(248,250,252)"))


def build_layered(args: argparse.Namespace) -> None:
    manifest_path = Path(args.manifest)
    renders_dir = Path(args.renders)
    output_path = Path(args.output)

    # Operational errors only — quality verdicts belong to the validator.
    if not manifest_path.exists():
        raise SystemExit(f"Manifest not found: {manifest_path}")
    try:
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as error:
        raise SystemExit(f"Manifest unparseable: {error}")
    for field in ("manifest_version", "mode", "slides"):
        if field not in manifest:
            raise SystemExit(f"Manifest missing required field: {field}")
    if manifest["manifest_version"] != 2:
        raise SystemExit(f"Unsupported manifest_version: {manifest['manifest_version']}")
    if manifest["mode"] != "layered":
        raise SystemExit(f"--manifest build requires mode=layered, got {manifest['mode']!r} "
                         "(flat decks use the v1 --layout path)")

    first = manifest["slides"][0] if manifest["slides"] else {}
    first_cw = float(first.get("canvasW") or manifest.get("canvasW") or 1920)
    first_ch = float(first.get("canvasH") or manifest.get("canvasH") or 1080)
    slide_w_in, slide_h_in = slide_dimensions(first_cw, first_ch)
    prs = Presentation()
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)
    blank_layout = prs.slide_layouts[6]

    text_count = picture_count = 0
    print(f"[build_hybrid_pptx] layered: {len(manifest['slides'])} slides → {output_path}")
    for entry in manifest["slides"]:
        cw = float(entry.get("canvasW") or manifest.get("canvasW") or 1920)
        ch = float(entry.get("canvasH") or manifest.get("canvasH") or 1080)
        base_png = renders_dir / entry["base"]["png"]
        if not base_png.exists():
            raise SystemExit(f"Render not found: {base_png}")

        slide = prs.slides.add_slide(blank_layout)
        if abs(cw / ch - first_cw / first_ch) > 1e-6:
            raise SystemExit("All slides in one PPTX must use the same canvas aspect ratio")
        slide.shapes.add_picture(
            str(base_png), 0, 0,
            width=Inches(slide_w_in), height=Inches(slide_h_in),
        )
        picture_count += 1

        # ONE merged z list: overlays, native shapes and text interleave
        # exactly as captured (C8).
        layers: list[tuple[int, str, dict]] = []
        for ov in entry.get("objects", []):
            layers.append((int(ov.get("z", 0)), "overlay", ov))
        for nv in entry.get("natives", []):
            layers.append((int(nv.get("z", 0)), "native", nv))
        for item in entry.get("text", []):
            if not item.get("text", "").strip():
                continue
            layers.append((int(item.get("z", 0)), "text", item))
        layers.sort(key=lambda t: t[0])

        vector_root = Path(args.vector_root) if args.vector_root else renders_dir
        svg_embedded = 0
        for _, kind, payload in layers:
            if kind == "overlay":
                png = renders_dir / payload["png"]
                if not png.exists():
                    raise SystemExit(f"Render not found: {png}")
                b = payload["bounds"]
                pic = slide.shapes.add_picture(
                    str(png),
                    Inches(b["x"] / cw * slide_w_in), Inches(b["y"] / ch * slide_h_in),
                    width=Inches(b["w"] / cw * slide_w_in),
                    height=Inches(b["h"] / ch * slide_h_in),
                )
                pic.name = f"Overlay: {payload['id']}"
                picture_count += 1
                # svgBlip: only when a vector source exists AND css effects do
                # not alter the rendered look (plan round 5 rule).
                src = payload.get("vector_source")
                if src and not payload.get("css_effects"):
                    svg_path = vector_root / src
                    if svg_path.exists():
                        embed_svg_blip(slide, pic, svg_path)
                        svg_embedded += 1
                    else:
                        print(f"  WARNING slide {entry['slide']:02d}: vector_source "
                              f"'{src}' not found under {vector_root} — PNG only")
            elif kind == "native":
                add_native_shape(slide, payload, cw, ch, slide_w_in, slide_h_in)
            else:
                add_text_box_v2(
                    slide, payload, args.font, cw, ch, slide_w_in, slide_h_in
                )
                text_count += 1
        print(f"  slide {entry['slide']:02d} ✓ ({len(entry.get('objects', []))} overlays, "
              f"{len(entry.get('natives', []))} natives, {svg_embedded} svgBlip)")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    # Informational audit only — exit verdicts live in validate_export_objects.py.
    print(f"[build_hybrid_pptx] layered done: {picture_count} pictures, {text_count} text boxes")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build(args: argparse.Namespace) -> None:
    layout_path = Path(args.layout)
    renders_dir = Path(args.renders)
    output_path = Path(args.output)

    if not layout_path.exists():
        raise FileNotFoundError(f"Layout file not found: {layout_path}")
    if not renders_dir.is_dir():
        raise FileNotFoundError(f"Renders directory not found: {renders_dir}")

    layout: list[dict] = json.loads(layout_path.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    print(f"[build_hybrid_pptx] Building {args.slides} slides → {output_path}")

    text_count = 0
    for slide_index in range(1, args.slides + 1):
        render_path = renders_dir / f"slide-{slide_index:02d}-bg.png"
        if not render_path.exists():
            raise FileNotFoundError(f"Render not found: {render_path}")

        slide = prs.slides.add_slide(blank_layout)
        # Full-slide background image
        slide.shapes.add_picture(str(render_path), 0, 0,
                                  width=SLIDE_W, height=SLIDE_H)

        # Native text boxes.
        #
        # This script consumes export-layout.json produced by capture-slides.js,
        # which captures with deck-stage `noscale` set. In that mode BOTH the box
        # coordinates AND the unscaled CSS fontSize live in the same authored
        # design space (canvasW x canvasH, normally 1920x1080), so a single
        # canvas dimension scales both correctly:
        #   box:  x / canvasW * 13.333in           font: css_px * 960 / canvasW
        # (For canvasW = 1920 the font factor is 960/1920 = 0.5, matching the
        #  proven Phase 1 build.)
        #
        # NOTE: this does NOT correctly process the legacy Phase-1
        # export-layout.json, whose coordinates were captured in 804x452.25
        # space while its fonts stayed in 1920 space (a split the old per-job
        # build patched by hand). Always regenerate layout via capture-slides.js.
        slide_data = next((s for s in layout if s.get("slide") == slide_index), None)
        if slide_data:
            cw = float(slide_data.get("canvasW") or args.layout_w)
            ch = float(slide_data.get("canvasH") or args.layout_h)
            if not slide_data.get("canvasW"):
                print(f"  WARNING slide {slide_index:02d}: layout has no canvasW; "
                      f"using --layout-w {args.layout_w}. Font sizes may be wrong "
                      f"if this is not freshly captured data.")
            for item in slide_data.get("items", []):
                if not item.get("text", "").strip():
                    continue
                add_text_box(slide, item, args.font, cw, ch)
                text_count += 1

        print(f"  slide {slide_index:02d} ✓")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    # Quick audit
    with zipfile.ZipFile(output_path) as zf:
        slide_xmls = [n for n in zf.namelist()
                      if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        pptx_runs = sum(
            zf.read(n).decode("utf-8", errors="ignore").count("<a:t>")
            for n in slide_xmls
        )

    print(f"\n[build_hybrid_pptx] Done")
    print(f"  output       → {output_path}")
    print(f"  slides       → {len(slide_xmls)}")
    print(f"  text boxes   → {text_count} layout items")
    print(f"  pptx <a:t>   → {pptx_runs} text runs")
    print(f"\nNOTE: Visual backgrounds are PNG images. Slide text is natively")
    print(f"      editable. Complex CSS effects (gradients, blends) in the")
    print(f"      background will not change when text is edited — that is expected.")


if __name__ == "__main__":
    import sys
    if "--manifest" in sys.argv:
        build_layered(parse_layered_args())
    else:
        build(parse_args())  # frozen v1 flat path
