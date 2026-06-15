#!/usr/bin/env python3
"""Test which export stack works on THIS machine, end to end.

Goal: prove whether the lightweight stack (markitdown + Playwright) can replace
the heavy system stack (LibreOffice + Poppler) for the two jobs this repo cares
about:

  A. Create a deck and export an EDITABLE .pptx
       → capture-slides.js  +  build_hybrid_pptx.py   (the current scripts)
  B. Export a deck to PDF
       → export-pdf.js  (Playwright)   vs   soffice  (LibreOffice)
  C. Read an existing .pptx back as text
       → markitdown        vs   soffice→pdf→pdftotext

It builds a tiny self-contained sample deck, serves it locally, runs YOUR real
scripts against it, validates the output, times each path, and prints a verdict
table. Nothing is installed — missing tools are reported as SKIP with the exact
install hint.

Usage:
    python3 slide-system/scripts/test_export_stack.py
    python3 slide-system/scripts/test_export_stack.py --keep   # keep work dir
    python3 slide-system/scripts/test_export_stack.py --json    # machine summary

Exit code: 0 if the lightweight stack covers jobs A+B+C, 1 otherwise.
"""

from __future__ import annotations

import argparse
import json
import shutil
import socket
import subprocess
import sys
import tempfile
import threading
import time
from http.server import HTTPServer, SimpleHTTPRequestHandler
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parents[2]
SCRIPTS = REPO_ROOT / "slide-system" / "scripts"
SLIDES = 2
WIDTH, HEIGHT = 1920, 1080
SAMPLE_WORDS = ["SUN.RISER", "Editable", "Hybrid", "Roadmap"]  # must survive round-trip


# --------------------------------------------------------------------------- #
# Small helpers
# --------------------------------------------------------------------------- #
class R:
    """One check result."""
    def __init__(self, name): self.name = name; self.status = "SKIP"; self.ms = None; self.note = ""
    def ok(self, ms=None, note=""): self.status = "PASS"; self.ms = ms; self.note = note; return self
    def fail(self, note=""): self.status = "FAIL"; self.note = note; return self
    def skip(self, note=""): self.status = "SKIP"; self.note = note; return self


def have(*names) -> str | None:
    for n in names:
        p = shutil.which(n)
        if p:
            return p
    return None


def py_mod(mod: str) -> bool:
    return subprocess.run([sys.executable, "-c", f"import {mod}"],
                          capture_output=True).returncode == 0


def run(cmd, cwd=None, timeout=180):
    """Run a command, return (rc, stdout+stderr, elapsed_ms)."""
    t0 = time.time()
    p = subprocess.run(cmd, cwd=cwd, capture_output=True, text=True, timeout=timeout)
    ms = int((time.time() - t0) * 1000)
    return p.returncode, (p.stdout or "") + (p.stderr or ""), ms


def free_port() -> int:
    s = socket.socket()
    s.bind(("127.0.0.1", 0))
    port = s.getsockname()[1]
    s.close()
    return port


# --------------------------------------------------------------------------- #
# Sample deck — plain HTML, no brand-pack dependency
# --------------------------------------------------------------------------- #
SAMPLE_HTML = """<!doctype html><html><head><meta charset="utf-8"><style>
  html,body{margin:0;padding:0;background:#FFFDF8}
  .slide{position:absolute;top:0;left:0;width:1920px;height:1080px;display:none;
         font-family:Arial,sans-serif;color:#171717;box-sizing:border-box;padding:120px}
  .slide.active{display:block}
  .slide.s1{background:#FF5533;color:#fff}
  h1{font-size:140px;font-weight:800;margin:0}
  p{font-size:56px;margin-top:40px}
  .tag{font-size:40px;color:#3333FF;font-weight:700}
  .accent{position:absolute;right:160px;top:160px;width:420px;height:420px;border-radius:60px;
          background:linear-gradient(135deg,#3333FF,#9933FF);box-shadow:0 24px 48px rgba(0,0,0,.4)}
</style></head><body>
  <div class="slide s1 active"><h1>SUN.RISER 2026</h1><p>Editable Hybrid Export Test</p>
       <div class="accent" data-export-layer="overlay" data-export-id="accent-card"></div></div>
  <div class="slide s2"><span class="tag">Section 02</span><h1>Roadmap</h1>
       <p>This text must stay editable in PowerPoint.</p>
       <div style="position:absolute;left:120px;bottom:120px;width:600px;height:24px;
                   background:#3333FF;border-radius:12px"
            data-export-native="rect" data-export-id="footer-bar"></div></div>
<script>
  function goToSlide(n){
    var s=document.querySelectorAll('.slide');
    s.forEach(function(el,i){el.classList.toggle('active', i===n);});
  }
  goToSlide(0);
</script></body></html>"""


def serve(directory: Path, port: int) -> HTTPServer:
    handler = lambda *a, **k: SimpleHTTPRequestHandler(*a, directory=str(directory), **k)
    httpd = HTTPServer(("127.0.0.1", port), handler)
    threading.Thread(target=httpd.serve_forever, daemon=True).start()
    return httpd


# --------------------------------------------------------------------------- #
# Validators
# --------------------------------------------------------------------------- #
def pptx_is_editable(pptx_path: Path) -> tuple[bool, str]:
    """True when the PPTX has NATIVE text boxes containing the sample words."""
    if not py_mod("pptx"):
        # Fall back to a raw-XML check: editable text lives in <a:t> runs.
        try:
            import zipfile
            xml = ""
            with zipfile.ZipFile(pptx_path) as z:
                for n in z.namelist():
                    if n.startswith("ppt/slides/slide") and n.endswith(".xml"):
                        xml += z.read(n).decode("utf-8", "ignore")
            hits = [w for w in SAMPLE_WORDS if w in xml]
            return (len(hits) > 0, f"raw-xml <a:t> words found: {hits}")
        except Exception as e:
            return (False, f"could not inspect: {e}")
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    runs, found = 0, set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run_ in para.runs:
                        runs += 1
                        for w in SAMPLE_WORDS:
                            if w.lower() in run_.text.lower():
                                found.add(w)
    ok = runs > 0 and len(found) > 0
    return (ok, f"{runs} native text runs, sample words: {sorted(found)}")


def pdf_pages(pdf_path: Path) -> int:
    """Page count without external tools (parse the PDF trailer best-effort)."""
    data = pdf_path.read_bytes()
    # Count /Type /Page (not /Pages) — good enough for a sanity check.
    return data.count(b"/Type /Page") or data.count(b"/Type/Page")


# --------------------------------------------------------------------------- #
# Phases
# --------------------------------------------------------------------------- #
def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--keep", action="store_true", help="Keep the work directory")
    ap.add_argument("--json", action="store_true", help="Print JSON summary only")
    args = ap.parse_args()

    results: list[R] = []
    node = have("node")
    npm_pw = (REPO_ROOT / "node_modules" / "playwright").is_dir()
    has_pptx = py_mod("pptx")
    has_markitdown = py_mod("markitdown") or have("markitdown") is not None
    soffice = have("soffice", "libreoffice")
    pdftotext = have("pdftotext")
    pdftoppm = have("pdftoppm")

    # ---- Tool probe ----
    probe = {
        "node": bool(node),
        "playwright (npm)": npm_pw,
        "python-pptx": has_pptx,
        "markitdown": has_markitdown,
        "libreoffice (soffice)": bool(soffice),
        "poppler (pdftotext/pdftoppm)": bool(pdftotext or pdftoppm),
    }

    work = Path(tempfile.mkdtemp(prefix="export-stack-test-"))
    deck_dir = work / "deck"; deck_dir.mkdir()
    out_dir = work / "out"; out_dir.mkdir()
    renders = out_dir / "renders"; renders.mkdir()
    (deck_dir / "index.html").write_text(SAMPLE_HTML, encoding="utf-8")

    port = free_port()
    httpd = serve(deck_dir, port)
    url = f"http://127.0.0.1:{port}/index.html"

    pptx_out = out_dir / "deck-editable.pptx"
    pdf_pw = out_dir / "deck-playwright.pdf"
    pdf_lo = out_dir / "deck-libreoffice.pdf"

    try:
        # ================= JOB A: editable PPTX (capture + build) =============
        a_cap = R("A1 capture-slides.js (Playwright → PNG + layout)")
        a_build = R("A2 build_hybrid_pptx.py (→ editable .pptx)")
        a_edit = R("A3 editable-text validation")
        if not node:
            a_cap.skip("install Node.js 18+ (nodejs.org)")
        elif not npm_pw:
            a_cap.skip("run ./slide-system/scripts/setup.sh (npm i playwright)")
        else:
            rc, log, ms = run([node, str(SCRIPTS / "capture-slides.js"),
                               "--url", url, "--slides", str(SLIDES),
                               "--out-dir", str(renders),
                               "--showJs", "goToSlide({n})", "--selector", ".slide.active"])
            layout = renders / "export-layout.json"
            pngs = list(renders.glob("slide-*-bg.png"))
            if rc == 0 and layout.exists() and len(pngs) == SLIDES:
                a_cap.ok(ms, f"{len(pngs)} PNGs + layout")
            else:
                a_cap.fail(log.strip().splitlines()[-1] if log.strip() else f"rc={rc}")

        if a_cap.status == "PASS":
            if not has_pptx:
                a_build.skip("pip install python-pptx Pillow")
            else:
                rc, log, ms = run([sys.executable, str(SCRIPTS / "build_hybrid_pptx.py"),
                                   "--layout", str(renders / "export-layout.json"),
                                   "--renders", str(renders), "--slides", str(SLIDES),
                                   "--output", str(pptx_out),
                                   "--font", "Arial", "--fallback-font", "Arial"])
                if rc == 0 and pptx_out.exists():
                    a_build.ok(ms, f"{pptx_out.stat().st_size // 1024} KB")
                    ok, note = pptx_is_editable(pptx_out)
                    (a_edit.ok(note=note) if ok else a_edit.fail(note))
                else:
                    a_build.fail(log.strip().splitlines()[-1] if log.strip() else f"rc={rc}")
        results += [a_cap, a_build, a_edit]

        # ================= JOB B: PDF — Playwright vs LibreOffice =============
        b_pw = R("B1 export-pdf.js (Playwright, HTML→PDF)")
        if not node:
            b_pw.skip("install Node.js 18+")
        elif not npm_pw:
            b_pw.skip("run setup.sh (npm i playwright)")
        else:
            rc, log, ms = run([node, str(SCRIPTS / "export-pdf.js"),
                               "--url", url, "--slides", str(SLIDES),
                               "--showJs", "goToSlide({n})", "--output", str(pdf_pw)])
            if rc == 0 and pdf_pw.exists():
                b_pw.ok(ms, f"{pdf_pw.stat().st_size // 1024} KB, ~{pdf_pages(pdf_pw)} page(s)")
            else:
                b_pw.fail(log.strip().splitlines()[-1] if log.strip() else f"rc={rc}")

        b_lo = R("B2 LibreOffice (PPTX→PDF, baseline)")
        if not soffice:
            b_lo.skip("not installed — and NOT required (B1 replaces it)")
        elif not pptx_out.exists():
            b_lo.skip("no .pptx to convert (job A did not produce one)")
        else:
            rc, log, ms = run([soffice, "--headless", "--convert-to", "pdf",
                               "--outdir", str(out_dir), str(pptx_out)], timeout=240)
            made = out_dir / (pptx_out.stem + ".pdf")
            if rc == 0 and made.exists():
                made.rename(pdf_lo)
                b_lo.ok(ms, f"{pdf_lo.stat().st_size // 1024} KB")
            else:
                b_lo.fail(log.strip().splitlines()[-1] if log.strip() else f"rc={rc}")
        results += [b_pw, b_lo]

        # ================= JOB C: read .pptx — markitdown vs LO+poppler =======
        c_mi = R("C1 markitdown (read .pptx text)")
        if not pptx_out.exists():
            c_mi.skip("no .pptx produced by job A")
        elif not has_markitdown:
            c_mi.skip('pip install "markitdown[pptx]"')
        else:
            rc, log, ms = run([sys.executable, "-m", "markitdown", str(pptx_out)])
            hits = [w for w in SAMPLE_WORDS if w.lower() in log.lower()]
            (c_mi.ok(ms, f"recovered words: {hits}") if rc == 0 and hits
             else c_mi.fail(f"rc={rc}, words={hits}"))

        c_lo = R("C2 LibreOffice+Poppler (read baseline)")
        if not (soffice and (pdftotext or pdftoppm)):
            c_lo.skip("not installed — and NOT required (C1 replaces it)")
        elif not pdf_lo.exists():
            c_lo.skip("no LibreOffice PDF from B2")
        elif not pdftotext:
            c_lo.skip("pdftotext (poppler) not present")
        else:
            rc, log, ms = run([pdftotext, str(pdf_lo), "-"])
            hits = [w for w in SAMPLE_WORDS if w.lower() in log.lower()]
            (c_lo.ok(ms, f"recovered words: {hits}") if rc == 0 and hits
             else c_lo.fail(f"rc={rc}"))
        results += [c_mi, c_lo]

        # ================= JOB D: 3-layer export (orchestrator end-to-end) ====
        # Slide 1 declares one overlay → the PPTX must hold 2 pictures (base +
        # the overlay as its own shape) plus native text. Slide 2 has none →
        # exactly 1 picture. Layered must NOT emit export-layout.json
        # (isolation rule #3) — and the flat run above (job A) proves the v1
        # path still works unchanged on the very same deck.
        d_run = R("D1 export_pptx.py --mode layered (full chain a→f)")
        d_struct = R("D2 object separation (overlay = own shape)")
        if not (node and npm_pw and has_pptx):
            d_run.skip("needs Node+playwright+python-pptx (see hints above)")
        else:
            layered_dir = out_dir / "layered"
            layered_pptx = out_dir / "deck-layered.pptx"
            rc, log, ms = run([sys.executable, str(SCRIPTS / "export_pptx.py"),
                               "--html", str(deck_dir / "index.html"),
                               "--slides", str(SLIDES),
                               "--out-dir", str(layered_dir),
                               "--output", str(layered_pptx),
                               "--mode", "layered", "--font", "Arial",
                               "--showJs", "goToSlide({n})",
                               "--selector", ".slide.active"], timeout=300)
            if rc == 0 and layered_pptx.exists():
                d_run.ok(ms, "gate PASS (validate_export_objects)")
                from pptx import Presentation
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                prs = Presentation(str(layered_pptx))
                pics = [sum(1 for s in sl.shapes
                            if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
                        for sl in prs.slides]
                names = [s.name for sl in prs.slides for s in sl.shapes]
                no_shim = not (layered_dir / "export-layout.json").exists()
                if (pics == [2, 1] and "Overlay: accent-card" in names
                        and "Native: footer-bar" in names and no_shim):
                    d_struct.ok(note=f"pictures/slide={pics}, overlay + native shapes named, no v1 shim")
                else:
                    d_struct.fail(f"pictures={pics}, shim_absent={no_shim}, names={names[:6]}")
            else:
                d_run.fail(log.strip().splitlines()[-1] if log.strip() else f"rc={rc}")
        results += [d_run, d_struct]

        # ================= JOB E: artwork decomposition (anti-B11) ============
        # A full-page artwork SVG must come apart into per-object overlays:
        # full-bleed group → base-candidate (kept out of the snippet),
        # overlapping neighbors → one cluster, a wide multi-object group →
        # split by its children's disjoint bboxes.
        e_run = R("E1 decompose_svg_objects.py (artwork → per-object overlays)")
        if not (node and npm_pw):
            e_run.skip("needs Node+playwright (see hints above)")
        else:
            fixture = out_dir / "decompose-fixture.svg"
            fixture.write_text(
                '<svg xmlns="http://www.w3.org/2000/svg" '
                'xmlns:ns1="http://www.inkscape.org/namespaces/inkscape" '
                'width="1920" height="1080" viewBox="0 0 1920 1080">\n'
                '<g ns1:groupmode="layer" ns1:label="Layer 1">\n'
                '<g><rect x="0" y="0" width="1920" height="1080" fill="#eee"/></g>\n'
                '<g><circle cx="140" cy="140" r="40" fill="red"/></g>\n'
                '<g><rect x="150" y="150" width="60" height="60" fill="blue"/></g>\n'
                '<g><rect x="100" y="900" width="50" height="50" fill="green"/>'
                '<rect x="1500" y="900" width="50" height="50" fill="purple"/></g>\n'
                '</g>\n</svg>\n', encoding="utf-8")
            decomp_dir = out_dir / "decompose"
            rc, log, ms = run([sys.executable,
                               str(SCRIPTS / "decompose_svg_objects.py"),
                               "--svg", str(fixture),
                               "--out-dir", str(decomp_dir),
                               "--prefix", "fx"], timeout=120)
            manifest_path = decomp_dir / "decompose-manifest.json"
            if rc == 0 and manifest_path.exists():
                manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
                objects = manifest["objects"]
                frags_exist = all((decomp_dir / o["file"]).exists() for o in objects)
                snippet = (decomp_dir / "snippet.html").read_text(encoding="utf-8")
                checks = (
                    len(objects) == 3                       # circle+rect, 2 split rects
                    and len(manifest["base_candidates"]) == 1  # full-page rect
                    and sum("children" in p for o in objects for p in o["parts"]) == 2
                    and frags_exist
                    and snippet.count("data-export-layer=\"overlay\"") == 3
                    and all(
                        'viewBox="0 0 ' in (decomp_dir / o["file"]).read_text()
                        and 'transform="translate(' in (decomp_dir / o["file"]).read_text()
                        for o in objects
                    ))
                if checks:
                    e_run.ok(ms, "1 base-candidate + 3 objects (cluster + 2 split)")
                else:
                    e_run.fail(f"objects={[(o['id'], o['parts']) for o in objects]}, "
                               f"base={len(manifest['base_candidates'])}")
            else:
                e_run.fail(log.strip().splitlines()[-1] if log.strip() else f"rc={rc}")
        results += [e_run]

    finally:
        httpd.shutdown()
        if not args.keep:
            shutil.rmtree(work, ignore_errors=True)

    # ---- Verdict ----
    light_jobs = {
        "A editable PPTX": next(r for r in results if r.name.startswith("A3")).status,
        "B HTML→PDF":      next(r for r in results if r.name.startswith("B1")).status,
        "C read PPTX":     next(r for r in results if r.name.startswith("C1")).status,
        "D layered 3-lớp": next(r for r in results if r.name.startswith("D2")).status,
        "E decompose artwork": next(r for r in results if r.name.startswith("E1")).status,
    }
    light_ok = all(s == "PASS" for s in light_jobs.values())

    summary = {
        "tools": probe,
        "results": [{"check": r.name, "status": r.status, "ms": r.ms, "note": r.note} for r in results],
        "lightweight_stack_jobs": light_jobs,
        "lightweight_replaces_heavy": light_ok,
        "work_dir": str(work) if args.keep else None,
    }

    if args.json:
        print(json.dumps(summary, indent=2))
        return 0 if light_ok else 1

    print("\n=== Tool probe ===")
    for k, v in probe.items():
        print(f"  [{'OK ' if v else '-- '}] {k}")

    print("\n=== Checks (markitdown + Playwright  vs  LibreOffice + Poppler) ===")
    for r in results:
        ms = f"{r.ms:>6} ms" if r.ms is not None else "      –  "
        print(f"  [{r.status:<4}] {ms}  {r.name}")
        if r.note:
            print(f"             ↳ {r.note}")

    print("\n=== Verdict: can the lightweight stack do the job? ===")
    for job, st in light_jobs.items():
        print(f"  [{st:<4}] {job}")
    if light_ok:
        print("\n  ✅ markitdown + Playwright cover A+B+C. LibreOffice + Poppler are NOT needed.")
    else:
        print("\n  ⚠ Some lightweight jobs did not PASS. Install the hinted package(s) above")
        print("    and re-run. SKIP usually means a missing pip/npm install, not a real failure.")
    if args.keep:
        print(f"\n  Work dir kept at: {work}")
    return 0 if light_ok else 1


if __name__ == "__main__":
    raise SystemExit(main())
